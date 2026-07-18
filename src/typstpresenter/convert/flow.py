"""
Idiomatic, human-editable slide emission ("flow mode", ``minimal=True``).

Instead of placing every shape at its absolute PPTX coordinates, slides are
rebuilt as normal Typst document flow: the title placeholder becomes a
``==`` heading, bullet paragraphs become native ``-`` list items, the
dominant body font size becomes one ``#set text(...)`` rule, bold/italic
runs become ``*...*`` / ``_..._`` markup, and pictures/tables/diagrams are
emitted in reading order. Simplicity and a coherent, overlap-free layout
take precedence over exact fidelity to the original coordinates.
"""

from __future__ import annotations

import re
from collections import Counter
from dataclasses import dataclass, replace
from pathlib import Path

from typstpresenter.convert.cetz import emit_cetz_shape, is_diagram_shape
from typstpresenter.convert.pptx_inherit import (
    resolve_bullet,
    resolve_font_size_pt,
    resolve_space_before,
)
from typstpresenter.convert.textbody import (
    PPT_LINE_PITCH_EM,
    autofit_scales,
    paragraph_run_chunks,
    typst_align,
)
from typstpresenter.verify.geometry import EMU_PER_PT, BBox

# Placeholder types that repeat on every slide (theme chrome, not content).
_CHROME_PH_TYPES = {"SLIDE_NUMBER", "FOOTER", "DATE"}

# Page margins of the emitted deck (shared with the emitter's config-page).
PAGE_MARGIN_X = 30.0

# Diagram cluster backend: "svg" (default since 2026-07-19, user-approved
# after the S4 scoreboard: S/F gates equal or better than CeTZ, exact
# geometry for all 187 presets) translates each cluster to an SVG file
# embedded via #image (package typstpresenter.diagram2svg), falling back
# to CeTZ for content SVG cannot represent yet (tables, hyperlink text,
# unembeddable images).  "cetz" draws CeTZ canvases directly.  Override
# via env TP_DIAGRAM_BACKEND or set flow.DIAGRAM_BACKEND programmatically.
import os as _os

DIAGRAM_BACKEND = _os.environ.get("TP_DIAGRAM_BACKEND", "svg")
PAGE_MARGIN_TOP = 24.0
PAGE_MARGIN_BOTTOM = 30.0

# Characters that are markup-active *anywhere* in Typst markup mode. Unlike
# the probed emitter we do not escape '-', '+', '=', '/' mid-line: readable
# text beats protection against rare smart-dash substitutions.
_ESCAPE_MIN = str.maketrans({c: f"\\{c}" for c in '\\#$*_`@<>[]'})
# ... but at the start of a line (or content block) they would start a
# list/enum/heading/term item; '1.' would start a numbered enum.
_LINE_START_RE = re.compile(r"^(\s*)(?:([-+=/])|(\d+)\.)")


def escape_flow(text: str) -> str:
    # \x0b (pptx in-paragraph break) becomes a hard break like a:br;
    # '//' would start a line comment and swallow the closing bracket
    text = text.replace("\x0b", "\n")
    return re.sub("//", r"/\\/", text.translate(_ESCAPE_MIN))


def guard_markup_start(line: str) -> str:
    """Escape list/enum/heading/term markers at line or content-block start."""
    m = _LINE_START_RE.match(line)
    if not m:
        return line
    if m.group(2):
        return f"{m.group(1)}\\{line[m.end(1):]}"
    return f"{m.group(1)}{m.group(3)}\\.{line[m.end():]}"


def _round_size(size: float) -> float:
    return round(size * 2) / 2


def _ph_type_name(shape) -> str | None:
    try:
        ph_type = shape.placeholder_format.type
    except (AttributeError, ValueError):
        return None
    return ph_type.name if ph_type is not None else None


@dataclass
class FlowBlock:
    """One flow element of a slide, still carrying its source geometry."""
    bbox: BBox
    kind: str          # text | image | table | diagram
    markup: str


# ------------------------------------------------------------ inline text --

# Colors that recur across the IBN lecture-slide corpus by a wide margin
# (measured over every run's character count); text in one of these gets a
# short named helper (`styles.typ`, imported once per output directory)
# instead of a fresh `#text(fill: rgb(...))` on every occurrence. One-off
# colors keep the literal `#text(fill: ...)` call -- not worth a name.
COLOR_PALETTE: dict[str, str] = {
    "FF0000": "red",
    "0070C0": "blue",
    "424456": "slate",
    "00B0F0": "cyan",
    "00B050": "green",
    "595959": "gray",
    "7030A0": "purple",
}


def _inline_chunk(text: str, style: tuple, context_size: float,
                  boundary_ok: bool) -> tuple[str, bool]:
    """Render one styled chunk; also report whether the result ends in a
    hash expression (whose parse would continue over a following '(' etc.)."""
    from typstpresenter.convert.textbody import typst_str

    size, bold, italic, underline, rgb, link, highlight = style

    def _linked(markup: str, ends_hash: bool) -> tuple[str, bool]:
        if highlight:
            markup = (f'#highlight(fill: rgb("#{highlight}"))'
                      f"[{guard_markup_start(markup)}]")
            ends_hash = True
        if link:
            return (f"#link({typst_str(link)})"
                    f"[{guard_markup_start(markup)}]", True)
        return markup, ends_hash

    inner = escape_flow(text).replace("\n", " \\ ")
    wrapped_hash = False
    if underline:
        inner = f"#underline[{guard_markup_start(inner)}]"
        wrapped_hash = True
    size_override = abs(size - context_size) > 0.26
    palette_fn = COLOR_PALETTE.get(rgb.upper()) if rgb and not size_override else None
    if not palette_fn and (size_override or rgb):
        args = []
        if size_override:
            args.append(f"size: {_round_size(size):g}pt")
        if rgb:
            args.append(f'fill: rgb("#{rgb}")')
        if bold:
            args.append('weight: "bold"')
        if italic:
            args.append('style: "italic"')
        return _linked(f"#text({', '.join(args)})[{guard_markup_start(inner)}]",
                       True)
    # pure bold/italic (optionally palette-colored) use native markup --
    # but only at word boundaries, otherwise the delimiters do not trigger.
    # A leading '/' right after the opening '*' reads as a block-comment
    # close ('*/') to Typst's tokenizer, so that case also falls back.
    star_ok = boundary_ok and not inner.startswith("/")
    if bold and italic:
        if star_ok:
            core, ends_hash = f"*_{inner}_*", False
        else:
            core, ends_hash = f"#strong[#emph[{guard_markup_start(inner)}]]", True
        if palette_fn:
            return _linked(f"#{palette_fn}[{guard_markup_start(core)}]", True)
        return _linked(core, ends_hash)
    if bold:
        if star_ok:
            core, ends_hash = f"*{inner}*", False
        else:
            core, ends_hash = f"#strong[{guard_markup_start(inner)}]", True
        if palette_fn:
            return _linked(f"#{palette_fn}[{guard_markup_start(core)}]", True)
        return _linked(core, ends_hash)
    if italic:
        if boundary_ok:
            core, ends_hash = f"_{inner}_", False
        else:
            core, ends_hash = f"#emph[{guard_markup_start(inner)}]", True
        if palette_fn:
            return _linked(f"#{palette_fn}[{guard_markup_start(core)}]", True)
        return _linked(core, ends_hash)
    if palette_fn:
        return _linked(f"#{palette_fn}[{guard_markup_start(inner)}]", True)
    return _linked(inner, wrapped_hash)


def _continues_code(text: str) -> bool:
    """True if `text` directly after a hash expression would extend its
    parse: a call chain '(', a content argument '[', or a field access '.'."""
    if not text:
        return False
    if text[0] in "([":
        return True
    return text[0] == "." and len(text) > 1 and (text[1].isalpha() or text[1] == "_")


def _paragraph_inline(chunks: list[tuple[str, tuple]], context_size: float) -> str:
    parts = []
    hash_end = False
    for i, (text, style) in enumerate(chunks):
        prev = chunks[i - 1][0][-1:] if i > 0 else " "
        nxt = chunks[i + 1][0][:1] if i + 1 < len(chunks) else " "
        boundary_ok = (
            not text[:1].isspace()
            and (not prev or not prev.isalnum())
            and (not nxt or not nxt.isalnum())
        )
        rendered, ends_hash = _inline_chunk(text, style, context_size, boundary_ok)
        if hash_end and _continues_code(rendered):
            # ';' ends the previous hash expression, renders as nothing
            parts.append(";")
        parts.append(rendered)
        hash_end = ends_hash
    return "".join(parts)


# ------------------------------------------------------------- text blocks --

def _paragraph_size(paragraph, shape, default_size: float, scale: float) -> float:
    run = paragraph.runs[0] if paragraph.runs else None
    resolved = resolve_font_size_pt(run, paragraph, shape)
    return (resolved if resolved is not None else default_size) * scale


def text_chunks_of_shape(shape, default_size: float) -> list[tuple[str, tuple]]:
    """All (text, style) chunks of a text frame (autofit scale applied)."""
    font_scale, _ = autofit_scales(shape)
    chunks: list[tuple[str, tuple]] = []
    for paragraph in shape.text_frame.paragraphs:
        chunks += paragraph_run_chunks(paragraph, shape, default_size, font_scale)
    return chunks


def dominant_size(chunk_lists: list[list[tuple[str, tuple]]],
                  fallback: float) -> float:
    """The font size covering the most characters, rounded to 0.5pt."""
    weights: Counter[float] = Counter()
    for chunks in chunk_lists:
        for text, style in chunks:
            weights[_round_size(style[0])] += len(text)
    if not weights:
        return fallback
    return weights.most_common(1)[0][0]


def _render_text_block(shape, default_size: float, context_size: float,
                       scale: float = 1.0) -> str:
    """A text frame as flowing markup: native lists, plain paragraphs."""
    font_scale, lnspc_red = autofit_scales(shape)
    font_scale *= scale
    lines: list[str] = []
    in_list = False
    for paragraph in shape.text_frame.paragraphs:
        chunks = paragraph_run_chunks(paragraph, shape, default_size, font_scale,
                                      breaks=True)
        if not chunks:
            in_list = False
            continue
        inline = _paragraph_inline(chunks, context_size)
        bullet = resolve_bullet(paragraph, shape)
        level = paragraph.level
        align = typst_align(paragraph, shape)
        if bullet and not align:
            marker = "+" if bullet == "1." else "-"
            if not in_list and lines:
                lines.append("")
            lines.append(f"{'  ' * level}{marker} {guard_markup_start(inline)}")
            in_list = True
        else:
            if lines:
                # plain (non-list) paragraphs keep the source's own
                # before-spacing instead of a fixed blank-line gap, so a
                # stack of short paragraphs (e.g. one number per line)
                # lines up with a neighboring bulleted column instead of
                # drifting apart under Typst's more generous default
                # paragraph spacing
                spc_bef = resolve_space_before(paragraph, shape)
                gap = None
                if spc_bef:
                    kind, value = spc_bef
                    para_size = _paragraph_size(paragraph, shape, default_size,
                                                font_scale)
                    if kind == "pct":
                        gap = value * PPT_LINE_PITCH_EM * para_size
                    else:
                        gap = value * font_scale * (1.0 - lnspc_red)
                lines.append(f"#v({gap:.2f}pt)" if gap is not None else "")
            line = guard_markup_start(inline)
            if align:
                line = f"#align({align})[{line}]"
            lines.append(line)
            in_list = False
    return "\n".join(lines)


def _paragraph_cell(paragraph, shape, default_size: float,
                    context_size: float, font_scale: float) -> str:
    """One paragraph rendered standalone (e.g. as a table cell)."""
    chunks = paragraph_run_chunks(paragraph, shape, default_size, font_scale,
                                  breaks=True)
    if not chunks:
        return ""
    inline = _paragraph_inline(chunks, context_size)
    line = guard_markup_start(inline)
    bullet = resolve_bullet(paragraph, shape)
    if bullet:
        marker = "+" if bullet == "1." else "-"
        line = f"{marker} {line}"
    return line


def _paired_value_table(shape_a, bbox_a: BBox, shape_b, bbox_b: BBox,
                        default_size: float, context_size: float,
                        scale: float) -> str | None:
    """Merge a label column and a narrow value column into one borderless
    table so each value sits on the same row as its label.

    Matches a wide text block A (labels, possibly preceded by header
    lines) beside a narrow block B of short single-line values whose
    source y positions line up row by row with A's trailing paragraphs.
    Rendering the two as independent grid columns cannot keep those rows
    aligned -- Typst's list pitch differs from PowerPoint's -- so a
    shared table is the only robust encoding.
    """
    if bbox_b.w >= 0.6 * bbox_a.w or bbox_b.x < bbox_a.x + 0.3 * bbox_a.w:
        return None
    if _vertical_overlap(bbox_a, bbox_b) < 0.5:
        return None
    paras_a = [p for p in shape_a.text_frame.paragraphs if p.text.strip()]
    paras_b = [p for p in shape_b.text_frame.paragraphs if p.text.strip()]
    n, m = len(paras_b), len(paras_a)
    if n < 3 or m < n:
        return None
    if any(len(p.text.strip()) > 30 for p in paras_b):
        return None  # long values would wrap: not a value column
    lead = m - n
    # B's top edge must sit where A's first paired label is expected;
    # otherwise the rows do not actually correspond
    positions = _estimate_paragraph_positions(shape_a, default_size, scale)
    nonempty = [i for i, p in enumerate(shape_a.text_frame.paragraphs)
                if p.text.strip()]
    est_lead = positions[nonempty[lead]][0]
    dy = bbox_b.y - bbox_a.y
    if abs(dy - est_lead) > 30:
        return None
    font_a, _ = autofit_scales(shape_a)
    font_b, _ = autofit_scales(shape_b)
    font_a *= scale
    font_b *= scale
    head = "\n\n".join(
        _paragraph_cell(p, shape_a, default_size, context_size, font_a)
        for p in paras_a[:lead])
    rows = []
    for pa, pb in zip(paras_a[lead:], paras_b):
        ca = _paragraph_cell(pa, shape_a, default_size, context_size, font_a)
        cb = _paragraph_cell(pb, shape_b, default_size, context_size, font_b)
        rows.append(f"  [{ca}], [{cb}],")
    label_w = max(bbox_b.x - bbox_a.x, 1.0)
    ratio = bbox_b.w / label_w
    lines = [head, ""] if head else []
    lines.append(f"#table(\n  columns: (1fr, {ratio:.2g}fr),\n"
                 "  stroke: none, inset: (x: 0pt, y: 2.5pt),")
    lines += rows
    lines.append(")")
    return "\n".join(lines)


def _estimate_paragraph_positions(
    shape, default_size: float, scale: float = 1.0,
) -> list[tuple[float, float]]:
    """Estimate the (y_start, y_end) of each paragraph within its text shape.

    Uses font size and PowerPoint's typical 1.22 em line pitch to compute
    how much vertical space each paragraph occupies. Returns coordinates
    relative to the shape's top edge (bbox.y must be added for absolute y).
    """
    from typstpresenter.convert.pptx_inherit import resolve_font_size_pt

    font_scale, _ = autofit_scales(shape)
    font_scale *= scale
    positions: list[tuple[float, float]] = []
    cursor = 0.0
    for paragraph in shape.text_frame.paragraphs:
        run = paragraph.runs[0] if paragraph.runs else None
        size = resolve_font_size_pt(run, paragraph, shape)
        size = (size if size is not None else default_size) * font_scale
        text = paragraph.text.strip()
        if not text:
            # empty paragraphs still consume about half a line
            positions.append((cursor, cursor + size * 0.6))
            cursor += size * 0.6
            continue
        # rough estimate: chars per line based on half-em average glyph width
        bbox_w = (shape.width or 1) / EMU_PER_PT
        chars_per_line = max(bbox_w / (0.5 * size), 1.0)
        n_lines = max(1, -(-len(text) // int(chars_per_line)))  # ceil division
        line_h = size * 1.22
        para_h = n_lines * line_h
        positions.append((cursor, cursor + para_h))
        cursor += para_h
    return positions


def _render_text_segments(
    shape, default_size: float, context_size: float,
    image_bboxes: list[BBox], shape_bbox: BBox,
    scale: float = 1.0,
) -> list[tuple[BBox, str]]:
    """Split a text shape into segments aligned with stacked images.

    Each image defines a row; paragraphs are assigned to the image whose
    y-position is closest. The returned sub-bboxes match the image heights
    so that the column-grouping logic pairs each image with its text segment.

    *image_bboxes* must be sorted by y and contain at least 2 entries.
    """
    para_positions = _estimate_paragraph_positions(shape, default_size, scale)
    paragraphs = list(shape.text_frame.paragraphs)
    n_seg = len(image_bboxes)
    boundaries = [ib.y for ib in image_bboxes]

    # assign each paragraph to a segment based on its absolute y midpoint
    assignments: list[int] = []
    for i, (y0, y1) in enumerate(para_positions):
        abs_mid = shape_bbox.y + (y0 + y1) / 2
        seg = n_seg - 1
        for k in range(n_seg - 1):
            if abs_mid < boundaries[k + 1]:
                seg = k
                break
        assignments.append(seg)

    # render each segment's paragraphs as a separate text block
    font_scale, _ = autofit_scales(shape)
    font_scale *= scale
    segments: list[tuple[BBox, str]] = []
    for seg_idx in range(n_seg):
        seg_paras = [paragraphs[i] for i, s in enumerate(assignments) if s == seg_idx]
        if not seg_paras:
            continue
        lines: list[str] = []
        in_list = False
        for paragraph in seg_paras:
            chunks = paragraph_run_chunks(paragraph, shape, default_size,
                                          font_scale, breaks=True)
            if not chunks:
                in_list = False
                continue
            inline = _paragraph_inline(chunks, context_size)
            bullet = resolve_bullet(paragraph, shape)
            level = paragraph.level
            align = typst_align(paragraph, shape)
            if bullet and not align:
                marker = "+" if bullet == "1." else "-"
                if not in_list and lines:
                    lines.append("")
                lines.append(f"{'  ' * level}{marker} {guard_markup_start(inline)}")
                in_list = True
            else:
                if lines:
                    lines.append("")
                line = guard_markup_start(inline)
                if align:
                    line = f"#align({align})[{line}]"
                lines.append(line)
                in_list = False
        if not lines:
            continue
        # sub-bbox matches the corresponding image's y and height so the
        # column grouping sees proper row pairs with matching vertical extent
        ib = image_bboxes[seg_idx]
        sub_bbox = BBox(shape_bbox.x, ib.y, shape_bbox.w, ib.h)
        segments.append((sub_bbox, "\n".join(lines)))
    return segments


# ---------------------------------------------------------------- pictures --

_TYPST_IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "svg", "webp"}


def _write_image_file(shape, eid: str, media_dir: Path) -> str | None:
    """Write the picture blob (converting if needed); returns the filename."""
    from typstpresenter.verify.pptx_geometry import picture_image

    image = picture_image(shape)
    if image is None:
        return None
    ext = image.ext.lower()
    media_dir.mkdir(parents=True, exist_ok=True)
    if ext in _TYPST_IMAGE_EXTS:
        filename = f"{eid}.{ext}"
        (media_dir / filename).write_bytes(image.blob)
        return filename
    try:
        import io

        from PIL import Image as PILImage

        with PILImage.open(io.BytesIO(image.blob)) as im:
            filename = f"{eid}.png"
            im.convert("RGBA").save(media_dir / filename)
        return filename
    except Exception:
        return None  # not renderable; drop rather than draw a broken frame


def _render_picture(shape, eid: str, bbox: BBox, media_dir: Path,
                    page_w: float, scale: float = 1.0) -> str:
    filename = _write_image_file(shape, eid, media_dir)
    if filename is None:
        return ""
    img = f'image("{media_dir.name}/{filename}", width: {bbox.w * scale:.0f}pt)'
    if abs(bbox.center[0] - page_w / 2) < page_w * 0.06:
        return f"#align(center, {img})"
    if bbox.center[0] > page_w * 0.6:
        return f"#align(right, {img})"
    return f"#{img}"


# ------------------------------------------------------------------ tables --

def _render_table(shape, default_size: float, context_size: float,
                  scale: float = 1.0) -> str:
    table = shape.table
    columns = ", ".join(
        f"{c.width * scale / EMU_PER_PT:.0f}pt" for c in table.columns)
    # dense tables get the source row heights (they fit the slide by
    # construction); auto rows grow with typst's line metrics and spill
    # dense tables off the page. Small tables stay editable with auto rows.
    rows = None
    if len(table.rows) > 8:
        rows = ", ".join(
            f"{r.height * scale / EMU_PER_PT:.0f}pt" for r in table.rows)
    cells = []
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            paras = []
            for paragraph in cell.text_frame.paragraphs:
                chunks = paragraph_run_chunks(paragraph, shape, default_size,
                                              scale, breaks=True)
                if chunks:
                    paras.append(_paragraph_inline(chunks, context_size))
            joined = " \\ ".join(paras)
            row_cells.append(f"[{guard_markup_start(joined)}]")
        cells.append("  " + ", ".join(row_cells) + ",")
    inset = 5.0 * scale if rows is None else min(5.0 * scale, 2.5)
    rows_cfg = f"  rows: ({rows}),\n" if rows else ""
    return (f"#table(\n  columns: ({columns}),\n{rows_cfg}"
            f"  inset: {inset:g}pt, stroke: 0.5pt,\n"
            + "\n".join(cells) + "\n)")


# ---------------------------------------------------------------- diagrams --

def _render_diagram_cluster(cluster: list[tuple], absorbed: list[tuple],
                            page_w: float, default_size: float,
                            context_size: float, media_dir: Path,
                            scale: float = 1.0,
                            slide_index: int = 0,
                            cluster_index: int = 0,
                            canvas_markers: bool = False) -> tuple[BBox, str]:
    """Diagram shapes of a slide, plus text/picture shapes lying in the
    diagram area ("absorbed"), as one flow-embedded CeTZ canvas.

    Absorbed shapes keep their original relative positions: scattered
    labels, hand-drawn tables and small icons are diagram annotations --
    stacking them as flowing paragraphs would garble the picture and blow
    up the slide height.
    """
    if DIAGRAM_BACKEND == "svg":
        from typstpresenter.diagram2svg.cluster import render_cluster_svg

        rendered = render_cluster_svg(
            cluster, absorbed, page_w, default_size, context_size, media_dir,
            scale, slide_index, cluster_index, canvas_markers,
            page_margin_x=PAGE_MARGIN_X)
        if rendered is not None:
            return rendered

    from typstpresenter.convert.cetz import _shape_rotation_deg, effective_bbox

    boxes = ([effective_bbox(shape, b) for shape, b in cluster]
             + [effective_bbox(sh, b) if k == "text" else b
                for k, sh, b, _ in absorbed])
    min_x = min(b.x for b in boxes)
    min_y = min(b.y for b in boxes)
    max_x = max(b.x2 for b in boxes)
    max_y = max(b.y2 for b in boxes)
    bounds = BBox(min_x, min_y, max_x - min_x, max_y - min_y)
    lines = [
        "cetz.canvas(length: 1pt, {",
        "  import cetz.draw: *",
        "  set-style(stroke: 0.75pt)",
        # anchor the canvas so labels sticking out do not shift the drawing
        f"  rect((0, 0), ({bounds.w:.2f}, {-bounds.h:.2f}), stroke: none)",
    ]
    if canvas_markers:
        available_w = page_w - 2 * PAGE_MARGIN_X
        fit_scale = min(available_w / bounds.w, 1.0) if bounds.w > 0 else 1.0
        total_scale = scale * fit_scale
        from typstpresenter.verify.pptx_geometry import element_id
        shape_ids = [element_id(slide_index, shape.shape_id) for shape, _ in cluster]
        for kind, sh, _, eid in absorbed:
            shape_ids.append(eid)
        shapes_str = ", ".join(f'"{sid}"' for sid in shape_ids)
        lines.append(
            f'  content((0, 0), [ #context metadata((kind: "canvas", id: "s{slide_index}-c{cluster_index}", '
            f'page: here().position().page, x: here().position().x.pt(), y: here().position().y.pt(), '
            f'w: {bounds.w * total_scale:.2f}, h: {bounds.h * total_scale:.2f}, shapes: ({shapes_str},)))<tp-canvas> ])'
        )
    # emit in source document order so overlays keep their z-order (an
    # annotation rectangle drawn on a screenshot must paint after it)
    entries = [("shape", shape, bbox, None) for shape, bbox in cluster]
    entries += list(absorbed)
    root = entries[0][1]._element.getroottree().getroot()
    z_order = {el: i for i, el in enumerate(root.iter())}
    entries.sort(key=lambda e: z_order.get(e[1]._element, 0))
    for kind, shape, bbox, eid in entries:
        if kind == "shape":
            rebased = replace(bbox, x=bbox.x - min_x, y=bbox.y - min_y)
            markup, _ = emit_cetz_shape(shape, "", rebased, probes=False,
                                        default_size=default_size)
            lines.append(markup)
            continue
        x, y = bbox.x - min_x, bbox.y - min_y
        if kind == "image":
            filename = _write_image_file(shape, eid, media_dir)
            if filename:
                lines.append(
                    f'  content(({x:.1f}, {-y:.1f}), anchor: "north-west", '
                    f'image("{media_dir.name}/{filename}", width: {bbox.w:.0f}pt))')
        elif kind == "table":
            markup = _render_table(shape, default_size, context_size)
            lines.append(
                f'  content(({x:.1f}, {-y:.1f}), anchor: "north-west", '
                f"box(width: {bbox.w:.1f}pt)[\n{markup}\n  ])")
        else:
            markup = _render_text_block(shape, default_size, context_size)
            if markup:
                # absorbed text must stay inside its source box: match
                # PowerPoint's line pitch instead of typst's airier default
                boxed = (f"box(width: {bbox.w:.1f}pt)[\n"
                         "#set par(leading: 0.59em, spacing: 0.59em)\n"
                         f"{markup}\n  ])")
                angle = _shape_rotation_deg(shape)
                if angle:
                    # rotated label: place by center, same angle convention
                    # as emit_cetz_shape (negate: PPTX is CW, canvas y flips)
                    ccx, ccy = x + bbox.w / 2, y + bbox.h / 2
                    cetz_angle = (-angle) % 360.0
                    lines.append(
                        f"  content(({ccx:.1f}, {-ccy:.1f}), "
                        f"angle: {cetz_angle:.0f}deg, {boxed}")
                else:
                    lines.append(
                        f'  content(({x:.1f}, {-y:.1f}), anchor: "north-west", '
                        f"{boxed}")
    lines.append("})")
    canvas = "\n".join(lines)
    # a canvas wider than the page content area would overflow regardless
    # of the slide-level calibration scale (e.g. a wide absorbed paragraph
    # next to shapes spread across most of the slide's width): shrink it
    # to fit as a hard floor, on top of whatever calibration already asked
    # for
    available_w = page_w - 2 * PAGE_MARGIN_X
    fit_scale = min(available_w / bounds.w, 1.0) if bounds.w > 0 else 1.0
    total_scale = scale * fit_scale
    if total_scale != 1.0:
        # uniform visual shrink (text included) for overflow calibration
        canvas = f"scale({total_scale * 100:.0f}%, reflow: true, {canvas})"
    # alignment is decided from the shapes' true source position (not the
    # canvas's own, possibly fitted, on-page footprint), but the returned
    # bounds must reflect the actual rendered size for callers doing
    # column/grid layout with it
    align_cx = bounds.center[0]
    fitted_bounds = replace(bounds, w=bounds.w * fit_scale, h=bounds.h * fit_scale)
    if abs(align_cx - page_w / 2) < page_w * 0.08:
        return fitted_bounds, f"#align(center, block({canvas}))"
    if align_cx > page_w * 0.6:
        return fitted_bounds, f"#align(right, block({canvas}))"
    return fitted_bounds, f"#{canvas}"


def _cluster_centers(coords: list[float], tolerance: float = 15.0) -> list[float]:
    """Cluster 1D coordinates by tolerance, returning sorted cluster means."""
    if not coords:
        return []
    coords = sorted(coords)
    clusters = [[coords[0]]]
    for c in coords[1:]:
        if c - clusters[-1][-1] <= tolerance:
            clusters[-1].append(c)
        else:
            clusters.append([c])
    return [sum(cl) / len(cl) for cl in clusters]


def _grid_index(coord: float, centers: list[float]) -> int:
    return min(range(len(centers)), key=lambda i: abs(centers[i] - coord))


def _spacing(centers: list[float], extents: dict[int, float]) -> float:
    if len(centers) <= 1:
        return 50.0
    gaps = [centers[i + 1] - centers[i] for i in range(len(centers) - 1)]
    return sum(gaps) / len(gaps)


def _connector_endpoints(shape, bbox: BBox) -> tuple[tuple[float, float], tuple[float, float]]:
    from typstpresenter.verify.geometry import EMU_PER_PT
    from typstpresenter.convert.cetz import _shape_flips
    try:
        bx = shape.begin_x / EMU_PER_PT
        by = shape.begin_y / EMU_PER_PT
        ex = shape.end_x / EMU_PER_PT
        ey = shape.end_y / EMU_PER_PT
        return (bx, by), (ex, ey)
    except Exception:
        pass
    flip_h, flip_v = _shape_flips(shape)
    x1, y1, x2, y2 = bbox.x, bbox.y, bbox.x2, bbox.y2
    bx, ex = (x2, x1) if flip_h else (x1, x2)
    by, ey = (y2, y1) if flip_v else (y1, y2)
    return (bx, by), (ex, ey)


def _detect_fletcher_diagram(diagram_shapes: list[tuple], absorbed: list[tuple],
                             page_w: float, default_size: float, context_size: float,
                             media_dir: Path, scale: float, slide_index: int,
                             canvas_markers: bool) -> FlowBlock | None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
    from typstpresenter.verify.geometry import BBox
    from typstpresenter.convert.cetz import _cetz_label_markup, _shape_flips
    from typstpresenter.convert.pptx_style import shape_fill_rgb, shape_line_rgb, shape_line_width_pt

    nodes = []
    connectors = []
    
    # Separate nodes and connectors in diagram_shapes
    for shape, bbox in diagram_shapes:
        is_conn = (shape.shape_type == MSO_SHAPE_TYPE.LINE
                   or shape.element.tag.endswith("}cxnSp"))
        if is_conn:
            connectors.append((shape, bbox))
        else:
            nodes.append((shape, bbox))
            
    # And text blocks in absorbed
    has_image_or_table = False
    for kind, shape, bbox, eid in absorbed:
        if kind in ("image", "table"):
            has_image_or_table = True
        elif kind == "text":
            nodes.append((shape, bbox))
            
    if has_image_or_table or len(nodes) < 2 or len(connectors) < 1:
        return None
        
    # Flowchart node heuristics to prevent misclassifying complex slide layouts
    node_heights = [b.h for _, b in nodes]
    node_widths = [b.w for _, b in nodes]
    if max(node_heights) >= 150.0 or max(node_widths) >= 250.0:
        return None
    if max(node_heights) / min(node_heights) >= 1.8:
        return None
        
    # Cluster coordinate centers
    col_centers = _cluster_centers([b.center[0] for _, b in nodes])
    row_centers = _cluster_centers([b.center[1] for _, b in nodes])
    
    if len(col_centers) < 1 or len(row_centers) < 1:
        return None
        
    # Map each node to its grid cell
    cell_of = {}
    cells = set()
    for shape, bbox in nodes:
        col = _grid_index(bbox.center[0], col_centers)
        row = _grid_index(bbox.center[1], row_centers)
        if (col, row) in cells:
            return None  # Overlap, fallback to CeTZ
        cells.add((col, row))
        cell_of[shape.shape_id] = (col, row)
        
    # Calculate spacing
    col_widths = {}
    row_heights = {}
    for shape, bbox in nodes:
        col, row = cell_of[shape.shape_id]
        col_widths[col] = max(col_widths.get(col, 0.0), bbox.w)
        row_heights[row] = max(row_heights.get(row, 0.0), bbox.h)
    
    spacing_x = _spacing(col_centers, col_widths)
    spacing_y = _spacing(row_centers, row_heights)
    
    # Map connectors to edges
    edges = []
    for shape, bbox in connectors:
        # get endpoints
        pts = _connector_endpoints(shape, bbox)
        if pts is None:
            continue
        (bx, by), (ex, ey) = pts
        
        # nearest nodes by center distance
        node_a = min(nodes, key=lambda n: (n[1].center[0] - bx)**2 + (n[1].center[1] - by)**2)
        node_b = min(nodes, key=lambda n: (n[1].center[0] - ex)**2 + (n[1].center[1] - ey)**2)
        
        if node_a[0] == node_b[0]:
            continue
            
        (col_a, row_a) = cell_of[node_a[0].shape_id]
        (col_b, row_b) = cell_of[node_b[0].shape_id]
        
        corner = None
        if col_a != col_b and row_a != row_b:
            leaves_horizontally = abs(bx - node_a[1].center[0]) >= node_a[1].w / 2 - 2
            corner = (col_b, row_a) if leaves_horizontally else (col_a, row_b)
            
        edges.append((col_a, row_a, corner, col_b, row_b))
        
    # Emit Fletcher markup
    lines = [
        "fletcher.diagram(",
        "  node-inset: 0pt,",
        f"  spacing: ({spacing_x * scale:.1f}pt, {spacing_y * scale:.1f}pt),",
        "  node-stroke: 0.75pt,"
    ]
    
    # Add nodes
    for shape, bbox in nodes:
        col, row = cell_of[shape.shape_id]
        is_diamond = False
        is_oval = False
        try:
            if hasattr(shape, "auto_shape_type"):
                if shape.auto_shape_type == MSO_SHAPE.DIAMOND:
                    is_diamond = True
                elif shape.auto_shape_type == MSO_SHAPE.OVAL:
                    is_oval = True
        except ValueError:
            pass
            
        if is_diamond:
            shape_arg = "shape: fletcher.shapes.diamond.with(fit: 1)"
        elif is_oval:
            shape_arg = "shape: fletcher.shapes.circle"
        else:
            shape_arg = "corner-radius: 4pt"
        w, h = (bbox.w / 2, bbox.h / 2) if is_diamond else (bbox.w, bbox.h)
        
        # Label text
        label_w = max(bbox.w - 7.2, 7.2)
        label = _cetz_label_markup(shape, "", probes=False, default_size=default_size, label_w=label_w)
        
        # Styling
        fill_rgb = shape_fill_rgb(shape)
        stroke_rgb = shape_line_rgb(shape) or "000000"
        thickness = shape_line_width_pt(shape)
        
        style_args = [shape_arg]
        if fill_rgb:
            style_args.append(f'fill: rgb("#{fill_rgb}")')
        if stroke_rgb:
            style_args.append(f'stroke: (paint: rgb("#{stroke_rgb}"), thickness: {thickness:g}pt)')
            
        lines.append(
            f"  node(({col}, {row}), [{label}], width: {w * scale:.2f}pt, height: {h * scale:.2f}pt, {', '.join(style_args)}),"
        )
        
    # Add edges
    for col_a, row_a, corner, col_b, row_b in edges:
        via = f"({corner[0]}, {corner[1]}), " if corner else ""
        lines.append(
            f'  edge(({col_a}, {row_a}), {via}({col_b}, {row_b}), "-|>"),'
        )
        
    lines.append(")")
    markup = f"#align(center,\n" + "\n".join(f"  {line}" for line in lines) + "\n)"
    
    # Union bbox for the flow block
    bounds = _union_bbox([b for _, b in nodes] + [b for _, b in connectors])
    return FlowBlock(bounds, "fletcher", markup)


# ------------------------------------------------------------ slide markup --

def _is_chrome(shape, page_h: float, bbox: BBox,
               page_w: float = 1e9) -> bool:
    """Slide chrome (page number, footer, date) that the theme should own."""
    if _ph_type_name(shape) in _CHROME_PH_TYPES:
        return True
    # unnamed text boxes hugging the bottom edge with tiny text are footers
    if (shape.has_text_frame and bbox.y > page_h * 0.92
            and len(shape.text_frame.text) < 60):
        return True
    # small ornaments inside the title band (course badges, QR links)
    # decorate every slide's header, they are not content
    if bbox.y2 <= 60 and bbox.h < 50:
        return True
    return False


def _off_page(bbox: BBox, page_w: float, page_h: float) -> bool:
    """Shapes parked outside the slide are working material, not content."""
    return (bbox.x >= page_w - 1 or bbox.y >= page_h - 1
            or bbox.x2 <= 1 or bbox.y2 <= 1)


def _overlap_frac(bbox: BBox, bounds: BBox) -> float:
    """Fraction of `bbox`'s area lying inside `bounds`."""
    area = bbox.w * bbox.h
    if area <= 0:
        return 0.0
    overlap = (max(0.0, min(bbox.x2, bounds.x2) - max(bbox.x, bounds.x))
               * max(0.0, min(bbox.y2, bounds.y2) - max(bbox.y, bounds.y)))
    return overlap / area


def flow_slide_blocks(slide, slide_index: int, page_w: float, page_h: float,
                      media_dir: Path, default_size: float,
                      context_size: float, scale: float = 1.0,
                      canvas_markers: bool = False) -> list[FlowBlock]:
    """Classify and render the content shapes of one slide (title excluded)."""
    from typstpresenter.verify.pptx_geometry import (
        element_id,
        is_picture,
        iter_flat_shapes,
    )

    diagram_shapes: list[tuple] = []
    candidates: list[tuple] = []  # (kind, shape, bbox, eid)
    for shape, bbox in iter_flat_shapes(slide.shapes):
        if shape == slide.shapes.title:
            continue
        if (_is_chrome(shape, page_h, bbox, page_w)
                or _off_page(bbox, page_w, page_h)):
            continue
        eid = element_id(slide_index, shape.shape_id)
        if is_diagram_shape(shape):
            diagram_shapes.append((shape, bbox))
        elif getattr(shape, "has_table", False):
            candidates.append(("table", shape, bbox, eid))
        elif is_picture(shape):
            candidates.append(("image", shape, bbox, eid))
        elif shape.has_text_frame and shape.text_frame.text.strip():
            candidates.append(("text", shape, bbox, eid))

    # Text/pictures lying in the diagram area -- or labels sitting on an
    # image (annotated screenshots) -- are annotations, not flowing prose:
    # absorb them into one canvas that keeps their relative positions.
    absorbed: list[tuple] = []
    sparse_kernel = False
    texts = [c for c in candidates if c[0] == "text"]
    images = [c for c in candidates if c[0] == "image"]

    def _is_small(bbox: BBox) -> bool:
        return bbox.w * bbox.h < 9000

    # labels genuinely sitting ON an image pull that image into the canvas;
    # large prose merely brushing an image stays in flow (rendered text
    # grows taller than its source box and would collide with the image)
    on_image: set[int] = set()
    labeled: set[int] = set()
    for ti, (_, _, tb, _) in enumerate(texts):
        for ii, (_, _, ib, _) in enumerate(images):
            # labels often sit flush against the image edge; a small
            # growth catches them without swallowing separate captions
            grown = BBox(ib.x - 12, ib.y - 12, ib.w + 24, ib.h + 24)
            if _overlap_frac(tb, grown) > (0.3 if _is_small(tb) else 0.75):
                on_image.add(ti)
                labeled.add(ii)
    # an image lying on a larger image (an inset/detail view) must keep its
    # overlay position: both go into the canvas, the larger one as kernel
    inset_images: set[int] = set()
    for ii, (_, _, ib, _) in enumerate(images):
        for jj, (_, _, jb, _) in enumerate(images):
            if (ii != jj and ib.w * ib.h < jb.w * jb.h
                    and _overlap_frac(ib, jb) > 0.5):
                inset_images.add(ii)
                labeled.add(jj)
    # an annotation drawn ON a picture (e.g. a rectangle framing one
    # component of a screenshot) keeps its position only if the picture
    # joins the same canvas
    for ii, (_, _, ib, _) in enumerate(images):
        if any(_overlap_frac(db, ib) > 0.8 for _, db in diagram_shapes):
            labeled.add(ii)
    # an image surrounded by several small labels is diagram art (arrows,
    # letters, callouts placed around a screenshot): keep the arrangement
    for ii, (_, _, ib, _) in enumerate(images):
        near = BBox(ib.x - 40, ib.y - 40, ib.w + 80, ib.h + 80)
        satellites = [ti for ti, (_, _, tb, _) in enumerate(texts)
                      if _is_small(tb) and _overlap_frac(tb, near) > 0.5]
        if len(satellites) >= 2:
            labeled.add(ii)

    def _union(boxes: list[BBox]) -> BBox:
        x = min(b.x for b in boxes)
        y = min(b.y for b in boxes)
        return BBox(x, y, max(b.x2 for b in boxes) - x,
                    max(b.y2 for b in boxes) - y)

    kernel_boxes = [b for _, b in diagram_shapes]
    kernel_boxes += [images[ii][2] for ii in labeled]
    if kernel_boxes:
        d_bounds = _union(kernel_boxes)
        # any-size absorption only against the drawn-shape area; bounds
        # that merely inherit an image must not swallow adjacent prose
        s_bounds = _union([b for _, b in diagram_shapes]) if diagram_shapes else None
        expanded = BBox(d_bounds.x - 30, d_bounds.y - 30,
                        d_bounds.w + 60, d_bounds.h + 60)
        # thin textless decorations (brackets, arrows) spanning the slide
        # must not swallow the text columns they decorate
        kernel_area = sum(b.w * b.h for _, b in diagram_shapes)
        kernel_text = any(
            getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
            for sh, _ in diagram_shapes)
        sparse_kernel = (
            s_bounds is not None and not kernel_text
            and kernel_area < 0.15 * s_bounds.w * s_bounds.h)
        large_text_threshold = 0.45 if sparse_kernel else 0.15

        absorbed = [images[ii] for ii in sorted(labeled)]
        on_image_entries = {id(texts[ti]) for ti in on_image}
        on_image_entries |= {id(images[ii]) for ii in inset_images}
        remaining = []
        for entry in candidates:
            kind, shape, bbox, eid = entry
            if entry in absorbed:
                continue
            if id(entry) in on_image_entries:
                absorbed.append(entry)
                continue
            small = _is_small(bbox)
            # large text blocks (whole columns) are prose, not annotations:
            # a sparse decoration kernel brushing them must not swallow them
            in_area = (
                (s_bounds is not None
                 and _overlap_frac(bbox, s_bounds)
                 > (0.15 if small else large_text_threshold))
                or (small and _overlap_frac(bbox, d_bounds) > 0.15)
            )
            # small captions hugging the diagram belong to it as well;
            # stacking them in flow detaches them from what they label
            nearby = small and _overlap_frac(bbox, expanded) > 0.5
            # row/column labels flanking the diagram (e.g. "A"/"B"/"C" axis
            # captions to its left) barely overlap its bounds by area, but
            # sit right beside it: catch these by vertical containment +
            # horizontal gap instead of an area fraction
            flanking = (
                small
                and min(bbox.y2, d_bounds.y2) - max(bbox.y, d_bounds.y)
                    > 0.6 * bbox.h
                and max(0.0, d_bounds.x - bbox.x2, bbox.x - d_bounds.x2) < 50.0
            )
            # a table drawn inside the diagram area (e.g. a stack layout
            # with a brace and labels around it) is part of the picture:
            # letting it flow would tear it out below the diagram
            absorb_table = (
                kind == "table" and s_bounds is not None
                and _overlap_frac(bbox, s_bounds) > 0.6)
            if absorb_table or (kind != "table"
                                and (in_area or nearby or flanking)):
                absorbed.append(entry)
            else:
                remaining.append(entry)
        candidates = remaining

    blocks: list[FlowBlock] = []
    # detect the "stacked images + tall text" pattern: multiple small images
    # at different y-positions, all in the same x-range, with a single tall
    # text block adjacent to them spanning most of their combined height.
    # Split the text into segments aligned with each image for row pairing.
    split_text_ids: set[str] = set()
    text_candidates = [(i, c) for i, c in enumerate(candidates) if c[0] == "text"]
    image_candidates = [(i, c) for i, c in enumerate(candidates) if c[0] == "image"]
    if len(image_candidates) >= 2 and text_candidates:
        # find groups of images stacked vertically in the same x-range
        img_by_x: dict[int, list[int]] = {}
        for idx, (_, _, ib, _) in image_candidates:
            # bucket by x (quantized to 30pt) to find columns of images
            xq = int(ib.x / 30)
            img_by_x.setdefault(xq, []).append(idx)
        for xq, img_idxs in img_by_x.items():
            if len(img_idxs) < 2:
                continue
            img_entries = [candidates[i] for i in img_idxs]
            img_bboxes = [c[2] for c in img_entries]
            # images must be stacked (non-overlapping y ranges, similar x)
            sorted_by_y = sorted(img_bboxes, key=lambda b: b.y)
            x_spread = max(b.x for b in sorted_by_y) - min(b.x for b in sorted_by_y)
            if x_spread > 50:
                continue  # not a column
            combined_y = sorted_by_y[0].y
            combined_y2 = sorted_by_y[-1].y2
            combined_h = combined_y2 - combined_y
            # find a tall text block next to these images
            for ti, (_, tshape, tbbox, teid) in text_candidates:
                if tbbox.h < combined_h * 0.6:
                    continue  # too short to span the images
                # text must be x-disjoint from the images (side by side)
                img_x_mid = sum(b.center[0] for b in sorted_by_y) / len(sorted_by_y)
                txt_x_mid = tbbox.center[0]
                if abs(img_x_mid - txt_x_mid) < max(sorted_by_y[0].w, tbbox.w) * 0.5:
                    continue  # overlapping x-range, not side-by-side
                # the text must vertically overlap the image stack
                y_overlap = (min(tbbox.y2, combined_y2) - max(tbbox.y, combined_y))
                if y_overlap < combined_h * 0.5:
                    continue
                # pattern matches! split text at image y-positions
                segments = _render_text_segments(
                    tshape, default_size, context_size,
                    sorted_by_y, tbbox, scale)
                if len(segments) >= 2:
                    split_text_ids.add(teid)
                    for sub_bbox, markup in segments:
                        if markup:
                            blocks.append(FlowBlock(sub_bbox, "text", markup))
    # a narrow column of short values whose rows line up with a wider
    # label column merges into a single borderless table (independent
    # grid columns cannot keep such rows aligned)
    paired_ids: set[str] = set()
    text_cands = [c for c in candidates
                  if c[0] == "text" and c[3] not in split_text_ids]
    for a in text_cands:
        for b in text_cands:
            if a[3] == b[3] or a[3] in paired_ids or b[3] in paired_ids:
                continue
            merged = _paired_value_table(a[1], a[2], b[1], b[2],
                                         default_size, context_size, scale)
            if merged:
                blocks.append(
                    FlowBlock(_union_bbox([a[2], b[2]]), "text", merged))
                paired_ids |= {a[3], b[3]}
    for kind, shape, bbox, eid in candidates:
        if eid in split_text_ids or eid in paired_ids:
            continue  # already split or merged above
        if kind == "table":
            markup = _render_table(shape, default_size, context_size, scale)
        elif kind == "image":
            markup = _render_picture(shape, eid, bbox, media_dir, page_w, scale)
        else:
            markup = _render_text_block(shape, default_size, context_size, scale)
        if markup:
            blocks.append(FlowBlock(bbox, kind, markup))
    if diagram_shapes or absorbed:
        fletcher_block = _detect_fletcher_diagram(
            diagram_shapes, absorbed, page_w, default_size, context_size,
            media_dir, scale, slide_index, canvas_markers)
        if fletcher_block is not None:
            blocks.append(fletcher_block)
        else:
            # G5: spatially disjoint diagram groups render as separate
            # clusters instead of one oversized sparse union canvas
            components = _spatial_components(diagram_shapes, absorbed)
            for ci, (comp_shapes, comp_absorbed) in enumerate(components):
                bounds, markup = _render_diagram_cluster(
                    comp_shapes, comp_absorbed, page_w, default_size,
                    context_size, media_dir, scale, slide_index=slide_index,
                    cluster_index=ci, canvas_markers=canvas_markers)
                # a decoration-only canvas (thin brackets/arrows, sparse
                # strips) lying across the text columns it decorates would
                # force those columns to stack; losing the decoration beats
                # halving the slide
                kernel_text = any(
                    getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
                    for sh, _ in comp_shapes)
                thin = bounds.w < 80 or bounds.h < 40
                droppable = ((sparse_kernel or thin) and not kernel_text
                             and not comp_absorbed)
                conflicts = droppable and any(
                    _decoration_conflict(bounds, b.bbox) for b in blocks)
                if not conflicts:
                    blocks.append(FlowBlock(bounds, "diagram", markup))
    # a picture or note box floated in the empty right portion of a wide
    # text box is a side figure: narrow the text's grouping bbox so both
    # form columns instead of the floater being stacked below the text
    for txt_block in blocks:
        if txt_block.kind != "text":
            continue
        tb = txt_block.bbox
        floated = [
            b.bbox.x for b in blocks
            if b is not txt_block
            and _overlap_frac(b.bbox, tb) > 0.7
            and b.bbox.x > tb.x + 0.5 * tb.w
            and ((b.kind == "image" and b.bbox.w >= 100)
                 or (b.kind == "text"
                     and b.bbox.w * b.bbox.h < 0.25 * tb.w * tb.h))
        ]
        if floated:
            new_w = max(min(floated) - 8 - tb.x, tb.w * 0.4)
            txt_block.bbox = replace(tb, w=new_w)

    blocks.sort(key=lambda b: (b.bbox.y, b.bbox.x))
    return _group_columns(blocks, page_w - 2 * PAGE_MARGIN_X, scale=scale)


def _spatial_components(diagram_shapes: list[tuple], absorbed: list[tuple],
                        pad: float = 25.0) -> list[tuple[list, list]]:
    """Partition a slide's diagram content into spatially disjoint groups.

    Two entries connect when their bboxes, inflated by `pad` pt, overlap.
    Components without a diagram shape (stray absorbed labels) and tiny
    single-shape specks merge into the nearest substantial component, so a
    split only happens between genuinely separate drawings (gap G5).
    """
    entries = ([("d", e, e[1]) for e in diagram_shapes]
               + [("a", e, e[2]) for e in absorbed])
    n = len(entries)
    if n == 0:
        return [(diagram_shapes, absorbed)]
    parent = list(range(n))

    def find(i: int) -> int:
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    def near(a: BBox, b: BBox) -> bool:
        return not (a.x2 + pad < b.x or b.x2 + pad < a.x
                    or a.y2 + pad < b.y or b.y2 + pad < a.y)

    for i in range(n):
        for j in range(i + 1, n):
            if near(entries[i][2], entries[j][2]):
                parent[find(i)] = find(j)

    groups: dict[int, list] = {}
    for i, entry in enumerate(entries):
        groups.setdefault(find(i), []).append(entry)
    comps = list(groups.values())
    if len(comps) <= 1:
        return [(diagram_shapes, absorbed)]

    def has_shape(c) -> bool:
        return any(tag == "d" for tag, _, _ in c)

    def substantial(c) -> bool:
        area = sum(b.w * b.h for _, _, b in c)
        return has_shape(c) and (len(c) >= 2 or area >= 2500)

    big = [c for c in comps if substantial(c)]
    if len(big) <= 1:
        return [(diagram_shapes, absorbed)]

    def center(c) -> tuple[float, float]:
        u = _union_bbox([b for _, _, b in c])
        return u.center

    for c in comps:
        if c in big:
            continue
        cx, cy = center(c)
        target = min(big, key=lambda g: (center(g)[0] - cx) ** 2
                     + (center(g)[1] - cy) ** 2)
        target.extend(c)

    big.sort(key=lambda c: (min(b.y for _, _, b in c), min(b.x for _, _, b in c)))
    return [
        ([e for tag, e, _ in c if tag == "d"],
         [e for tag, e, _ in c if tag == "a"])
        for c in big
    ]


def _decoration_conflict(canvas: BBox, other: BBox) -> bool:
    """Does a decoration canvas overlap a flow block enough to break its
    column layout?"""
    x_overlap = min(canvas.x2, other.x2) - max(canvas.x, other.x)
    narrower = min(canvas.w, other.w)
    if narrower <= 0 or x_overlap / narrower < 0.3:
        return False
    return _vertical_overlap(canvas, other) > 0.3


def _vertical_overlap(a: BBox, b: BBox) -> float:
    """Overlap of the y ranges as a fraction of the smaller height."""
    overlap = min(a.y2, b.y2) - max(a.y, b.y)
    smaller = min(a.h, b.h)
    return overlap / smaller if smaller > 0 else 0.0


def _column_disjoint(a: FlowBlock, b: FlowBlock) -> bool:
    # PPT columns routinely overlap by a few points (text boxes are drawn
    # generously, diagram labels overhang); tolerate a sliver relative to
    # the narrower block -- more for diagrams, whose bounds include the
    # empty corners of their bounding box
    factor = 0.3 if "diagram" in (a.kind, b.kind) else 0.12
    slack = max(8.0, factor * min(a.bbox.w, b.bbox.w))
    return a.bbox.x2 <= b.bbox.x + slack or b.bbox.x2 <= a.bbox.x + slack


def _column_of(block: FlowBlock, row: list[FlowBlock]) -> int | None:
    """The single row column whose x-range contains `block`, if any."""
    hits = []
    for k, member in enumerate(row):
        overlap = (min(block.bbox.x2, member.bbox.x2)
                   - max(block.bbox.x, member.bbox.x))
        frac = overlap / block.bbox.w if block.bbox.w > 0 else 0.0
        if frac > 0.5:
            hits.append(k)
        elif frac > 0.2:
            return None  # straddles a column boundary
    return hits[0] if len(hits) == 1 else None


def _union_bbox(boxes: list[BBox]) -> BBox:
    x = min(b.x for b in boxes)
    y = min(b.y for b in boxes)
    return BBox(x, y, max(b.x2 for b in boxes) - x, max(b.y2 for b in boxes) - y)


def _fit_to_cell(block: FlowBlock, cell_pt: float) -> str:
    """Fixed-size content (canvas, picture) wider than its grid cell gets a
    uniform visual shrink so it cannot overflow the cell."""
    if block.kind in ("diagram", "image") and block.bbox.w > cell_pt > 0:
        factor = cell_pt / block.bbox.w
        return f"#scale({factor * 100:.0f}%, reflow: true)[{block.markup}]"
    return block.markup


def _row_markup(row: list[FlowBlock], extras: dict[int, list[FlowBlock]],
                content_w: float, scale: float = 1.0) -> FlowBlock:
    row = sorted(row, key=lambda b: b.bbox.x)
    all_blocks = list(row) + [b for cell in extras.values() for b in cell]
    bounds = _union_bbox([b.bbox for b in all_blocks])
    if not extras and all(b.kind == "image" for b in row):
        # a bare row of pictures: center it as a whole, with the source
        # gaps as gutters, instead of spreading it over fr-columns
        gaps = [row[k + 1].bbox.x - row[k].bbox.x2 for k in range(len(row) - 1)]
        gutter = max(sum(gaps) / len(gaps), 4.0) if gaps else 16.0
        cells = ", ".join(f"[{b.markup}]" for b in row)
        markup = (f"#align(center, grid(\n  columns: {len(row)}, "
                  f"column-gutter: {gutter:.0f}pt,\n  {cells},\n))")
        return FlowBlock(bounds, "grid", markup)
    # a cell must be wide enough for its widest occupant: a canvas or
    # image stacked below a narrow column member (an "extra") would
    # otherwise be fit-shrunk to the narrow member's share
    widths = [
        max([b.bbox.w] + [e.bbox.w for e in extras.get(k, ())] + [1.0])
        for k, b in enumerate(row)
    ]
    base = widths[0]
    total = sum(widths)
    usable = content_w - 16.0 * (len(row) - 1)
    columns = ", ".join(f"{w / base:.2g}fr" for w in widths)
    cells = []
    # a row starting well right of the margin keeps its offset as a fixed
    # leading column, so its horizontal placement matches the source
    offset = bounds.x - PAGE_MARGIN_X
    if offset > 40:
        columns = f"{offset * scale:.0f}pt, " + columns
        cells.append("  [],")
        usable -= offset
    for k, member in enumerate(row):
        cell_pt = usable * widths[k] / total
        parts = []
        # a column that starts deeper than the row's top (e.g. a label
        # line only present in the neighboring column) keeps that offset,
        # otherwise its rows would ride up out of alignment with the
        # column(s) it corresponds to
        dy = member.bbox.y - bounds.y
        if dy > 10:
            parts.append(f"#v({dy * scale:.0f}pt)")
        parts.append(_fit_to_cell(member, cell_pt))
        parts += [_fit_to_cell(b, cell_pt) for b in
                  sorted(extras.get(k, ()), key=lambda b: b.bbox.y)]
        indented = "\n".join(f"    {line}" for line in
                             "\n\n".join(parts).split("\n"))
        cells.append(f"  [\n{indented}\n  ],")
    markup = (f"#grid(\n  columns: ({columns}), column-gutter: 16pt,\n"
              + "\n".join(cells) + "\n)")
    return FlowBlock(bounds, "grid", markup)


def _group_columns(blocks: list[FlowBlock],
                   content_w: float, scale: float = 1.0) -> list[FlowBlock]:
    """Merge side-by-side blocks into one #grid block (columns layout).

    Blocks that share most of their vertical range but occupy disjoint
    horizontal ranges were columns in the original slide; stacking them
    would double the slide height and shuffle the reading order. Blocks
    further down that stay within one column's x-range (a picture below
    the text of its column) belong into that column's cell.
    """
    merged: list[FlowBlock] = []
    i = 0
    while i < len(blocks):
        row = [blocks[i]]
        j = i + 1
        while j < len(blocks):
            cand = blocks[j]
            if (all(_vertical_overlap(cand.bbox, m.bbox) > 0.4 for m in row)
                    and all(_column_disjoint(cand, m) for m in row)):
                row.append(cand)
                j += 1
            else:
                break
        extras: dict[int, list[FlowBlock]] = {}
        if len(row) > 1:
            row.sort(key=lambda b: b.bbox.x)
            row_bottom = max(m.bbox.y2 for m in row)
            while j < len(blocks):
                cand = blocks[j]
                col = _column_of(cand, row)
                if col is None or cand.bbox.y >= row_bottom + 40:
                    break
                # a block that sits entirely below ALL row members belongs
                # to a separate row, not this column's continuation
                if not any(_vertical_overlap(cand.bbox, m.bbox) > 0
                           for m in row):
                    break
                extras.setdefault(col, []).append(cand)
                j += 1
        if len(row) == 1:
            merged.append(row[0])
        else:
            merged.append(_row_markup(row, extras, content_w, scale=scale))
        i = j if j > i + 1 else i + 1
    return merged


def has_center_title(slide) -> bool:
    """True for title-layout slides (their title is a CENTER_TITLE box)."""
    return any(_ph_type_name(shape) == "CENTER_TITLE"
               for shape in slide.shapes)


def is_section_divider(slide, page_w: float, page_h: float) -> bool:
    """A slide whose only content is a title, itself centered on the page
    (horizontally and vertically) -- a "section break" slide built from a
    custom layout rather than PowerPoint's own Section Header placeholder
    type (which `has_center_title` already covers)."""
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    from typstpresenter.convert.pptx_inherit import resolve_alignment, resolve_anchor
    from typstpresenter.verify.pptx_geometry import iter_flat_shapes

    title = slide.shapes.title
    if title is None or not getattr(title, "has_text_frame", False):
        return False
    if not title.text_frame.text.strip():
        return False
    if has_center_title(slide):
        return False
    for _shape, _bbox in _content_text_shapes(slide, page_h, page_w):
        return False
    for shape, bbox in iter_flat_shapes(slide.shapes):
        if (shape == title or _is_chrome(shape, page_h, bbox, page_w)
                or _off_page(bbox, page_w, page_h)):
            continue
        return False
    paragraph = next((p for p in title.text_frame.paragraphs if p.text.strip()),
                      None)
    if paragraph is None:
        return False
    return (resolve_alignment(paragraph, title) == PP_ALIGN.CENTER
            and resolve_anchor(title) == MSO_ANCHOR.MIDDLE)


def slide_title_size(slide, fallback: float = 40.0) -> float:
    """Resolved (inherited) title size of a slide, autofit applied."""
    title = slide.shapes.title
    if title is None or not getattr(title, "has_text_frame", False):
        return fallback
    font_scale, _ = autofit_scales(title)
    sizes = []
    for paragraph in title.text_frame.paragraphs:
        if not paragraph.text.strip():
            continue
        run = paragraph.runs[0] if paragraph.runs else None
        size = resolve_font_size_pt(run, paragraph, title)
        if size is not None:
            sizes.append(size * font_scale)
    return _round_size(max(sizes)) if sizes else fallback


def slide_title_link(slide) -> str | None:
    """A hyperlink carried by the title (run link or shape click action)."""
    from pptx.oxml.ns import qn

    title = slide.shapes.title
    if title is None or not getattr(title, "has_text_frame", False):
        return None
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                if run.hyperlink.address:
                    return run.hyperlink.address
            except (AttributeError, KeyError):
                pass
    for el in title.element.iter(qn("a:hlinkClick")):
        rid = el.get(qn("r:id"))
        if rid:
            try:
                return title.part.rels[rid].target_ref
            except KeyError:
                pass
    return None


def deck_title_size(prs) -> float | None:
    """Most common resolved title size of the content slides (title-layout
    and section-divider slides excluded); sizes the theme's heading
    preamble."""
    page_w = prs.slide_width / EMU_PER_PT
    page_h = prs.slide_height / EMU_PER_PT
    weights: Counter[float] = Counter()
    for slide in prs.slides:
        title = slide.shapes.title
        if title is None or has_center_title(slide):
            continue
        if not getattr(title, "has_text_frame", False):
            continue
        if not title.text_frame.text.strip():
            continue
        if is_section_divider(slide, page_w, page_h):
            continue
        weights[slide_title_size(slide)] += 1
    if not weights:
        return None
    return weights.most_common(1)[0][0]


def _title_trailing_badge(slide) -> tuple[str, str] | None:
    """Split off a highlighted "badge" line (e.g. a link) that the source
    glues under the title with a line break, so it renders as its own line
    below the heading instead of being joined into the same line."""
    from pptx.oxml.ns import qn

    from typstpresenter.convert.textbody import _run_highlight, run_style

    title = slide.shapes.title
    if title is None or not getattr(title, "has_text_frame", False):
        return None
    paragraphs = title.text_frame.paragraphs
    if len(paragraphs) != 1:
        return None
    paragraph = paragraphs[0]
    children = list(paragraph._p)
    br_positions = [i for i, c in enumerate(children) if c.tag == qn("a:br")]
    if not br_positions:
        return None
    split = br_positions[-1]
    run_iter = iter(paragraph.runs)
    before, after = [], []
    for i, child in enumerate(children):
        if child.tag == qn("a:r"):
            run = next(run_iter)
            (before if i < split else after).append(run)
    if not before or not after or not all(_run_highlight(r) for r in after):
        return None
    heading_text = "".join(r.text for r in before).strip()
    trailing_text = "".join(r.text for r in after).strip()
    if not heading_text or not trailing_text:
        return None
    default_size = slide_title_size(slide)
    style = run_style(after[0], paragraph, title, default_size)
    markup, _ = _inline_chunk(trailing_text, style, style[0], True)
    return heading_text, markup


def flow_slide_markup(slide, slide_index: int, page_w: float, page_h: float,
                      media_dir: Path, default_size: float,
                      doc_size: float, scale: float = 1.0,
                      calibration_marker: bool = False,
                      canvas_markers: bool = False,
                      heading_size: float | None = None,
                      trailing_marker: str = "") -> list[str]:
    """The complete markup of one slide (heading + flowing content).

    ``trailing_marker``, if set, is extra markup (the deck-end calibration
    sentinel) that must land *inside* this slide's own content -- appending
    it as bare top-level markup after a bracket-scoped slide (title-slide,
    centered-slide) does not keep it on that slide's last page; touying
    treats it as new content and starts an extra, duplicated page. Only the
    deck's last slide is ever called with a non-empty value.
    """
    from typstpresenter.convert.textbody import typst_str

    from typstpresenter.convert.emitter import slide_title_text

    title = slide_title_text(slide)
    if title and has_center_title(slide):
        return _title_slide_markup(slide, slide_index, title, page_w, page_h,
                                   media_dir, default_size, doc_size, scale,
                                   calibration_marker, trailing_marker)
    if title and is_section_divider(slide, page_w, page_h):
        return _section_divider_markup(slide, slide_index, title, scale,
                                       calibration_marker, trailing_marker)
    badge_markup = None
    if title:
        split = _title_trailing_badge(slide)
        if split:
            title, badge_markup = split
    body: list[str] = []
    heading = escape_flow(title) if title else ""
    if title:
        # slides whose resolved title size deviates from the deck-wide
        # heading size get an inline override inside the heading body
        own_size = slide_title_size(slide)
        if heading_size is not None and abs(own_size - heading_size) > 2:
            heading = f"#text(size: {own_size:g}pt)[{heading}]"
        link = slide_title_link(slide)
        if link:
            heading = f"#link({typst_str(link)})[{heading}]"
    body.append(f"== {heading}" if heading else "== ")
    body.append("")
    if badge_markup:
        body.append(badge_markup)
        body.append("")
    if calibration_marker:
        # temporary, invisible: lets `typst query` report the page each
        # slide starts on; the final output is emitted without markers.
        # place(hide(..)) keeps it out of flow, and the blank line after the
        # heading keeps touying from splitting the slide -- a bare or
        # heading-adjacent top-level element would create an extra page
        body.append(f"#place(hide(context metadata((s: {slide_index}, "
                    "p: here().position().page))))")
        body.append("")

    # slide-dominant body size: one #set rule instead of per-run sizes
    chunk_lists = [
        text_chunks_of_shape(shape, default_size)
        for shape, bbox in _content_text_shapes(slide, page_h, page_w)
    ]
    slide_size = dominant_size(chunk_lists, doc_size)
    context_size = _round_size(slide_size * scale)
    set_line = None
    if abs(context_size - doc_size) > 0.26:
        set_line = f"#set text(size: {context_size:g}pt)"

    content: list[str] = []
    prev: FlowBlock | None = None
    for block in flow_slide_blocks(slide, slide_index, page_w, page_h,
                                   media_dir, default_size, context_size,
                                   scale, canvas_markers=canvas_markers):
        # preserve deliberate vertical whitespace of the source layout
        # (sparse slides place content low on purpose); the reference is
        # where flow content normally starts (below the heading, if any)
        if prev is None:
            start = 130.0 if title else 60.0
            if block.bbox.y > start + 25:
                content += [f"#v({(block.bbox.y - start) * scale:.0f}pt)", ""]
        else:
            gap = block.bbox.y - prev.bbox.y2
            if gap > 45:
                content += [f"#v({(gap - 30) * scale:.0f}pt)", ""]
        content.append(block.markup)
        content.append("")
        prev = block
    if set_line:
        # a bare #set at slide level (or a `#[..]` scope containing #align)
        # makes touying split the slide in two; the block() function is the
        # only tested wrapper that keeps one slide. Full width so #align
        # still centers on the page.
        body += ["#block(width: 100%)[", set_line, ""] + content + ["]", ""]
    else:
        body += content
    if trailing_marker:
        body.append(trailing_marker)
    return body


def _title_slide_markup(slide, slide_index: int, title: str, page_w: float,
                        page_h: float, media_dir: Path, default_size: float,
                        doc_size: float, scale: float,
                        calibration_marker: bool,
                        trailing_marker: str = "") -> list[str]:
    """A CENTER_TITLE slide as `#title-slide[..]`: big centered title, the
    remaining content (subtitle, authors, date) centered below it."""
    from typstpresenter.convert.textbody import typst_str

    title_size = _round_size(slide_title_size(slide) * scale)
    body = ["#title-slide["]
    if calibration_marker:
        body.append(f"  #place(hide(context metadata((s: {slide_index}, "
                    "p: here().position().page))))")
    heading = f"#text(size: {title_size:g}pt)[{escape_flow(title)}]"
    link = slide_title_link(slide)
    if link:
        heading = f"#link({typst_str(link)})[{heading}]"
    body.append(f"  {heading}")
    body.append("")

    chunk_lists = [
        text_chunks_of_shape(shape, default_size)
        for shape, bbox in _content_text_shapes(slide, page_h, page_w)
    ]
    slide_size = dominant_size(chunk_lists, doc_size)
    context_size = _round_size(slide_size * scale)
    if abs(context_size - doc_size) > 0.26:
        body.append(f"  #set text(size: {context_size:g}pt)")
        body.append("")
    for block in flow_slide_blocks(slide, slide_index, page_w, page_h,
                                   media_dir, default_size, context_size,
                                   scale):
        body.append(block.markup)
        body.append("")
    if trailing_marker:
        body.append(f"  {trailing_marker}")
    body += ["]", ""]
    return body


def _section_divider_markup(slide, slide_index: int, title: str, scale: float,
                            calibration_marker: bool,
                            trailing_marker: str = "") -> list[str]:
    """A section-break slide: big bold title, centered on the page. Uses
    touying's `centered-slide` (from `themes.simple`) so the style stays
    uniform and easy to re-theme in one place."""
    from typstpresenter.convert.textbody import typst_str

    title_size = _round_size(slide_title_size(slide) * scale)
    body = ["#centered-slide["]
    if calibration_marker:
        body.append(f"  #place(hide(context metadata((s: {slide_index}, "
                    "p: here().position().page))))")
    heading = f'#text(size: {title_size:g}pt, weight: "bold")[{escape_flow(title)}]'
    link = slide_title_link(slide)
    if link:
        heading = f"#link({typst_str(link)})[{heading}]"
    body.append(f"  {heading}")
    if trailing_marker:
        body.append(f"  {trailing_marker}")
    body += ["]", ""]
    return body


def _content_text_shapes(slide, page_h: float, page_w: float = 1e9):
    from typstpresenter.verify.pptx_geometry import iter_flat_shapes

    for shape, bbox in iter_flat_shapes(slide.shapes):
        if (shape == slide.shapes.title
                or _is_chrome(shape, page_h, bbox, page_w)
                or _off_page(bbox, page_w, page_h)):
            continue
        if (not is_diagram_shape(shape) and shape.has_text_frame
                and shape.text_frame.text.strip()):
            yield shape, bbox


def document_body_size(prs, default_size: float) -> float:
    """Deck-wide dominant body text size (for the preamble #set text)."""
    page_h = prs.slide_height / EMU_PER_PT
    page_w = prs.slide_width / EMU_PER_PT
    chunk_lists = []
    for slide in prs.slides:
        for shape, _ in _content_text_shapes(slide, page_h, page_w):
            chunk_lists.append(text_chunks_of_shape(shape, default_size))
    return dominant_size(chunk_lists, default_size)
