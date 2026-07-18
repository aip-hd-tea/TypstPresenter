"""Autoshapes and connectors as CeTZ drawing commands (y axis flipped)."""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE

from typstpresenter.convert.pptx_inherit import resolve_anchor
from typstpresenter.convert.pptx_style import (
    shape_fill_rgb,
    shape_font_rgb,
    shape_line_rgb,
    shape_line_width_pt,
)
from typstpresenter.convert.textbody import (
    PPT_LINE_PITCH_EM,
    TYPST_LINE_HEIGHT_EM,
    autofit_scales,
    paragraph_runs_markup,
)
from typstpresenter.verify.geometry import BBox


def is_diagram_shape(shape) -> bool:
    return shape.shape_type in (
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.FREEFORM,
        MSO_SHAPE_TYPE.LINE,
    ) or shape.element.tag.endswith("}cxnSp")


def _freeform_points(shape) -> list[tuple[float, float]] | None:
    """Absolute-pt polygon vertices of a straight-edged freeform shape.

    Only the common case -- a single closed path built from moveTo/lnTo
    (no curves) -- is extracted exactly; anything with cubicBezTo,
    quadBezTo or arcTo segments, or more than one subpath, falls back to
    None so the caller draws the shape's plain bounding box instead.
    """
    from pptx.oxml.ns import qn

    from typstpresenter.verify.geometry import EMU_PER_PT

    custgeom = shape.element.spPr.find(qn("a:custGeom"))
    if custgeom is None:
        return None
    path_lst = custgeom.find(qn("a:pathLst"))
    if path_lst is None:
        return None
    paths = path_lst.findall(qn("a:path"))
    if len(paths) != 1:
        return None
    path = paths[0]
    path_w = int(path.get("w") or 0)
    path_h = int(path.get("h") or 0)
    if not path_w or not path_h:
        return None

    points: list[tuple[float, float]] = []
    for child in path:
        tag = child.tag.split("}")[-1]
        if tag == "close":
            continue
        if tag not in ("moveTo", "lnTo"):
            return None  # curve segment -- not a straight polygon
        pt = child.find(qn("a:pt"))
        if pt is None:
            return None
        points.append((int(pt.get("x")), int(pt.get("y"))))
    if len(points) < 3:
        return None

    sx = shape.width / path_w
    sy = shape.height / path_h
    return [
        ((shape.left + lx * sx) / EMU_PER_PT, (shape.top + ly * sy) / EMU_PER_PT)
        for lx, ly in points
    ]


def _shape_flips(shape) -> tuple[bool, bool]:
    from pptx.oxml.ns import qn

    xfrm = shape.element.spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return False, False
    return xfrm.get("flipH") == "1", xfrm.get("flipV") == "1"


def _connector_is_elbow(shape) -> bool:
    from pptx.oxml.ns import qn

    prst_geom = shape.element.spPr.find(qn("a:prstGeom"))
    prst = prst_geom.get("prst") if prst_geom is not None else ""
    return bool(prst) and prst.startswith("bentConnector")


def _shape_rotation_deg(shape) -> float:
    """Clockwise rotation in degrees (PPTX/OOXML convention), or 0.0."""
    from pptx.oxml.ns import qn

    xfrm = shape.element.spPr.find(qn("a:xfrm")) if shape.element.spPr is not None else None
    rot = xfrm.get("rot") if xfrm is not None else None
    return int(rot) / 60000.0 if rot else 0.0


def effective_bbox(shape, bbox: BBox) -> BBox:
    """`bbox` widened to the axis-aligned footprint of a rotated shape.

    A shape's PPTX bbox is always its *unrotated* extent; a diagram
    canvas that sizes itself from those raw boxes underestimates a
    rotated shape's true ink (e.g. a tall bar rotated 90 degrees becomes
    wide), causing the canvas to overflow its declared bounds.
    """
    import math

    angle = _shape_rotation_deg(shape)
    if not angle or shape.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE,
                                             MSO_SHAPE_TYPE.TEXT_BOX):
        return bbox
    rad = math.radians(angle)
    cos_a, sin_a = math.cos(rad), math.sin(rad)
    cx, cy = bbox.center
    corners = [(bbox.x, bbox.y), (bbox.x2, bbox.y), (bbox.x2, bbox.y2), (bbox.x, bbox.y2)]
    rotated = []
    for x, y in corners:
        dx, dy = x - cx, y - cy
        rotated.append((cx + dx * cos_a - dy * sin_a, cy + dx * sin_a + dy * cos_a))
    xs = [p[0] for p in rotated]
    ys = [p[1] for p in rotated]
    return BBox(min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys))


def _cetz_style_args(shape) -> str:
    """fill/stroke arguments resolved from explicit colors or the theme."""
    fill = shape_fill_rgb(shape)
    stroke_rgb = shape_line_rgb(shape)
    fill_arg = f'fill: rgb("#{fill}")' if fill else "fill: none"
    if stroke_rgb:
        width = shape_line_width_pt(shape)
        stroke_arg = f'stroke: (paint: rgb("#{stroke_rgb}"), thickness: {width:g}pt)'
    else:
        stroke_arg = "stroke: none"
    return f"{fill_arg}, {stroke_arg}"


def _cetz_label_markup(shape, eid: str, probes: bool, default_size: float,
                       label_w: float) -> str:
    """Shape label as width-constrained box (multi-paragraph, styled, probed)."""
    parts = []
    if shape.has_text_frame:
        # shape text defaults to the style's font color (often white)
        font_rgb = shape_font_rgb(shape)
        font_scale, lnspc_red = autofit_scales(shape)
        for paragraph in shape.text_frame.paragraphs:
            runs = paragraph_runs_markup(paragraph, shape, default_size,
                                         scale=font_scale,
                                         default_color=font_rgb)
            if runs:
                parts.append(f"#par[{''.join(runs)}]")
        if parts:
            # PowerPoint line metrics, as in emit_text_body: without this,
            # multi-paragraph labels get typst's default paragraph spacing
            # and grow far taller than the shape
            leading = max(
                PPT_LINE_PITCH_EM * (1.0 - lnspc_red) - TYPST_LINE_HEIGHT_EM, 0.1
            )
            parts.insert(
                0, f"#set par(leading: {leading:.3f}em, spacing: {leading:.3f}em);"
            )
    content = "".join(parts)
    boxed = f"box(width: {label_w:.2f}pt, align(center)[{content}])" if content else ""
    if not probes:
        return f"#{boxed}" if boxed else ""
    if boxed:
        # label probe measures the constrained label so overflow beyond the
        # shape can be classified as a source condition, not a translation bug
        return f'#tp-label-probe("{eid}", {boxed})'
    return f'#tp-node-probe("{eid}")'


def emit_cetz_shape(shape, eid: str, bbox: BBox, probes: bool,
                    default_size: float = 18.0):
    """One PPTX autoshape/connector as CeTZ drawing commands.

    Returns (markup, label_expectation); the expectation is
    (center_x, anchor_y, anchor_kind) for drift calibration, or None.
    """
    from pptx.enum.shapes import MSO_SHAPE

    lines = []
    x1, y1, x2, y2 = bbox.x, -bbox.y, bbox.x2, -bbox.y2
    cx, cy = bbox.center
    # NOTE: this used to be `shape.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE,)`,
    # which also matched FREEFORM shapes -- they were silently drawn as a
    # bogus diagonal "connector" line across their own bounding box instead
    # of their actual polygon. Line/connector detection must be explicit.
    is_connector = (shape.shape_type == MSO_SHAPE_TYPE.LINE
                    or shape.element.tag.endswith("}cxnSp"))

    if is_connector:
        flip_h, flip_v = _shape_flips(shape)
        bx, ex = (x2, x1) if flip_h else (x1, x2)
        by, ey = (y2, y1) if flip_v else (y1, y2)
        stroke_rgb = shape_line_rgb(shape) or "000000"
        width = shape_line_width_pt(shape)
        stroke = f'stroke: (paint: rgb("#{stroke_rgb}"), thickness: {width:g}pt)'
        if _connector_is_elbow(shape):
            mx = (bx + ex) / 2
            points = (f"({bx:.2f}, {by:.2f}), ({mx:.2f}, {by:.2f}), "
                      f"({mx:.2f}, {ey:.2f}), ({ex:.2f}, {ey:.2f})")
        else:
            points = f"({bx:.2f}, {by:.2f}), ({ex:.2f}, {ey:.2f})"
        lines.append(f"    line({points}, mark: (end: \">\"), {stroke})")
        if probes:
            lines.append(f'    content(({cx:.2f}, {-cy:.2f}), [#tp-node-probe("{eid}")])')
        return "\n".join(lines), None

    auto = None
    try:
        auto = shape.auto_shape_type
    except (ValueError, AttributeError):
        pass
    style = _cetz_style_args(shape)
    if auto == MSO_SHAPE.OVAL:
        lines.append(
            f"    circle(({cx:.2f}, {-cy:.2f}), "
            f"radius: ({bbox.w / 2:.2f}, {bbox.h / 2:.2f}), {style})"
        )
    elif auto == MSO_SHAPE.DIAMOND:
        lines.append(
            f"    line(({cx:.2f}, {y1:.2f}), ({x2:.2f}, {-cy:.2f}), "
            f"({cx:.2f}, {y2:.2f}), ({x1:.2f}, {-cy:.2f}), close: true, {style})"
        )
    elif auto == MSO_SHAPE.ISOSCELES_TRIANGLE:
        lines.append(
            f"    line(({cx:.2f}, {y1:.2f}), ({x2:.2f}, {y2:.2f}), "
            f"({x1:.2f}, {y2:.2f}), close: true, {style})"
        )
    elif auto in (MSO_SHAPE.LEFT_BRACE, MSO_SHAPE.RIGHT_BRACE,
                  MSO_SHAPE.LEFT_BRACKET, MSO_SHAPE.RIGHT_BRACKET):
        # braces/brackets are open strokes -- a theme-filled rect here
        # would paint a solid bar over whatever the brace annotates
        stroke_rgb = (shape_line_rgb(shape) or shape_fill_rgb(shape)
                      or "000000")
        width = shape_line_width_pt(shape)
        stroke = f'stroke: (paint: rgb("#{stroke_rgb}"), thickness: {width:g}pt)'
        left = auto in (MSO_SHAPE.LEFT_BRACE, MSO_SHAPE.LEFT_BRACKET)
        xo, xt = (x2, x1) if left else (x1, x2)  # open side, tip side
        if auto in (MSO_SHAPE.LEFT_BRACE, MSO_SHAPE.RIGHT_BRACE):
            xs = (x1 + x2) / 2  # spine in the middle, tip pokes out at cy
            pts = (f"({xo:.2f}, {y1:.2f}), ({xs:.2f}, {y1:.2f}), "
                   f"({xs:.2f}, {-cy:.2f}), ({xt:.2f}, {-cy:.2f}), "
                   f"({xs:.2f}, {-cy:.2f}), ({xs:.2f}, {y2:.2f}), "
                   f"({xo:.2f}, {y2:.2f})")
        else:
            pts = (f"({xo:.2f}, {y1:.2f}), ({xt:.2f}, {y1:.2f}), "
                   f"({xt:.2f}, {y2:.2f}), ({xo:.2f}, {y2:.2f})")
        lines.append(f"    line({pts}, {stroke})")
    elif auto is None and shape.shape_type == MSO_SHAPE_TYPE.FREEFORM and (
        freeform_pts := _freeform_points(shape)
    ):
        pts = ", ".join(f"({x:.2f}, {-y:.2f})" for x, y in freeform_pts)
        lines.append(f"    line({pts}, close: true, {style})")
    else:
        radius = 4.0 if auto == MSO_SHAPE.ROUNDED_RECTANGLE else 0.0
        lines.append(
            f"    rect(({x1:.2f}, {y1:.2f}), ({x2:.2f}, {y2:.2f}), "
            f"radius: {radius:g}, {style})"
        )
    # constrain the label to the shape width so long text wraps like
    # PowerPoint instead of spreading beyond the shape
    label_w = max(bbox.w - 7.2, 7.2)
    label = _cetz_label_markup(shape, eid, probes, default_size, label_w)
    expectation = None
    if label:
        from pptx.enum.text import MSO_ANCHOR

        # respect the text frame's vertical anchor: top/bottom-anchored
        # panels must not be centered (their text grows from the edge).
        # The OOXML default when no anchor is set anywhere is "t" (top).
        anchor = resolve_anchor(shape) if shape.has_text_frame else None
        if anchor == MSO_ANCHOR.BOTTOM:
            lines.append(
                f'    content(({cx:.2f}, {y2 + 3.6:.2f}), anchor: "south", [{label}])'
            )
            expectation = (cx, bbox.y2 - 3.6, "south")
        elif anchor == MSO_ANCHOR.MIDDLE:
            lines.append(f"    content(({cx:.2f}, {-cy:.2f}), [{label}])")
            expectation = (cx, cy, "center")
        else:  # OOXML default: top
            lines.append(
                f'    content(({cx:.2f}, {y1 - 3.6:.2f}), anchor: "north", [{label}])'
            )
            expectation = (cx, bbox.y + 3.6, "north")

    # PPTX rotates clockwise around the shape center; CeTZ's rotate() is
    # the standard CCW-positive convention, and our canvas y is negated,
    # which flips the rotation sense once more -- net effect: negate the
    # angle. Only applied in flow mode: rotating the group would also
    # rotate the probe anchors the (legacy) probed emitter's calibration
    # measures, which assumes axis-aligned content.
    if not probes:
        angle = _shape_rotation_deg(shape)
        if angle:
            body = "\n".join(f"  {line}" for line in lines)
            lines = [
                f"  group({{",
                f"    rotate({-angle:.2f}deg, origin: ({cx:.2f}, {-cy:.2f}))",
                body,
                "  })",
            ]
    return "\n".join(lines), expectation
