"""Public API: translate PPTX shapes to an SVG document.

L1 scope: preset auto-shapes and straight connectors (prstGeom), solid
fill/stroke, flips.  Rotation is parsed and emitted (verified in L3).
Unsupported geometry (custGeom) falls back to a bounding-box rectangle
and is reported in the result.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx.oxml.ns import qn

from typstpresenter.diagram2svg.presets import (
    evaluate_custgeom,
    evaluate_preset,
    evaluate_text_rect,
    has_preset,
)
from typstpresenter.diagram2svg.style import resolve_style
from typstpresenter.diagram2svg.svg_writer import ArrowEnd, SvgDocument, SvgShape
from typstpresenter.diagram2svg.text import layout_shape_text
from typstpresenter.verify.geometry import BBox
from typstpresenter.verify.pptx_geometry import element_id, iter_flat_shapes


@dataclass
class SvgResult:
    svg: str
    shape_count: int
    fallbacks: list[str] = field(default_factory=list)  # ids drawn as bbox rect
    skipped: list[str] = field(default_factory=list)  # ids not drawn at all


def _spPr(shape):
    return shape.element.find(qn("p:spPr"))


def _preset_of(shape) -> tuple[str | None, dict[str, str]]:
    """(preset name, avLst overrides) from the shape's prstGeom, if any."""
    sp_pr = _spPr(shape)
    if sp_pr is None:
        return None, {}
    prst = sp_pr.find(qn("a:prstGeom"))
    if prst is None:
        return None, {}
    overrides: dict[str, str] = {}
    av = prst.find(qn("a:avLst"))
    if av is not None:
        for gd in av.findall(qn("a:gd")):
            overrides[gd.get("name")] = gd.get("fmla")
    return prst.get("prst"), overrides


def _xfrm_of(shape) -> tuple[float, bool, bool]:
    """(rotation deg CW, flipH, flipV)."""
    sp_pr = _spPr(shape)
    xfrm = sp_pr.find(qn("a:xfrm")) if sp_pr is not None else None
    if xfrm is None:
        return 0.0, False, False
    rot = float(xfrm.get("rot", "0")) / 60000.0
    return rot, xfrm.get("flipH") == "1", xfrm.get("flipV") == "1"


def _arrow_ends(shape) -> tuple[ArrowEnd | None, ArrowEnd | None]:
    """(headEnd, tailEnd) markers from spPr/a:ln."""
    sp_pr = _spPr(shape)
    ln = sp_pr.find(qn("a:ln")) if sp_pr is not None else None
    if ln is None:
        return None, None

    def _end(tag: str) -> ArrowEnd | None:
        el = ln.find(qn(tag))
        if el is None or el.get("type") in (None, "none"):
            return None
        return ArrowEnd(
            kind=el.get("type"),
            w=el.get("w", "med"),
            length=el.get("len", "med"),
        )

    return _end("a:headEnd"), _end("a:tailEnd")


def _has_geometry(shape) -> bool:
    """Shapes we translate: auto-shapes and connectors with prstGeom/custGeom."""
    sp_pr = _spPr(shape)
    if sp_pr is None:
        return False
    return (
        sp_pr.find(qn("a:prstGeom")) is not None
        or sp_pr.find(qn("a:custGeom")) is not None
    )


def shapes_to_svg(
    pairs: list[tuple[object, BBox]],
    region: BBox,
    slide_index: int,
) -> SvgResult:
    """Translate (shape, absolute pt bbox) pairs into one SVG document.

    `pairs` come from iter_flat_shapes (groups already flattened); order
    is source z-order and is preserved as SVG paint order.
    """
    doc = SvgDocument(x=region.x, y=region.y, w=region.w, h=region.h)
    result = SvgResult(svg="", shape_count=0)

    for shape, bbox in pairs:
        sid = element_id(slide_index, shape.shape_id)
        if not _has_geometry(shape):
            result.skipped.append(sid)
            continue
        svg_shape, fell_back = build_svg_shape(shape, bbox, sid)
        if fell_back:
            result.fallbacks.append(sid)
        doc.shapes.append(svg_shape)
        result.shape_count += 1

    result.svg = doc.to_string()
    return result


def build_svg_shape(
    shape, bbox: BBox, sid: str, default_size: float = 18.0
) -> tuple[SvgShape, bool]:
    """Translate one geometric shape; returns (SvgShape, bbox_fallback?)."""
    preset, overrides = _preset_of(shape)
    w_emu = bbox.w * 12700.0
    h_emu = bbox.h * 12700.0
    text_rect = None
    paths = None
    fell_back = False
    if preset and has_preset(preset):
        paths = evaluate_preset(preset, w_emu, h_emu, overrides)
        rect_emu = evaluate_text_rect(preset, w_emu, h_emu, overrides)
        if rect_emu:
            text_rect = tuple(v / 12700.0 for v in rect_emu)
    else:
        cust = _spPr(shape).find(qn("a:custGeom"))
        if cust is not None:
            try:
                paths = evaluate_custgeom(cust, w_emu, h_emu)
            except (KeyError, ValueError, IndexError):
                paths = None
    if not paths:
        # unknown preset or unevaluable custGeom: bbox rectangle fallback
        paths = evaluate_preset("rect", w_emu, h_emu)
        fell_back = True
    rot, flip_h, flip_v = _xfrm_of(shape)
    head, tail = _arrow_ends(shape)
    return SvgShape(
        id=sid,
        x_pt=bbox.x,
        y_pt=bbox.y,
        w_pt=bbox.w,
        h_pt=bbox.h,
        paths=paths,
        style=resolve_style(shape),
        rot_deg=rot,
        flip_h=flip_h,
        flip_v=flip_v,
        head=head,
        tail=tail,
        text_elements=layout_shape_text(shape, bbox.w, bbox.h, text_rect,
                                        default_size=default_size),
    ), fell_back


def slide_to_svg(slide, slide_index: int, slide_w_pt: float, slide_h_pt: float) -> SvgResult:
    """Translate a whole slide's geometric shapes (full-slide viewBox)."""
    pairs = list(iter_flat_shapes(slide.shapes))
    region = BBox(x=0, y=0, w=slide_w_pt, h=slide_h_pt)
    return shapes_to_svg(pairs, region, slide_index)


def pptx_to_svgs(pptx_path: Path | str, out_dir: Path | str) -> list[tuple[Path, SvgResult]]:
    """Convert every slide of a deck to slide{N}.svg in out_dir."""
    import pptx as _pptx

    prs = _pptx.Presentation(str(pptx_path))
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    slide_w = prs.slide_width / 12700.0
    slide_h = prs.slide_height / 12700.0
    results = []
    for i, slide in enumerate(prs.slides):
        res = slide_to_svg(slide, i, slide_w, slide_h)
        path = out / f"slide{i + 1}.svg"
        path.write_text(res.svg, encoding="utf-8")
        results.append((path, res))
    return results
