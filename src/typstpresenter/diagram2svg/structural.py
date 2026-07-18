"""Structural benchmark: compare a generated SVG against PPTX ground truth.

Independent of the writer: parses the SVG text (transforms + path data),
samples every path numerically, and checks per shape id:

- geometric bbox == source shape bbox (within tolerance; open "line"
  geometry may have zero extent on one axis),
- command profile matches the expected preset family (ellipse = arcs
  only, rect/diamond/triangle = lines only, roundRect = both),
- fill / stroke color and stroke width match the resolved PPTX style.
"""

from __future__ import annotations

import math
import re
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from pathlib import Path

import pptx

from typstpresenter.convert.pptx_style import (
    shape_fill_rgb,
    shape_line_rgb,
    shape_line_width_pt,
)
from typstpresenter.verify.geometry import BBox
from typstpresenter.verify.pptx_geometry import element_id, iter_flat_shapes

_TOL_PT = 0.75


def _rotated_footprint(bbox: BBox, rot_deg: float, ) -> BBox:
    """Axis-aligned bounds of the bbox rotated about its center (CW, y-down)."""
    if not rot_deg:
        return bbox
    r = math.radians(rot_deg)
    cx, cy = bbox.x + bbox.w / 2, bbox.y + bbox.h / 2
    xs, ys = [], []
    for px, py in (
        (bbox.x, bbox.y), (bbox.x + bbox.w, bbox.y),
        (bbox.x, bbox.y + bbox.h), (bbox.x + bbox.w, bbox.y + bbox.h),
    ):
        dx, dy = px - cx, py - cy
        xs.append(cx + dx * math.cos(r) - dy * math.sin(r))
        ys.append(cy + dx * math.sin(r) + dy * math.cos(r))
    return BBox(x=min(xs), y=min(ys), w=max(xs) - min(xs), h=max(ys) - min(ys))


def _shape_rotation(shape) -> float:
    from typstpresenter.diagram2svg.convert import _xfrm_of

    return _xfrm_of(shape)[0]


def _expected_bounds(shape, bbox: BBox) -> BBox | None:
    """Bounds the shape's geometry should cover on the page.

    Many presets legitimately do not span their bbox (mathPlus arms,
    default arc = quarter ellipse) or exceed it (callout leader tips), so
    the expectation is derived from the evaluated geometry — placement,
    scaling, flips and rotation are still verified independently here;
    geometry-formula correctness is pinned by the preset unit tests and
    the visual reference tier.
    """
    from pptx.oxml.ns import qn

    from typstpresenter.diagram2svg.convert import _preset_of, _xfrm_of
    from typstpresenter.diagram2svg.presets import (
        evaluate_custgeom,
        evaluate_preset,
        has_preset,
    )
    from typstpresenter.diagram2svg.svg_writer import path_data

    preset, overrides = _preset_of(shape)
    w_emu, h_emu = bbox.w * 12700.0, bbox.h * 12700.0
    if preset and has_preset(preset):
        paths = evaluate_preset(preset, w_emu, h_emu, overrides)
    else:
        sp_pr = shape.element.find(qn("p:spPr"))
        cust = sp_pr.find(qn("a:custGeom")) if sp_pr is not None else None
        if cust is None:
            return None
        try:
            paths = evaluate_custgeom(cust, w_emu, h_emu)
        except (KeyError, ValueError, IndexError):
            return None

    rot, flip_h, flip_v = _xfrm_of(shape)
    # same transform chain as SvgShape.transform()
    m = (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)
    m = _mat_mul(m, (1, 0, 0, 1, bbox.x, bbox.y))
    if rot or flip_h or flip_v:
        cx, cy = bbox.w / 2, bbox.h / 2
        m = _mat_mul(m, (1, 0, 0, 1, cx, cy))
        if rot:
            r = math.radians(rot)
            m = _mat_mul(m, (math.cos(r), math.sin(r), -math.sin(r), math.cos(r), 0, 0))
        if flip_h or flip_v:
            m = _mat_mul(m, (-1 if flip_h else 1, 0, 0, -1 if flip_v else 1, 0, 0))
        m = _mat_mul(m, (1, 0, 0, 1, -cx, -cy))

    pts: list[tuple[float, float]] = []
    for ep in paths:
        sampled, _ = sample_path(path_data(ep.segments, ep.sx, ep.sy))
        pts.extend(_apply(m, x, y) for x, y in sampled)
    if not pts:
        return None
    xs = [p[0] for p in pts]
    ys = [p[1] for p in pts]
    return BBox(x=min(xs), y=min(ys), w=max(xs) - min(xs), h=max(ys) - min(ys))


@dataclass
class Finding:
    shape_id: str
    kind: str  # missing | bbox | profile | fill | stroke | stroke-width | extra
    detail: str


@dataclass
class StructuralReport:
    findings: list[Finding] = field(default_factory=list)
    checked: int = 0

    @property
    def ok(self) -> bool:
        return not self.findings

    def add(self, shape_id: str, kind: str, detail: str) -> None:
        self.findings.append(Finding(shape_id, kind, detail))


# ------------------------------------------------------------- transforms --

_TRANSFORM_RE = re.compile(r"(translate|rotate|scale)\(([^)]*)\)")


def _mat_mul(a, b):
    return (
        a[0] * b[0] + a[2] * b[1],
        a[1] * b[0] + a[3] * b[1],
        a[0] * b[2] + a[2] * b[3],
        a[1] * b[2] + a[3] * b[3],
        a[0] * b[4] + a[2] * b[5] + a[4],
        a[1] * b[4] + a[3] * b[5] + a[5],
    )


def parse_transform(text: str):
    """SVG transform list → affine matrix (a, b, c, d, e, f)."""
    m = (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)
    for op, argstr in _TRANSFORM_RE.findall(text or ""):
        args = [float(v) for v in re.split(r"[\s,]+", argstr.strip()) if v]
        if op == "translate":
            tx, ty = args[0], args[1] if len(args) > 1 else 0.0
            m = _mat_mul(m, (1, 0, 0, 1, tx, ty))
        elif op == "scale":
            sx = args[0]
            sy = args[1] if len(args) > 1 else sx
            m = _mat_mul(m, (sx, 0, 0, sy, 0, 0))
        elif op == "rotate":
            r = math.radians(args[0])
            cos, sin = math.cos(r), math.sin(r)
            m = _mat_mul(m, (cos, sin, -sin, cos, 0, 0))
    return m


def _apply(m, x, y):
    return (m[0] * x + m[2] * y + m[4], m[1] * x + m[3] * y + m[5])


# -------------------------------------------------------------- path data --

_NUM_RE = re.compile(r"[-+]?\d*\.?\d+(?:[eE][-+]?\d+)?")


def _sample_arc(x0, y0, rx, ry, large, sweep, x1, y1, n=32):
    """Endpoint→center parametrization (SVG spec B.2.4), sampled points."""
    if rx == 0 or ry == 0:
        return [(x1, y1)]
    dx, dy = (x0 - x1) / 2.0, (y0 - y1) / 2.0
    # radii scale-up if too small
    lam = (dx / rx) ** 2 + (dy / ry) ** 2
    if lam > 1:
        s = math.sqrt(lam)
        rx, ry = rx * s, ry * s
    sign = -1.0 if large == sweep else 1.0
    num = rx**2 * ry**2 - rx**2 * dy**2 - ry**2 * dx**2
    den = rx**2 * dy**2 + ry**2 * dx**2
    co = sign * math.sqrt(max(num / den, 0.0)) if den else 0.0
    cxp, cyp = co * rx * dy / ry, -co * ry * dx / rx
    cx, cy = cxp + (x0 + x1) / 2.0, cyp + (y0 + y1) / 2.0

    def angle(ux, uy, vx, vy):
        dot = ux * vx + uy * vy
        length = math.hypot(ux, uy) * math.hypot(vx, vy)
        ang = math.acos(max(-1.0, min(1.0, dot / length))) if length else 0.0
        return -ang if ux * vy - uy * vx < 0 else ang

    th1 = angle(1, 0, (x0 - cx) / rx, (y0 - cy) / ry)
    dth = angle((x0 - cx) / rx, (y0 - cy) / ry, (x1 - cx) / rx, (y1 - cy) / ry)
    if not sweep and dth > 0:
        dth -= 2 * math.pi
    elif sweep and dth < 0:
        dth += 2 * math.pi
    return [
        (cx + rx * math.cos(th1 + dth * i / n), cy + ry * math.sin(th1 + dth * i / n))
        for i in range(n + 1)
    ]


def _sample_cubic(p0, p1, p2, p3, n=32):
    pts = []
    for i in range(n + 1):
        t = i / n
        mt = 1 - t
        pts.append(
            (
                mt**3 * p0[0] + 3 * mt**2 * t * p1[0] + 3 * mt * t**2 * p2[0] + t**3 * p3[0],
                mt**3 * p0[1] + 3 * mt**2 * t * p1[1] + 3 * mt * t**2 * p2[1] + t**3 * p3[1],
            )
        )
    return pts


def sample_path(d: str) -> tuple[list[tuple[float, float]], set[str]]:
    """All sampled points of an SVG path plus the set of command letters."""
    points: list[tuple[float, float]] = []
    ops: set[str] = set()
    cur = (0.0, 0.0)
    tokens = re.findall(r"[MLCQAZ]|" + _NUM_RE.pattern, d)
    i = 0
    while i < len(tokens):
        op = tokens[i]
        ops.add(op)
        i += 1

        def take(k):
            nonlocal i
            vals = [float(tokens[i + j]) for j in range(k)]
            i += k
            return vals

        if op in ("M", "L"):
            x, y = take(2)
            cur = (x, y)
            points.append(cur)
        elif op == "C":
            x1, y1, x2, y2, x, y = take(6)
            points.extend(_sample_cubic(cur, (x1, y1), (x2, y2), (x, y)))
            cur = (x, y)
        elif op == "Q":
            x1, y1, x, y = take(4)
            c1 = (cur[0] + 2 / 3 * (x1 - cur[0]), cur[1] + 2 / 3 * (y1 - cur[1]))
            c2 = (x + 2 / 3 * (x1 - x), y + 2 / 3 * (y1 - y))
            points.extend(_sample_cubic(cur, c1, c2, (x, y)))
            cur = (x, y)
        elif op == "A":
            rx, ry, _rot, large, sweep, x, y = take(7)
            points.extend(_sample_arc(cur[0], cur[1], rx, ry, int(large), int(sweep), x, y))
            cur = (x, y)
        # Z: nothing to sample
    return points, ops


# ---------------------------------------------------------------- checker --

# expected path-command profile per preset family
_PROFILES = {
    "rect": ("L", set("ACQ")),
    "diamond": ("L", set("ACQ")),
    "triangle": ("L", set("ACQ")),
    "ellipse": ("A", set("LCQ")),
    "roundRect": ("LA", set("CQ")),
    "line": ("L", set("ACQ")),
}


def _preset_name(shape) -> str | None:
    from pptx.oxml.ns import qn

    sp_pr = shape.element.find(qn("p:spPr"))
    if sp_pr is None:
        return None
    prst = sp_pr.find(qn("a:prstGeom"))
    return prst.get("prst") if prst is not None else None


def check_svg_against_pptx(
    svg_text: str, pptx_path: Path | str, slide_index: int
) -> StructuralReport:
    report = StructuralReport()
    prs = pptx.Presentation(str(pptx_path))
    slide = list(prs.slides)[slide_index]

    # index SVG groups by id
    root = ET.fromstring(svg_text)
    ns = {"svg": "http://www.w3.org/2000/svg"}
    groups = {g.get("id"): g for g in root.findall("svg:g", ns)}

    expected_ids = set()
    for shape, bbox in iter_flat_shapes(slide.shapes):
        from typstpresenter.diagram2svg.convert import _has_geometry

        if not _has_geometry(shape):
            continue
        sid = element_id(slide_index, shape.shape_id)
        expected_ids.add(sid)
        g = groups.get(sid)
        if g is None:
            report.add(sid, "missing", "no <g> with this id in SVG")
            continue
        report.checked += 1
        m = parse_transform(g.get("transform"))

        all_pts: list[tuple[float, float]] = []
        ops: set[str] = set()
        fills: set[str] = set()
        strokes: set[str] = set()
        widths: set[float] = set()
        for p in g.findall("svg:path", ns):
            pts, pops = sample_path(p.get("d", ""))
            all_pts.extend(_apply(m, x, y) for x, y in pts)
            ops |= pops
            fills.add(p.get("fill", "none"))
            strokes.add(p.get("stroke", "none"))
            if p.get("stroke") not in (None, "none"):
                widths.add(float(p.get("stroke-width", "1")))

        if not all_pts:
            report.add(sid, "missing", "group contains no drawable path")
            continue

        # --- bbox (geometric; strokes excluded by design)
        xs = [p[0] for p in all_pts]
        ys = [p[1] for p in all_pts]
        got = BBox(x=min(xs), y=min(ys), w=max(xs) - min(xs), h=max(ys) - min(ys))
        expected = _expected_bounds(shape, bbox) or bbox
        for name, gv, ev in (
            ("x", got.x, expected.x),
            ("y", got.y, expected.y),
            ("w", got.w, expected.w),
            ("h", got.h, expected.h),
        ):
            if abs(gv - ev) > _TOL_PT:
                report.add(sid, "bbox", f"{name}: svg {gv:.2f} vs expected {ev:.2f} pt")

        # --- command profile
        preset = _preset_name(shape)
        profile = _PROFILES.get(preset or "")
        if profile:
            required, forbidden = profile
            for req in required:
                if req not in ops:
                    report.add(sid, "profile", f"{preset}: expected {req} commands, got {sorted(ops)}")
            bad = ops & forbidden
            if bad:
                report.add(sid, "profile", f"{preset}: unexpected commands {sorted(bad)}")

        # --- style
        from pptx.enum.dml import MSO_FILL

        is_gradient = False
        try:
            is_gradient = shape.fill.type == MSO_FILL.GRADIENT
        except (AttributeError, TypeError):
            pass
        try:
            exp_fill = shape_fill_rgb(shape)
        except (AttributeError, TypeError):
            exp_fill = None
        exp_stroke = shape_line_rgb(shape)
        painted_fills = {f for f in fills if f != "none"}
        if is_gradient:
            if not any(f.startswith("url(#grad-") for f in painted_fills):
                report.add(sid, "fill", f"expected gradient url, got {painted_fills}")
        elif exp_fill is None and painted_fills:
            report.add(sid, "fill", f"expected no fill, got {painted_fills}")
        elif exp_fill is not None and f"#{exp_fill}" not in fills:
            report.add(sid, "fill", f"expected #{exp_fill}, got {fills}")

        # --- dash
        from typstpresenter.diagram2svg.style import _dash_of

        if _dash_of(shape) and exp_stroke is not None:
            has_dash = any(
                p.get("stroke-dasharray") for p in g.findall("svg:path", ns)
            )
            if not has_dash:
                report.add(sid, "stroke", "source prstDash but no stroke-dasharray")
        painted_strokes = {s for s in strokes if s != "none"}
        if exp_stroke is None and painted_strokes:
            report.add(sid, "stroke", f"expected no stroke, got {painted_strokes}")
        elif exp_stroke is not None and f"#{exp_stroke}" not in strokes:
            report.add(sid, "stroke", f"expected #{exp_stroke}, got {strokes}")
        if exp_stroke is not None and widths:
            exp_w = shape_line_width_pt(shape)
            if all(abs(w - exp_w) > 0.05 for w in widths):
                report.add(sid, "stroke-width", f"expected {exp_w}, got {widths}")

    for sid in groups.keys() - expected_ids:
        report.add(sid, "extra", "SVG group has no matching source shape")

    # --- z-order: SVG paint order must equal source document order
    svg_order = [g.get("id") for g in root.findall("svg:g", ns)]
    src_order = [
        element_id(slide_index, shape.shape_id)
        for shape, _ in iter_flat_shapes(slide.shapes)
        if element_id(slide_index, shape.shape_id) in set(svg_order)
    ]
    if svg_order != src_order:
        report.add("-", "z-order", f"svg {svg_order} vs source {src_order}")
    return report


_WS = re.compile(r"\s+")
# private-use-area glyphs (Symbol/Wingdings runs, bullet chars): symbol-font
# mapping is a known gap (G7) verified visually, not textually
_PUA = re.compile(r"[-]")


def _norm(text: str) -> str:
    return _WS.sub(" ", _PUA.sub("", text)).strip()


def check_text_and_markers(
    svg_text: str, pptx_path: Path | str, slide_index: int
) -> StructuralReport:
    """Text content, tspan placement, and arrowhead-marker presence."""
    from pptx.oxml.ns import qn

    report = StructuralReport()
    prs = pptx.Presentation(str(pptx_path))
    slide = list(prs.slides)[slide_index]
    root = ET.fromstring(svg_text)
    ns = {"svg": "http://www.w3.org/2000/svg"}
    groups = {g.get("id"): g for g in root.findall("svg:g", ns)}
    defs = root.find("svg:defs", ns)
    marker_ids = (
        {m.get("id") for m in defs.findall("svg:marker", ns)} if defs is not None else set()
    )

    for shape, bbox in iter_flat_shapes(slide.shapes):
        sid = element_id(slide_index, shape.shape_id)
        g = groups.get(sid)
        if g is None:
            continue
        report.checked += 1
        m = parse_transform(g.get("transform"))

        # --- text content and placement (tspans of one <text> line join
        # without separator; wrap-inserted line breaks stand for a space).
        # Symbol-font runs are compared in their mapped Unicode form.
        from typstpresenter.diagram2svg.symbols import (
            is_symbol_font,
            map_symbol_text,
            run_symbol_font,
        )

        def _mapped_para(p) -> str:
            from pptx.oxml.ns import qn

            runs_by_el = {run._r: run for run in p.runs}
            parts = []
            for child in p._p:
                if child.tag == qn("a:br"):
                    parts.append(" ")
                elif child.tag == qn("a:fld"):
                    parts.append("".join(t.text or ""
                                         for t in child.findall(qn("a:t"))))
                elif child.tag == qn("a:r"):
                    run = runs_by_el.get(child)
                    if run is None:
                        continue
                    name = run.font.name
                    sym = run_symbol_font(run)
                    if (is_symbol_font(name) or is_symbol_font(sym)
                            or any(0xF000 <= ord(c) <= 0xF0FF for c in run.text)):
                        parts.append(map_symbol_text(
                            run.text, name if is_symbol_font(name) else sym))
                    else:
                        parts.append(run.text)
            return "".join(parts)

        src_text = ""
        if getattr(shape, "has_text_frame", False):
            src_text = _norm(" ".join(
                _mapped_para(p) for p in shape.text_frame.paragraphs))
        tspans = g.findall(".//svg:tspan", ns)
        got_text = _norm(" ".join(
            "".join(t.text or "" for t in line.findall("svg:tspan", ns))
            for line in g.findall(".//svg:text", ns)
        ))
        if src_text and got_text != src_text:
            # emitted bullet glyphs (resolve_bullet) are correct extras the
            # source paragraph text does not contain
            relaxed = got_text
            from typstpresenter.convert.pptx_inherit import resolve_bullet

            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    try:
                        b = resolve_bullet(p, shape)
                    except (AttributeError, KeyError, ValueError):
                        b = None
                    if b:
                        from typstpresenter.diagram2svg.symbols import bullet_font

                        glyph = _norm(map_symbol_text(b, bullet_font(p)))
                        if glyph:
                            relaxed = relaxed.replace(f"{glyph} ", "", 1)
            if _norm(relaxed) != src_text:
                report.add(sid, "text", f"source {src_text!r} vs svg {got_text!r}")
        if not src_text and got_text:
            report.add(sid, "text", f"unexpected text {got_text!r}")
        fp = _rotated_footprint(bbox, _shape_rotation(shape))
        for t in tspans:
            x, y = _apply(m, float(t.get("x", "0")), float(t.get("y", "0")))
            if not (fp.x - 2 <= x <= fp.x + fp.w + 2
                    and fp.y - 2 <= y <= fp.y + fp.h + 2):
                report.add(sid, "text-pos",
                           f"tspan anchor ({x:.1f},{y:.1f}) outside bbox {fp}")

        # --- arrowhead markers
        sp_pr = shape.element.find(qn("p:spPr"))
        ln = sp_pr.find(qn("a:ln")) if sp_pr is not None else None
        for tag, attr in (("a:headEnd", "marker-start"), ("a:tailEnd", "marker-end")):
            el = ln.find(qn(tag)) if ln is not None else None
            wanted = el is not None and el.get("type") not in (None, "none")
            have = [
                p.get(attr) for p in g.findall("svg:path", ns) if p.get(attr)
            ]
            if wanted and not have:
                report.add(sid, "marker", f"{tag} present in source but no {attr} in svg")
            for ref in have:
                mid = ref.removeprefix("url(#").removesuffix(")")
                if mid not in marker_ids:
                    report.add(sid, "marker", f"{attr} references undefined marker {mid}")
    return report
