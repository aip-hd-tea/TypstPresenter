"""
Method F: tolerant structural fidelity between PPTX source and rendered PDF.

Method S only checks that the flow-mode output is *sane* (compiles, no
overlaps, one page per slide). Method F additionally checks that the output
still *corresponds* to the source -- not pixel-exact, but structurally:

- every visible source picture appears on its page, at a comparable
  relative size and in a comparable region of the page,
- the slide title is present and rendered at a size comparable to the
  resolved source size (placeholder inheritance included),
- title-layout slides (CENTER_TITLE) are actually centered,
- body text that matches the source is not rendered drastically smaller,
- hyperlinks survive the translation.

All checks carry generous tolerances: flow mode is *allowed* to move
content, the benchmark only flags translations a human would call wrong
(missing image, half-size fonts, right-column content landing left).
"""

from __future__ import annotations

import math
import re
from dataclasses import dataclass, field
from pathlib import Path

import fitz
import pptx

from typstpresenter.verify.geometry import EMU_PER_PT, BBox

# images smaller than this (pt^2) are decoration (icons inside diagram
# canvases); matching them against the many small PDF images is noise
_IMG_MIN_AREA = 2000.0
# rendered/source linear scale of an image, relative to the slide scale
_IMG_SCALE_ISSUE = 0.7
_IMG_SCALE_WARN = 0.8
# center drift as a fraction of the page width/height
_DX_ISSUE, _DX_WARN = 0.25, 0.13
_DY_ISSUE, _DY_WARN = 0.32, 0.18
# rendered/source font size ratios, relative to the slide scale
_TITLE_RATIO_ISSUE = 0.8
_BODY_RATIO_ISSUE = 0.55
_BODY_RATIO_WARN = 0.72
# uniform whole-slide shrink (flow calibration): relative layout is intact,
# so it is one slide-level finding, not one per element
_SLIDE_SCALE_ISSUE = 0.62
_SLIDE_SCALE_WARN = 0.82
# minimum normalized length for text matching (shorter strings collide)
_MATCH_MIN_CHARS = 6


@dataclass
class FidelityReport:
    pdf_path: Path
    issues: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return not self.issues

    def summary(self) -> str:
        status = "OK" if self.ok else f"{len(self.issues)} issues"
        lines = [f"{status} ({len(self.warnings)} warnings)"]
        lines += [f"  issue: {i}" for i in self.issues]
        lines += [f"  warning: {w}" for w in self.warnings]
        return "\n".join(lines)


# ------------------------------------------------------------ source facts --

@dataclass
class _SlideFacts:
    title: str = ""
    title_size: float | None = None
    centered_title: bool = False
    images: list[BBox] = field(default_factory=list)
    # (paragraph text, expected size in pt) of body text shapes
    texts: list[tuple[str, float]] = field(default_factory=list)
    links: set[str] = field(default_factory=set)
    # the source itself overstuffs a text box / the page; no layout can
    # render this at full size
    overfull: bool = False


def _estimated_text_height(shape, bbox: BBox, default_size: float,
                           font_scale: float) -> float:
    """Rough laid-out height of a text frame at PowerPoint line pitch."""
    from typstpresenter.convert.pptx_inherit import resolve_font_size_pt

    total = 0.0
    for paragraph in shape.text_frame.paragraphs:
        run = paragraph.runs[0] if paragraph.runs else None
        size = resolve_font_size_pt(run, paragraph, shape)
        size = (size if size is not None else default_size) * font_scale
        chars_per_line = max(bbox.w / (0.5 * size), 1.0)
        lines = max(1.0, math.ceil(len(paragraph.text) / chars_per_line))
        total += lines * size * 1.22
    return total


def _norm(text: str) -> str:
    return re.sub(r"\s+", " ", text.replace("\x0b", " ")).strip().casefold()


def extract_slide_facts(pptx_path: Path | str,
                        default_size: float = 18.0) -> list[_SlideFacts]:
    from typstpresenter.convert.cetz import is_diagram_shape
    from typstpresenter.convert.flow import _is_chrome, _off_page
    from typstpresenter.convert.pptx_inherit import resolve_font_size_pt
    from typstpresenter.convert.textbody import autofit_scales
    from typstpresenter.verify.pptx_geometry import (
        is_picture,
        iter_flat_shapes,
        picture_image,
    )

    prs = pptx.Presentation(str(pptx_path))
    page_w = prs.slide_width / EMU_PER_PT
    page_h = prs.slide_height / EMU_PER_PT
    facts = []
    for slide in prs.slides:
        f = _SlideFacts()
        for shape in slide.shapes:
            try:
                ph_type = shape.placeholder_format.type
            except (AttributeError, ValueError):
                ph_type = None
            if ph_type is not None and ph_type.name == "CENTER_TITLE":
                f.centered_title = True
        for shape, bbox in iter_flat_shapes(slide.shapes):
            is_title = shape == slide.shapes.title
            if not is_title and (_is_chrome(shape, page_h, bbox, page_w)
                                 or _off_page(bbox, page_w, page_h)):
                continue
            if is_picture(shape):
                image = picture_image(shape)
                if (image is not None and bbox.w * bbox.h >= _IMG_MIN_AREA
                        and image.ext.lower() != "svg"):
                    f.images.append(bbox)
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.text_frame.text.strip():
                continue
            font_scale, _ = autofit_scales(shape)
            if not is_title and bbox.h > 200:
                est = _estimated_text_height(shape, bbox, default_size,
                                             font_scale)
                if est > bbox.h * 1.15:
                    f.overfull = True
            for paragraph in shape.text_frame.paragraphs:
                text = _norm(paragraph.text)
                run = paragraph.runs[0] if paragraph.runs else None
                size = resolve_font_size_pt(run, paragraph, shape)
                size = (size if size is not None else default_size) * font_scale
                if is_title:
                    if text:
                        f.title = (f.title + " " + text).strip()
                        f.title_size = max(f.title_size or 0.0, size)
                elif len(text) >= _MATCH_MIN_CHARS and not is_diagram_shape(shape):
                    f.texts.append((text, size))
                for run in paragraph.runs:
                    address = run.hyperlink.address
                    if address:
                        f.links.add(address)
        facts.append(f)
    return facts


# --------------------------------------------------------------- PDF facts --

def _page_spans(page: fitz.Page) -> list[tuple[str, float, BBox]]:
    spans = []
    for block in page.get_text("dict")["blocks"]:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = _norm(span["text"])
                if text:
                    r = fitz.Rect(span["bbox"])
                    spans.append((text, span["size"],
                                  BBox(r.x0, r.y0, r.width, r.height)))
    return spans


def _page_images(page: fitz.Page) -> list[BBox]:
    boxes = []
    for block in page.get_text("dict")["blocks"]:
        if block.get("type") == 1:
            r = fitz.Rect(block["bbox"])
            boxes.append(BBox(r.x0, r.y0, r.width, r.height))
    return boxes


def _match_images(src: list[BBox], out: list[BBox],
                  page_w: float, page_h: float) -> list[tuple[BBox, BBox]]:
    """Greedy source->output assignment by center distance and scale."""
    pairs = []
    remaining = list(out)
    for sb in sorted(src, key=lambda b: -b.w * b.h):
        if not remaining:
            break
        def cost(ob: BBox) -> float:
            dx = abs(sb.center[0] - ob.center[0]) / page_w
            dy = abs(sb.center[1] - ob.center[1]) / page_h
            scale = ob.w / sb.w if sb.w else 1.0
            return dx + dy + abs(math.log(max(scale, 1e-3)))
        best = min(remaining, key=cost)
        remaining.remove(best)
        pairs.append((sb, best))
    return pairs


def _matched_text_size(text: str, spans: list[tuple[str, float, BBox]]
                       ) -> float | None:
    """Median rendered size of the spans matching a source paragraph.

    A span must cover a real part of the paragraph, otherwise short spans
    of other elements ("Ninjas" in a table cell) pollute the statistics of
    a paragraph mentioning the same word.
    """
    min_len = max(_MATCH_MIN_CHARS, 0.25 * len(text))
    sizes = sorted(size for span_text, size, _ in spans
                   if (len(span_text) >= min_len and span_text in text)
                   or text in span_text)
    return sizes[len(sizes) // 2] if sizes else None


# ------------------------------------------------------------------ checks --

def verify_fidelity(pptx_path: Path | str,
                    pdf_path: Path | str,
                    default_size: float = 18.0) -> FidelityReport:
    """Compare the compiled PDF against the PPTX source, with tolerances."""
    pdf_path = Path(pdf_path)
    report = FidelityReport(pdf_path=pdf_path)
    facts = extract_slide_facts(pptx_path, default_size)
    doc = fitz.open(pdf_path)
    if len(doc) != len(facts):
        report.issues.append(
            f"page count {len(doc)} != slide count {len(facts)}; "
            "skipping per-slide fidelity checks")
        return report

    for index, (f, page) in enumerate(zip(facts, doc)):
        n = index + 1
        page_w, page_h = page.rect.width, page.rect.height
        spans = _page_spans(page)
        images = _page_images(page)
        image_pairs = _match_images(f.images, images, page_w, page_h)

        # the slide's overall scale: flow calibration shrinks slides that
        # do not fit as flowing content *uniformly*, which preserves the
        # relative layout -- judge elements against the slide scale and
        # report heavy uniform shrink once per slide
        samples = [ob.w / sb.w for sb, ob in image_pairs if sb.w]
        for text, expected in f.texts:
            rendered = _matched_text_size(text, spans)
            if rendered and expected:
                samples.append(rendered / expected)
        samples.sort()
        slide_scale = min(samples[len(samples) // 2], 1.0) if samples else 1.0
        if slide_scale < _SLIDE_SCALE_ISSUE and not f.overfull:
            report.issues.append(
                f"slide {n}: content uniformly shrunk to {slide_scale:.0%} "
                "of the source scale (does not fit as flow)")
        elif slide_scale < _SLIDE_SCALE_WARN:
            note = " (source box is overstuffed)" if f.overfull else ""
            report.warnings.append(
                f"slide {n}: content uniformly shrunk to {slide_scale:.0%} "
                f"of the source scale{note}")
        # a shrunk slide necessarily re-positions content; keep flagging
        # relocations, but not as hard issues
        pos_issues = (report.issues if slide_scale >= _SLIDE_SCALE_WARN
                      else report.warnings)

        # --- title presence, size, centering
        if f.title:
            min_len = min(4, len(f.title))
            matched = [s for s in spans
                       if len(s[0]) >= min_len
                       and (s[0] in f.title or f.title in s[0])]
            if not matched:
                report.issues.append(f"slide {n}: title text not found in PDF")
            else:
                rendered = max(s[1] for s in matched)
                if f.title_size and rendered / f.title_size < _TITLE_RATIO_ISSUE:
                    report.issues.append(
                        f"slide {n}: title rendered at {rendered:.0f}pt, "
                        f"source is {f.title_size:.0f}pt")
                if f.centered_title:
                    tb = matched[0][2]
                    for _, size, bb in matched:
                        if size == rendered:
                            tb = bb
                            break
                    cx, cy = tb.center[0] / page_w, tb.center[1] / page_h
                    if abs(cx - 0.5) > 0.15 or not 0.18 < cy < 0.7:
                        report.issues.append(
                            f"slide {n}: title-slide heading not centered "
                            f"(center at {cx:.2f}, {cy:.2f})")

        # --- images: completeness, scale relative to slide scale, region
        if len(images) < len(f.images):
            report.issues.append(
                f"slide {n}: only {len(images)} of {len(f.images)} "
                "source images rendered")
        for sb, ob in image_pairs:
            scale = ob.w / sb.w if sb.w else 1.0
            relative = scale / slide_scale
            where = f"image at source ({sb.x:.0f},{sb.y:.0f})"
            if relative < _IMG_SCALE_ISSUE:
                report.issues.append(
                    f"slide {n}: {where} rendered at {scale:.0%} of its size "
                    f"(slide scale {slide_scale:.0%})")
            elif relative < _IMG_SCALE_WARN:
                report.warnings.append(
                    f"slide {n}: {where} rendered at {scale:.0%} of its size "
                    f"(slide scale {slide_scale:.0%})")
            dx = abs(sb.center[0] - ob.center[0]) / page_w
            dy = abs(sb.center[1] - ob.center[1]) / page_h
            if dx > _DX_ISSUE or dy > _DY_ISSUE:
                pos_issues.append(
                    f"slide {n}: {where} moved by ({dx:.0%}, {dy:.0%}) "
                    "of the page")
            elif dx > _DX_WARN or dy > _DY_WARN:
                report.warnings.append(
                    f"slide {n}: {where} moved by ({dx:.0%}, {dy:.0%}) "
                    "of the page")

        # --- body text size relative to the slide scale
        for text, expected in f.texts:
            rendered = _matched_text_size(text, spans)
            if not rendered or not expected:
                continue
            ratio = rendered / expected / slide_scale
            if ratio < _BODY_RATIO_ISSUE:
                report.issues.append(
                    f"slide {n}: text '{text[:32]}' rendered at "
                    f"{rendered:.0f}pt, source is {expected:.0f}pt "
                    f"(slide scale {slide_scale:.0%})")
            elif ratio < _BODY_RATIO_WARN:
                report.warnings.append(
                    f"slide {n}: text '{text[:32]}' rendered at "
                    f"{rendered:.0f}pt, source is {expected:.0f}pt "
                    f"(slide scale {slide_scale:.0%})")

        # --- hyperlinks
        uris = {link.get("uri") for link in page.get_links() if link.get("uri")}
        for url in sorted(f.links - uris):
            report.issues.append(f"slide {n}: hyperlink lost: {url}")

    return report
