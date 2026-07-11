"""
Resolve inherited text properties of PPTX placeholder shapes.

python-pptx only reports *explicit* run/paragraph formatting; placeholders
usually inherit font size, alignment and anchoring through the chain

    slide shape -> slide-layout placeholder -> slide-master placeholder
    -> master text styles (p:txStyles) -> spec default

This module walks that chain on the underlying XML. It is used by both the
emitter (to translate faithfully) and the ground-truth extraction (so the
verification compares against the same interpretation).
"""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

_ALGN = {
    "l": PP_ALIGN.LEFT,
    "ctr": PP_ALIGN.CENTER,
    "r": PP_ALIGN.RIGHT,
    "just": PP_ALIGN.JUSTIFY,
}
_ANCHOR = {"t": MSO_ANCHOR.TOP, "ctr": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}

# Placeholder types treated with the master's title style (idx-independent).
_TITLE_TYPE_NAMES = {"TITLE", "CENTER_TITLE"}


def _ph_type_name(shape) -> str | None:
    try:
        ph_type = shape.placeholder_format.type
    except (AttributeError, ValueError):
        return None
    return ph_type.name if ph_type is not None else None


def _inheritance_chain(shape):
    """Yield layout and master placeholder shapes this shape inherits from."""
    if not getattr(shape, "is_placeholder", False):
        return
    idx = shape.placeholder_format.idx
    type_name = _ph_type_name(shape)
    slide = shape.part.slide
    layout = getattr(slide, "slide_layout", None)
    if layout is None:
        return
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == idx or (
            type_name in _TITLE_TYPE_NAMES
            and _ph_type_name(ph) in _TITLE_TYPE_NAMES
        ):
            yield ph
            break
    master = layout.slide_master
    for ph in master.placeholders:
        master_type = _ph_type_name(ph)
        if type_name in _TITLE_TYPE_NAMES and master_type in _TITLE_TYPE_NAMES:
            yield ph
            break
        if type_name not in _TITLE_TYPE_NAMES and master_type == "BODY":
            yield ph
            break


def _props_from_lvl(lvl) -> dict:
    """Extract size/alignment/bullet properties from an a:lvlNpPr element."""
    props: dict = {}
    if lvl.get("algn"):
        props["algn"] = lvl.get("algn")
    defRPr = lvl.find(qn("a:defRPr"))
    if defRPr is not None and defRPr.get("sz"):
        props["sz"] = int(defRPr.get("sz")) / 100.0
    if lvl.find(qn("a:buNone")) is not None:
        props["bullet"] = None
    else:
        bu_char = lvl.find(qn("a:buChar"))
        if bu_char is not None:
            props["bullet"] = bu_char.get("char")
        elif lvl.find(qn("a:buAutoNum")) is not None:
            props["bullet"] = "1."
    spc = _space_before_of(lvl)
    if spc is not None:
        props["spc_bef"] = spc
    return props


def _space_before_of(pPr) -> tuple[str, float] | None:
    """spcBef as ("pct", fraction-of-line) or ("pt", points), if present."""
    spcBef = pPr.find(qn("a:spcBef"))
    if spcBef is None:
        return None
    pct = spcBef.find(qn("a:spcPct"))
    if pct is not None and pct.get("val"):
        return ("pct", int(pct.get("val")) / 100000.0)
    pts = spcBef.find(qn("a:spcPts"))
    if pts is not None and pts.get("val"):
        return ("pt", int(pts.get("val")) / 100.0)
    return None


def _lst_style_props(shape, level: int) -> dict:
    """Properties from a placeholder's own a:lstStyle."""
    txBody = shape.text_frame._txBody
    lst = txBody.find(qn("a:lstStyle"))
    if lst is None:
        return {}
    lvl = lst.find(qn(f"a:lvl{level + 1}pPr"))
    if lvl is None:
        return {}
    return _props_from_lvl(lvl)


def _master_txstyle_props(shape, level: int) -> dict:
    """Size/alignment from the slide master's p:txStyles."""
    type_name = _ph_type_name(shape)
    slide = shape.part.slide
    layout = getattr(slide, "slide_layout", None)
    if layout is None:
        return {}
    master = layout.slide_master
    txStyles = master.element.find(qn("p:txStyles"))
    if txStyles is None:
        return {}
    style_tag = "p:titleStyle" if type_name in _TITLE_TYPE_NAMES else "p:bodyStyle"
    style = txStyles.find(qn(style_tag))
    if style is None:
        return {}
    lvl = style.find(qn(f"a:lvl{level + 1}pPr"))
    if lvl is None:
        return {}
    return _props_from_lvl(lvl)


def _inherited_prop(shape, level: int, key: str):
    if not getattr(shape, "is_placeholder", False):
        return None
    for ancestor in _inheritance_chain(shape):
        props = _lst_style_props(ancestor, level)
        if key in props:
            return props[key]
    props = _master_txstyle_props(shape, level)
    return props.get(key)


def resolve_font_size_pt(run, paragraph, shape) -> float | None:
    """Effective font size in pt, or None if nothing in the chain sets one."""
    if run is not None and run.font.size is not None:
        return run.font.size.pt
    if paragraph.font.size is not None:
        return paragraph.font.size.pt
    return _inherited_prop(shape, paragraph.level, "sz")


def resolve_bullet(paragraph, shape) -> str | None:
    """Effective bullet character (or numbering approximation) if any."""
    from pptx.oxml.ns import qn as _qn

    pPr = paragraph._p.find(_qn("a:pPr"))
    if pPr is not None:
        if pPr.find(_qn("a:buNone")) is not None:
            return None
        bu_char = pPr.find(_qn("a:buChar"))
        if bu_char is not None:
            return bu_char.get("char")
        if pPr.find(_qn("a:buAutoNum")) is not None:
            return "1."
    return _inherited_prop(shape, paragraph.level, "bullet")


def resolve_alignment(paragraph, shape) -> PP_ALIGN | None:
    if paragraph.alignment is not None:
        return paragraph.alignment
    algn = _inherited_prop(shape, paragraph.level, "algn")
    return _ALGN.get(algn) if algn else None


def resolve_space_before(paragraph, shape) -> tuple[str, float] | None:
    """Effective paragraph space-before, explicit or inherited."""
    pPr = paragraph._p.find(qn("a:pPr"))
    if pPr is not None:
        spc = _space_before_of(pPr)
        if spc is not None:
            return spc
    return _inherited_prop(shape, paragraph.level, "spc_bef")


def resolve_autofit(shape) -> str:
    """Effective autofit mode: "shrink" (normAutofit), "resize" (spAutoFit)
    or "none". PowerPoint placeholders inherit this through bodyPr."""
    if not getattr(shape, "has_text_frame", False):
        return "none"

    def _fit_of(candidate) -> str | None:
        bodyPr = candidate.text_frame._txBody.find(qn("a:bodyPr"))
        if bodyPr is None:
            return None
        if bodyPr.find(qn("a:normAutofit")) is not None:
            return "shrink"
        if bodyPr.find(qn("a:spAutoFit")) is not None:
            return "resize"
        if bodyPr.find(qn("a:noAutofit")) is not None:
            return "none"
        return None

    fit = _fit_of(shape)
    if fit:
        return fit
    for ancestor in _inheritance_chain(shape):
        fit = _fit_of(ancestor)
        if fit:
            return fit
    return "none"


def resolve_anchor(shape) -> MSO_ANCHOR | None:
    """Effective vertical anchor of a text frame."""
    tf = shape.text_frame
    if tf.vertical_anchor is not None:
        return tf.vertical_anchor
    if not getattr(shape, "is_placeholder", False):
        return None
    for ancestor in _inheritance_chain(shape):
        bodyPr = ancestor.text_frame._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None and bodyPr.get("anchor"):
            return _ANCHOR.get(bodyPr.get("anchor"))
    return None
