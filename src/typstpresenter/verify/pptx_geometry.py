"""
Extract element geometry (ground truth) from a PPTX file via python-pptx.
"""

from __future__ import annotations

from pathlib import Path

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from typstpresenter.verify.geometry import (
    EMU_PER_PT,
    BBox,
    DocGeometry,
    ElementGeometry,
    ElementKind,
    SlideGeometry,
)


def element_id(slide_index: int, shape_id: int) -> str:
    """Stable element id shared between PPTX ground truth and Typst probes."""
    return f"s{slide_index}-e{shape_id}"


def is_picture(shape) -> bool:
    """True for any shape carrying an image.

    Pictures inserted into a content placeholder keep shape_type
    PLACEHOLDER, but their XML element is a ``p:pic`` like any picture.
    Movies are ``p:pic`` too; they render as their poster frame.
    OLE control objects and graphic frames containing OLE objects often
    contain a nested preview image as a p:pic/a:blip.
    """
    from pptx.oxml.ns import qn

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.element.tag == qn("p:pic"):
        return True

    is_ole = False
    if shape.shape_type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
        is_ole = True
    elif shape.element.tag.endswith("graphicFrame") and "oleObj" in shape.element.xml:
        is_ole = True

    if is_ole:
        try:
            if list(shape.element.iter(qn("a:blip"))):
                return True
        except Exception:
            pass
    return False


def picture_image(shape):
    """The Image part of a picture shape (a movie's poster frame, an OLE
    preview image), or None."""
    from pptx.oxml.ns import qn

    try:
        return shape.image
    except (AttributeError, ValueError, KeyError):
        pass
    try:
        return shape.poster_frame
    except AttributeError:
        pass

    is_ole = False
    if shape.shape_type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
        is_ole = True
    elif shape.element.tag.endswith("graphicFrame") and "oleObj" in shape.element.xml:
        is_ole = True

    if is_ole:
        try:
            for blip in shape.element.iter(qn("a:blip")):
                rId = blip.get(qn("r:embed"))
                if rId and rId in shape.part.related_parts:
                    part = shape.part.related_parts[rId]
                    if hasattr(part, "image"):
                        return part.image
        except Exception:
            pass
    return None


def _kind_of(shape) -> ElementKind:
    st = shape.shape_type
    if is_picture(shape):
        return ElementKind.IMAGE
    if st == MSO_SHAPE_TYPE.GROUP:
        return ElementKind.GROUP
    if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        return ElementKind.SHAPE
    if st in (MSO_SHAPE_TYPE.LINE,) or shape.shape_type is None and shape.element.tag.endswith("cxnSp"):
        return ElementKind.CONNECTOR
    if getattr(shape, "has_table", False):
        return ElementKind.TABLE
    if shape.has_text_frame:
        return ElementKind.TEXT
    return ElementKind.OTHER


def _shape_text(shape) -> str:
    if getattr(shape, "has_table", False):
        return "\n".join(
            cell.text.strip() for row in shape.table.rows for cell in row.cells
        ).strip()
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def _is_left_top_aligned(shape) -> bool:
    """True if the text starts at the box's top-left corner (the only case
    where ink-anchor checks in Method A are meaningful)."""
    from typstpresenter.convert.pptx_inherit import resolve_alignment, resolve_anchor

    if not shape.has_text_frame:
        return True
    if resolve_anchor(shape) not in (None, MSO_ANCHOR.TOP):
        return False
    return all(
        resolve_alignment(p, shape) in (None, PP_ALIGN.LEFT)
        for p in shape.text_frame.paragraphs
    )


def _autofit_of(shape) -> str:
    from typstpresenter.convert.pptx_inherit import resolve_autofit

    return resolve_autofit(shape)


def _bbox_of(shape) -> BBox | None:
    if shape.left is None or shape.top is None:
        return None
    return BBox(
        x=shape.left / EMU_PER_PT,
        y=shape.top / EMU_PER_PT,
        w=(shape.width or 0) / EMU_PER_PT,
        h=(shape.height or 0) / EMU_PER_PT,
    )


# --------------------------------------------------------- group flattening --

# Affine per-axis transform (sx, sy, tx, ty) mapping child EMU coordinates
# to absolute EMU: abs = (x * sx + tx, y * sy + ty).
_IDENTITY = (1.0, 1.0, 0.0, 0.0)


def _group_child_transform(group, outer=_IDENTITY):
    """Compose the group's child coordinate mapping with an outer transform.

    Children of a group live in their own coordinate space defined by
    a:chOff/a:chExt and are scaled into the group's a:off/a:ext.
    """
    from pptx.oxml.ns import qn

    grp_sp_pr = group.element.find(qn("p:grpSpPr"))
    xfrm = grp_sp_pr.find(qn("a:xfrm")) if grp_sp_pr is not None else None
    if xfrm is None or shape_off_missing(group):
        return outer
    ch_off = xfrm.find(qn("a:chOff"))
    ch_ext = xfrm.find(qn("a:chExt"))
    ox, oy = group.left, group.top
    ew, eh = group.width or 0, group.height or 0
    cx0 = int(ch_off.get("x")) if ch_off is not None else 0
    cy0 = int(ch_off.get("y")) if ch_off is not None else 0
    cw = int(ch_ext.get("cx")) if ch_ext is not None else ew
    ch = int(ch_ext.get("cy")) if ch_ext is not None else eh
    sxl = ew / cw if cw else 1.0
    syl = eh / ch if ch else 1.0
    txl = ox - cx0 * sxl
    tyl = oy - cy0 * syl
    sx, sy, tx, ty = outer
    return (sx * sxl, sy * syl, sx * txl + tx, sy * tyl + ty)


def shape_off_missing(shape) -> bool:
    return shape.left is None or shape.top is None


def iter_flat_shapes(shapes, transform=_IDENTITY):
    """Yield (shape, absolute BBox in pt) with groups flattened recursively."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            child_transform = _group_child_transform(shape, transform)
            yield from iter_flat_shapes(shape.shapes, child_transform)
            continue
        if shape_off_missing(shape):
            continue
        sx, sy, tx, ty = transform
        yield shape, BBox(
            x=(shape.left * sx + tx) / EMU_PER_PT,
            y=(shape.top * sy + ty) / EMU_PER_PT,
            w=((shape.width or 0) * sx) / EMU_PER_PT,
            h=((shape.height or 0) * sy) / EMU_PER_PT,
        )


def extract_pptx_geometry(path: Path | str) -> DocGeometry:
    """
    Read a *.pptx file and return the geometry of all placed shapes.

    Shapes without a position (e.g. unplaced placeholders) are skipped.
    Group shapes are recorded as a single element (no recursion into
    children, whose coordinates live in a scaled group space).
    """
    prs = pptx.Presentation(str(path))
    slide_w = prs.slide_width / EMU_PER_PT
    slide_h = prs.slide_height / EMU_PER_PT

    doc = DocGeometry(source=str(path))
    for slide_index, slide in enumerate(prs.slides):
        sg = SlideGeometry(index=slide_index, width=slide_w, height=slide_h)
        for shape, bbox in iter_flat_shapes(slide.shapes):
            kind = _kind_of(shape)
            text = _shape_text(shape)
            # Empty text boxes / empty placeholders render nothing.
            if kind in (ElementKind.TEXT, ElementKind.OTHER) and not text:
                continue
            meta = {
                "shape_name": shape.name,
                "align_left_top": _is_left_top_aligned(shape),
                "autofit": _autofit_of(shape),
            }
            if kind == ElementKind.IMAGE:
                image = picture_image(shape)
                ext = image.ext.lower() if image is not None else ""
                # typst cannot render these; the emitter converts or draws a
                # placeholder, so Method A must not require image ink
                meta["unrenderable"] = ext not in (
                    "png", "jpg", "jpeg", "gif", "svg", "webp")
            if kind == ElementKind.SHAPE:
                # shapes without fill and outline leave no ink; only their
                # text (if any) is visually verifiable
                from typstpresenter.convert.pptx_style import (
                    shape_fill_rgb,
                    shape_line_rgb,
                )
                meta["invisible"] = (
                    shape_fill_rgb(shape) is None and shape_line_rgb(shape) is None
                )
            sg.elements.append(
                ElementGeometry(
                    kind=kind,
                    bbox=bbox,
                    id=element_id(slide_index, shape.shape_id),
                    text=text,
                    meta=meta,
                )
            )
        doc.slides.append(sg)
    return doc
