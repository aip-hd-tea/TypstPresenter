"""
Baseline PPTX -> Touying emitter.

This is deliberately the simplest possible faithful translation: every
shape is placed at its absolute PPTX coordinates on a Touying slide whose
page size equals the PPTX slide size (margin 0), so PPTX pt coordinates
map 1:1 to Typst pt coordinates. It serves two purposes:

1. It produces paired (PPTX, Typst) documents to exercise and evaluate
   the verification tools (Methods A and B).
2. It is the seed of the real converter: later development replaces
   absolute placement with idiomatic Touying/CeTZ/Fletcher constructs,
   while the verification tools guard the visual fidelity.

Fault injection deliberately corrupts the emitted geometry (shift, resize,
extra text, dropped elements) to measure whether the verification methods
detect the corresponding problems.

Autoshapes and connectors are emitted as one CeTZ canvas per slide with
node probes, so diagram translation to CeTZ is covered as well.
"""

from __future__ import annotations

from dataclasses import dataclass, replace
from pathlib import Path

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from typstpresenter.verify.geometry import EMU_PER_PT, BBox
from typstpresenter.verify.method_b import PROBE_PRELUDE
from typstpresenter.verify.pptx_geometry import element_id, iter_flat_shapes
from typstpresenter.verify.pptx_inherit import (
    resolve_alignment,
    resolve_anchor,
    resolve_autofit,
    resolve_bullet,
    resolve_font_size_pt,
    resolve_space_before,
)
from typstpresenter.verify.pptx_style import (
    shape_fill_rgb,
    shape_font_rgb,
    shape_line_rgb,
    shape_line_width_pt,
)

TOUYING_VERSION = "0.6.1"
CETZ_VERSION = "0.5.2"

# Every Typst markup-active character must be escaped: braces/brackets and
# inline markers, but also '/' (starts a comment after another '/'), and the
# line-start markers =,-,+ (heading/list/enum). Escaping also suppresses
# smart-dash substitution -- PPTX text is literal.
_ESCAPE = str.maketrans({c: f"\\{c}" for c in '\\#$*_`@<>"~[]{}/=-+\''})


def escape_typst(text: str) -> str:
    return text.translate(_ESCAPE)


@dataclass(frozen=True)
class Fault:
    """A deliberate corruption of one emitted element (for evaluation)."""
    element_id: str
    dx: float = 0.0        # shift in pt
    dy: float = 0.0
    scale_w: float = 1.0   # box resize factors
    scale_h: float = 1.0
    extra_text: str = ""   # appended to the content (provokes overflow)
    drop: bool = False     # omit the element entirely

    def expected_issues(self) -> set[str]:
        """Issue kinds a verifier with access to box geometry (Method B)
        should report for this fault."""
        expected = set()
        if self.drop:
            return {"missing"}
        if self.dx or self.dy:
            expected.add("moved")
        if self.scale_w != 1.0 or self.scale_h != 1.0:
            expected.add("resized")
        if self.extra_text or self.scale_h < 1.0:
            expected.add("overflow")
        return expected

    def expected_issues_ink(self) -> set[str]:
        """Issue kinds detectable from rendered ink alone (Method A).

        Boxes without fill or stroke leave no trace in the PDF, so pure
        box resizing (and the overflow it causes within the original
        bounds) is invisible to ink-based verification.
        """
        expected = set()
        if self.drop:
            return {"missing"}
        if self.dx or self.dy:
            expected.add("moved")
        if self.extra_text:
            expected.add("overflow")
        return expected


# Line metrics for the emitted font stack (Calibri): PowerPoint's single
# line pitch is ascent+descent+lineGap ~ 1.22 em; typst renders a line as
# cap-height+descender ~ 0.632 em plus `leading`. Matching the two gives
# the leading that reproduces PowerPoint's vertical rhythm.
PPT_LINE_PITCH_EM = 1.22
TYPST_LINE_HEIGHT_EM = 0.632


def _autofit_scales(shape) -> tuple[float, float]:
    """(fontScale, lnSpcReduction) from a:normAutofit, defaults (1, 0)."""
    from pptx.oxml.ns import qn

    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    fit = bodyPr.find(qn("a:normAutofit")) if bodyPr is not None else None
    if fit is None:
        return 1.0, 0.0
    font_scale = int(fit.get("fontScale") or 100000) / 100000.0
    lnspc_red = int(fit.get("lnSpcReduction") or 0) / 100000.0
    return font_scale, lnspc_red


def _run_size_pt(run, paragraph, shape, default: float, scale: float = 1.0) -> float:
    resolved = resolve_font_size_pt(run, paragraph, shape)
    return (resolved if resolved is not None else default) * scale


def _run_color(run) -> str | None:
    color = run.font.color
    try:
        if color and color.type is not None and color.rgb is not None:
            return str(color.rgb)
    except AttributeError:
        pass
    return None


def _run_style(run, paragraph, shape, default_size: float, scale: float = 1.0,
               default_color: str | None = None) -> tuple:
    """Style signature of a run: (size, bold, italic, underline, color)."""
    return (
        round(_run_size_pt(run, paragraph, shape, default_size, scale), 2),
        bool(run.font.bold),
        bool(run.font.italic),
        bool(run.font.underline),
        _run_color(run) or default_color,
    )


def _styled_text(style: tuple, text: str) -> str:
    size, bold, italic, underline, rgb = style
    args = [f"size: {size:g}pt"]
    if bold:
        args.append('weight: "bold"')
    if italic:
        args.append('style: "italic"')
    if rgb:
        args.append(f'fill: rgb("#{rgb}")')
    inner = escape_typst(text)
    if underline:
        inner = f"#underline[{inner}]"
    return f"#text({', '.join(args)})[{inner}]"


def _paragraph_runs_markup(paragraph, shape, default_size: float,
                           scale: float = 1.0,
                           default_color: str | None = None) -> list[str]:
    """Runs of a paragraph as Typst markup, merging equal-styled neighbors.

    Merging matters beyond cleanliness: a chain of separate #text() calls
    gives the line breaker a break opportunity at every run boundary, which
    shreds formula-styled text (many tiny runs) into one word per line.
    """
    from pptx.oxml.ns import qn

    chunks: list[tuple[str, tuple]] = []
    for run in paragraph.runs:
        if not run.text:
            continue
        style = _run_style(run, paragraph, shape, default_size, scale, default_color)
        if chunks and chunks[-1][1] == style:
            chunks[-1] = (chunks[-1][0] + run.text, style)
        else:
            chunks.append((run.text, style))
    # fields (slide number, date) are not exposed as runs; emit their
    # cached text with the paragraph's default styling
    for fld in paragraph._p.findall(qn("a:fld")):
        t = fld.find(qn("a:t"))
        if t is not None and t.text:
            size = round(_run_size_pt(None, paragraph, shape, default_size, scale), 2)
            chunks.append((t.text, (size, False, False, False, default_color)))
    return [_styled_text(style, text) for text, style in chunks]


def _typst_align(paragraph, shape) -> str | None:
    from pptx.enum.text import PP_ALIGN

    alignment = resolve_alignment(paragraph, shape)
    return {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}.get(alignment)


def _shape_bbox(shape) -> BBox | None:
    if shape.left is None or shape.top is None:
        return None
    return BBox(
        shape.left / EMU_PER_PT,
        shape.top / EMU_PER_PT,
        (shape.width or 0) / EMU_PER_PT,
        (shape.height or 0) / EMU_PER_PT,
    )


def _emit_text_body(shape, default_size: float, extra_text: str,
                    extra_scale: float = 1.0) -> str:
    """Render the paragraphs of a text frame as Typst markup."""
    from pptx.enum.text import MSO_ANCHOR

    font_scale, lnspc_red = _autofit_scales(shape)
    font_scale *= extra_scale
    leading_em = max(
        PPT_LINE_PITCH_EM * (1.0 - lnspc_red) - TYPST_LINE_HEIGHT_EM, 0.1
    )
    parts: list[str] = [
        f"#set par(leading: {leading_em:.3f}em, spacing: {leading_em:.3f}em)"
    ]
    first_paragraph = True
    for paragraph in shape.text_frame.paragraphs:
        runs = _paragraph_runs_markup(paragraph, shape, default_size, font_scale)
        if not runs:
            continue
        para_size = _run_size_pt(
            paragraph.runs[0] if paragraph.runs else None,
            paragraph, shape, default_size, font_scale,
        )
        spc_bef = resolve_space_before(paragraph, shape)
        if spc_bef and not first_paragraph:
            kind, value = spc_bef
            if kind == "pct":
                gap = value * PPT_LINE_PITCH_EM * para_size  # para_size is scaled
            else:
                # PowerPoint's autofit scales point-based paragraph spacing
                # along with the text; without this, many-paragraph bodies
                # keep a fixed spacing floor and shrink cannot converge
                gap = value * font_scale * (1.0 - lnspc_red)
            parts.append(f"#v({gap:.2f}pt)")
        prefix = ""
        indent = paragraph.level * 18.0
        if indent:
            prefix += f"#h({indent:g}pt)"
        bullet = resolve_bullet(paragraph, shape)
        if bullet:
            if len(bullet) == 1 and 0xE000 <= ord(bullet) <= 0xF8FF:
                bullet = "•"
            # size the bullet like its paragraph so it doesn't set the pitch
            prefix += f"#text(size: {para_size:.4g}pt)[{escape_typst(bullet)} ]"
        line = f"#par(hanging-indent: 0pt)[{prefix}{''.join(runs)}]"
        align = _typst_align(paragraph, shape)
        if align:
            line = f"#align({align})[{line}]"
        parts.append(line)
        first_paragraph = False
    if extra_text:
        parts.append(f"#par[#text(size: {default_size:g}pt)[{escape_typst(extra_text)}]]")
    content = "\n".join(parts)
    anchor = resolve_anchor(shape)
    if anchor == MSO_ANCHOR.MIDDLE:
        content = f"#align(horizon)[\n{content}\n]"
    elif anchor == MSO_ANCHOR.BOTTOM:
        content = f"#align(bottom)[\n{content}\n]"
    return content


def _is_diagram_shape(shape) -> bool:
    return shape.shape_type in (
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.FREEFORM,  # approximated by its bounding box for now
        MSO_SHAPE_TYPE.LINE,
    ) or shape.element.tag.endswith("}cxnSp")


# formats the typst `image` function can load directly
_TYPST_IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "svg", "webp"}


def _emit_picture(shape, eid: str, media_dir: Path) -> str:
    """Write the picture blob and return typst markup for it.

    Formats typst cannot load (wmf, emf, tiff, bmp) are converted to PNG
    via Pillow when possible; otherwise a placeholder frame keeps the
    geometry intact.
    """
    image = shape.image
    ext = image.ext.lower()
    if ext in _TYPST_IMAGE_EXTS:
        filename = f"{eid}.{ext}"
        (media_dir / filename).write_bytes(image.blob)
        return f'#image("{media_dir.name}/{filename}", width: 100%, height: 100%)'
    try:
        import io

        from PIL import Image as PILImage

        with PILImage.open(io.BytesIO(image.blob)) as im:
            filename = f"{eid}.png"
            im.convert("RGBA").save(media_dir / filename)
        return f'#image("{media_dir.name}/{filename}", width: 100%, height: 100%)'
    except Exception:
        return ('#rect(width: 100%, height: 100%, stroke: 0.5pt + gray)'
                f'// unsupported image format: {ext}')


def _emit_table(shape, default_size: float) -> str:
    """A PPTX table as a Typst table with the exact column/row extents."""
    table = shape.table
    columns = ", ".join(f"{c.width / EMU_PER_PT:.2f}pt" for c in table.columns)
    rows = ", ".join(f"{r.height / EMU_PER_PT:.2f}pt" for r in table.rows)
    cells = []
    for row in table.rows:
        for cell in row.cells:
            runs = []
            for paragraph in cell.text_frame.paragraphs:
                runs += _paragraph_runs_markup(paragraph, shape, default_size)
            cells.append(f"  [{''.join(runs)}],")
    header = (f"#table(\n  columns: ({columns}),\n  rows: ({rows}),\n"
              "  inset: 3pt, stroke: 0.5pt,\n")
    return header + "\n".join(cells) + "\n)"


def _shape_flips(shape) -> tuple[bool, bool]:
    from pptx.oxml.ns import qn

    xfrm = shape.element.spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return False, False
    return xfrm.get("flipH") == "1", xfrm.get("flipV") == "1"


def _connector_is_elbow(shape) -> bool:
    from pptx.oxml.ns import qn

    prst_geom = shape.element.spPr.find(qn("a:prstGeom"))
    prst = prst_geom.get("prst") if prst_geom is not None else ""
    return bool(prst) and prst.startswith("bentConnector")


def _cetz_style_args(shape) -> str:
    """fill/stroke arguments resolved from explicit colors or the theme."""
    fill = shape_fill_rgb(shape)
    stroke_rgb = shape_line_rgb(shape)
    fill_arg = f'fill: rgb("#{fill}")' if fill else "fill: none"
    if stroke_rgb:
        width = shape_line_width_pt(shape)
        stroke_arg = f'stroke: (paint: rgb("#{stroke_rgb}"), thickness: {width:g}pt)'
    else:
        stroke_arg = "stroke: none"
    return f"{fill_arg}, {stroke_arg}"


def _cetz_label_markup(shape, eid: str, probes: bool, default_size: float,
                       label_w: float) -> str:
    """Shape label as width-constrained box (multi-paragraph, styled, probed)."""
    parts = []
    if shape.has_text_frame:
        # shape text defaults to the style's font color (often white)
        font_rgb = shape_font_rgb(shape)
        font_scale, lnspc_red = _autofit_scales(shape)
        for paragraph in shape.text_frame.paragraphs:
            runs = _paragraph_runs_markup(paragraph, shape, default_size,
                                          scale=font_scale,
                                          default_color=font_rgb)
            if runs:
                parts.append(f"#par[{''.join(runs)}]")
        if parts:
            # PowerPoint line metrics, as in _emit_text_body: without this,
            # multi-paragraph labels get typst's default paragraph spacing
            # and grow far taller than the shape
            leading = max(
                PPT_LINE_PITCH_EM * (1.0 - lnspc_red) - TYPST_LINE_HEIGHT_EM, 0.1
            )
            parts.insert(
                0, f"#set par(leading: {leading:.3f}em, spacing: {leading:.3f}em);"
            )
    content = "".join(parts)
    boxed = f"box(width: {label_w:.2f}pt, align(center)[{content}])" if content else ""
    if not probes:
        return f"#{boxed}" if boxed else ""
    if boxed:
        # label probe measures the constrained label so overflow beyond the
        # shape can be classified as a source condition, not a translation bug
        return f'#tp-label-probe("{eid}", {boxed})'
    return f'#tp-node-probe("{eid}")'


def _emit_cetz_shape(shape, eid: str, bbox: BBox, probes: bool,
                     default_size: float = 18.0) -> str:
    """One PPTX autoshape/connector as CeTZ drawing commands (y axis flipped)."""
    from pptx.enum.shapes import MSO_SHAPE

    lines = []
    x1, y1, x2, y2 = bbox.x, -bbox.y, bbox.x2, -bbox.y2
    cx, cy = bbox.center
    is_connector = shape.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE,)

    if is_connector:
        flip_h, flip_v = _shape_flips(shape)
        bx, ex = (x2, x1) if flip_h else (x1, x2)
        by, ey = (y2, y1) if flip_v else (y1, y2)
        stroke_rgb = shape_line_rgb(shape) or "000000"
        width = shape_line_width_pt(shape)
        stroke = f'stroke: (paint: rgb("#{stroke_rgb}"), thickness: {width:g}pt)'
        if _connector_is_elbow(shape):
            mx = (bx + ex) / 2
            points = (f"({bx:.2f}, {by:.2f}), ({mx:.2f}, {by:.2f}), "
                      f"({mx:.2f}, {ey:.2f}), ({ex:.2f}, {ey:.2f})")
        else:
            points = f"({bx:.2f}, {by:.2f}), ({ex:.2f}, {ey:.2f})"
        lines.append(f"    line({points}, mark: (end: \">\"), {stroke})")
        if probes:
            lines.append(f'    content(({cx:.2f}, {-cy:.2f}), [#tp-node-probe("{eid}")])')
        return "\n".join(lines), None

    auto = None
    try:
        auto = shape.auto_shape_type
    except (ValueError, AttributeError):
        pass
    style = _cetz_style_args(shape)
    if auto == MSO_SHAPE.OVAL:
        lines.append(
            f"    circle(({cx:.2f}, {-cy:.2f}), "
            f"radius: ({bbox.w / 2:.2f}, {bbox.h / 2:.2f}), {style})"
        )
    elif auto == MSO_SHAPE.DIAMOND:
        lines.append(
            f"    line(({cx:.2f}, {y1:.2f}), ({x2:.2f}, {-cy:.2f}), "
            f"({cx:.2f}, {y2:.2f}), ({x1:.2f}, {-cy:.2f}), close: true, {style})"
        )
    elif auto == MSO_SHAPE.ISOSCELES_TRIANGLE:
        lines.append(
            f"    line(({cx:.2f}, {y1:.2f}), ({x2:.2f}, {y2:.2f}), "
            f"({x1:.2f}, {y2:.2f}), close: true, {style})"
        )
    else:
        radius = 4.0 if auto == MSO_SHAPE.ROUNDED_RECTANGLE else 0.0
        lines.append(
            f"    rect(({x1:.2f}, {y1:.2f}), ({x2:.2f}, {y2:.2f}), "
            f"radius: {radius:g}, {style})"
        )
    # constrain the label to the shape width so long text wraps like
    # PowerPoint instead of spreading beyond the shape
    label_w = max(bbox.w - 7.2, 7.2)
    label = _cetz_label_markup(shape, eid, probes, default_size, label_w)
    expectation = None
    if label:
        from pptx.enum.text import MSO_ANCHOR

        # respect the text frame's vertical anchor: top/bottom-anchored
        # panels must not be centered (their text grows from the edge).
        # The OOXML default when no anchor is set anywhere is "t" (top).
        anchor = resolve_anchor(shape) if shape.has_text_frame else None
        if anchor == MSO_ANCHOR.BOTTOM:
            lines.append(
                f'    content(({cx:.2f}, {y2 + 3.6:.2f}), anchor: "south", [{label}])'
            )
            expectation = (cx, bbox.y2 - 3.6, "south")
        elif anchor == MSO_ANCHOR.MIDDLE:
            lines.append(f"    content(({cx:.2f}, {-cy:.2f}), [{label}])")
            expectation = (cx, cy, "center")
        else:  # OOXML default: top
            lines.append(
                f'    content(({cx:.2f}, {y1 - 3.6:.2f}), anchor: "north", [{label}])'
            )
            expectation = (cx, bbox.y + 3.6, "north")
    return "\n".join(lines), expectation


def emit_touying(
    pptx_path: Path | str,
    out_path: Path | str,
    probes: bool = True,
    faults: tuple[Fault, ...] = (),
    default_font_size: float = 18.0,
    minimal: bool = False,
) -> Path:
    """
    Emit a Touying presentation with absolute positioning from a PPTX file.

    Every text/image element becomes a ``tp-probe(...)`` call (designed
    geometry included), autoshapes/connectors become one CeTZ canvas per
    slide with ``tp-node-probe`` markers. With ``probes=False`` the same
    layout is emitted without instrumentation.
    """
    pptx_path = Path(pptx_path)
    out_path = Path(out_path)
    if minimal:
        probes = False
    prs = pptx.Presentation(str(pptx_path))
    page_w = prs.slide_width / EMU_PER_PT
    page_h = prs.slide_height / EMU_PER_PT
    faults_by_id = {f.element_id: f for f in faults}

    margin_cfg = "margin: (top: 3em, bottom: 0pt, left: 0pt, right: 0pt)" if minimal else "margin: 0pt"
    header_cfg = "" if minimal else "  header: none,"
    header = [
        "// Auto-generated by typstpresenter baseline emitter -- edit freely.",
        f'#import "@preview/touying:{TOUYING_VERSION}": *',
        "#import themes.simple: *",
        f'#import "@preview/cetz:{CETZ_VERSION}"',
        "",
        "#show: simple-theme.with(",
        f"  config-page(width: {page_w:g}pt, height: {page_h:g}pt, {margin_cfg}),",
        "  config-common(handout: true),",
        header_cfg,
        "  footer: none,",
        ")",
        '#set text(font: ("Calibri", "Arial", "Liberation Sans"))',
        # PPTX text is literal; PowerPoint does not substitute quotes at render
        "#set smartquote(enabled: false)",
        "",
    ]
    if probes:
        header.append(PROBE_PRELUDE)

    media_dir = out_path.parent / f"{out_path.stem}_media"
    # elements whose bodyPr requests shrink-on-overflow; PowerPoint computes
    # the actual fontScale at edit time, so we must derive it ourselves
    autofit_ids: set[str] = set()
    # label probe id -> (slide, expected label center x, expected anchor y,
    # anchor kind); used to measure and correct per-slide canvas drift
    label_expect: dict[str, tuple] = {}

    def _get_slide_title(slide) -> str:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text.strip().replace("\n", " ")
        return ""

    def build_body(shrink: dict[str, float],
                   canvas_adjust: dict[int, tuple[float, float]]) -> list[str]:
        body: list[str] = []
        for slide_index, slide in enumerate(prs.slides):
            title = _get_slide_title(slide)
            if minimal and title:
                body.append(f"== {escape_typst(title)}")
            else:
                body.append("== ")

            cetz_lines: list[str] = []
            # shapes parked outside the slide inflate the canvas bounds and
            # would shift the whole canvas; track the overhang to compensate
            overhang_left = overhang_top = 0.0
            for shape, bbox in iter_flat_shapes(slide.shapes):
                if minimal and shape == slide.shapes.title:
                    continue
                eid = element_id(slide_index, shape.shape_id)
                fault = faults_by_id.get(eid)
                if fault:
                    if fault.drop:
                        continue
                    bbox = replace(
                        bbox,
                        x=bbox.x + fault.dx, y=bbox.y + fault.dy,
                        w=bbox.w * fault.scale_w, h=bbox.h * fault.scale_h,
                    )

                if _is_diagram_shape(shape):
                    markup, expectation = _emit_cetz_shape(shape, eid, bbox, probes)
                    cetz_lines.append(markup)
                    if expectation is not None:
                        label_expect[eid] = (slide_index, *expectation)
                    overhang_left = min(overhang_left, bbox.x)
                    overhang_top = min(overhang_top, bbox.y)
                    continue

                if getattr(shape, "has_table", False):
                    content = _emit_table(shape, default_font_size)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    media_dir.mkdir(parents=True, exist_ok=True)
                    content = _emit_picture(shape, eid, media_dir)
                elif shape.has_text_frame and shape.text_frame.text.strip():
                    if resolve_autofit(shape) == "shrink":
                        autofit_ids.add(eid)
                    content = _emit_text_body(
                        shape, default_size=default_font_size,
                        extra_text=fault.extra_text if fault else "",
                        extra_scale=shrink.get(eid, 1.0),
                    )
                else:
                    continue

                geometry = f"{bbox.x:.2f}pt, {bbox.y:.2f}pt, {bbox.w:.2f}pt, {bbox.h:.2f}pt"
                if probes:
                    body.append(f'  #tp-probe("{eid}", {geometry})[\n{content}\n  ]')
                else:
                    body.append(
                        f"  #place(top + left, dx: {bbox.x:.2f}pt, dy: {bbox.y:.2f}pt, "
                        f"block(width: {bbox.w:.2f}pt, height: {bbox.h:.2f}pt)[\n{content}\n  ])"
                    )

            if cetz_lines:
                adjust_x, adjust_y = canvas_adjust.get(slide_index, (0.0, 0.0))
                body.append(
                    f"  #place(top + left, dx: {overhang_left + adjust_x:.2f}pt, "
                    f"dy: {overhang_top + adjust_y:.2f}pt, cetz.canvas(length: 1pt, {{"
                )
                body.append("    import cetz.draw: *")
                body.append("    set-style(stroke: 0.75pt)")
                # invisible full-page rect anchors the canvas at the page
                # origin, otherwise CeTZ shrink-wraps it and everything shifts
                body.append(f"    rect((0, 0), ({page_w:.2f}, {-page_h:.2f}), stroke: none)")
                body.append("\n".join(cetz_lines))
                body.append("  }))")
            body.append("")
        return body

    shrink: dict[str, float] = {}
    canvas_adjust: dict[int, tuple[float, float]] = {}

    def write() -> None:
        out_path.write_text(
            "\n".join(header + build_body(shrink, canvas_adjust)),
            encoding="utf-8", newline="\n",
        )

    write()

    # Calibration: like PowerPoint at edit time, measure the laid-out content
    # via the probes and (a) shrink autofit text until it fits its box,
    # (b) correct per-slide canvas drift caused by labels or shapes whose
    # ink extends beyond the page (cetz sizes the canvas to its content).
    if probes and (autofit_ids or label_expect):
        from typstpresenter.verify.typst_tools import query

        for _ in range(8):
            measured = query(out_path, "<tp-probe>").value

            needed = {
                p["id"]: p["h"] / p["content-h"]
                for p in measured
                if "w" in p and p["id"] in autofit_ids
                and p["content-h"] > p["h"] + 0.5
            }
            for eid, ratio in needed.items():
                # extra 3% margin compounds across rounds, absorbing the
                # sub-linear response of reflowing text
                shrink[eid] = shrink.get(eid, 1.0) * ratio * 0.97

            drift: dict[int, list[tuple[float, float]]] = {}
            for p in measured:
                expect = label_expect.get(p["id"])
                if expect is None or "label-h" not in p:
                    continue
                slide_index, expect_cx, expect_y, kind = expect
                measured_cx = p["x"] + p["label-w"] / 2
                if kind == "north":
                    measured_y = p["y"]
                elif kind == "south":
                    measured_y = p["y"] + p["label-h"]
                else:
                    measured_y = p["y"] + p["label-h"] / 2
                drift.setdefault(slide_index, []).append(
                    (measured_cx - expect_cx, measured_y - expect_y))
            adjusted = False
            for slide_index, deltas in drift.items():
                dxs = sorted(d[0] for d in deltas)
                dys = sorted(d[1] for d in deltas)
                dx, dy = dxs[len(dxs) // 2], dys[len(dys) // 2]
                if abs(dx) > 0.5 or abs(dy) > 0.5:
                    prev = canvas_adjust.get(slide_index, (0.0, 0.0))
                    canvas_adjust[slide_index] = (prev[0] - dx, prev[1] - dy)
                    adjusted = True

            if not needed and not adjusted:
                break
            write()

    return out_path
