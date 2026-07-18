"""
Generate the evaluation corpus: paired (PPTX, Typst) documents.

Covers:
(i)  simple presentations with common layouts (title+content, two columns,
     2x2 grid),
(ii) diagram slides translated to CeTZ (via the emitter) and Fletcher
     (dedicated generator),
plus fault-injected variants with known ground truth, used to measure the
detection quality of verification Methods A and B.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

import pptx
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt

from typstpresenter.convert.emitter import Fault, emit_touying, escape_typst
from typstpresenter.verify.geometry import EMU_PER_PT, DocGeometry
from typstpresenter.verify.pptx_geometry import extract_pptx_geometry

FLETCHER_VERSION = "0.5.8"


@dataclass
class CorpusCase:
    name: str
    pptx_path: Path
    typ_path: Path
    truth: DocGeometry
    # element id -> issue kinds Method B (box geometry) must report;
    # empty for clean cases
    expected_issues_b: dict[str, set[str]] = field(default_factory=dict)
    # element id -> issue kinds detectable from rendered ink (Method A)
    expected_issues_a: dict[str, set[str]] = field(default_factory=dict)
    # injected text exists that is legitimately reported as 'extra' by A
    allows_extra_text: bool = False
    # dense real decks exceed Method A's heuristic matching; they are
    # gated by Method B only (see docs/verification-methods.md)
    verify_with_a: bool = True
    kind: str = "layout"  # layout | diagram-cetz | diagram-fletcher


# ------------------------------------------------------------ PPTX builders --

def _add_textbox(slide, x_in, y_in, w_in, h_in, text: str, size_pt: float,
                 bold: bool = False):
    from pptx.util import Inches
    box = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    tf.text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
        if not paragraph.runs:
            paragraph.font.size = Pt(size_pt)
    return box


LOREM = (
    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Aenean commodo "
    "ligula eget dolor. Aenean massa. Cum sociis natoque penatibus et magnis "
    "dis parturient montes, nascetur ridiculus mus."
)


def build_title_content(path: Path) -> None:
    prs = pptx.Presentation()  # default 10 x 7.5 in
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_textbox(slide, 0.5, 0.4, 9.0, 1.0, "Title and Content", 36, bold=True)
    _add_textbox(slide, 0.5, 1.8, 9.0, 5.0,
                 f"First point about the topic\nSecond point with more words\n{LOREM}", 18)
    prs.save(str(path))


def build_two_columns(path: Path) -> None:
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.4, 9.0, 1.0, "Two Columns", 36, bold=True)
    _add_textbox(slide, 0.5, 1.8, 4.25, 4.8, f"Left column.\n{LOREM}", 16)
    _add_textbox(slide, 5.25, 1.8, 4.25, 4.8, f"Right column.\n{LOREM}", 16)
    prs.save(str(path))


def build_grid(path: Path) -> None:
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.3, 9.0, 0.9, "2x2 Grid", 32, bold=True)
    cells = [(0.5, 1.5), (5.25, 1.5), (0.5, 4.4), (5.25, 4.4)]
    for i, (x, y) in enumerate(cells):
        _add_textbox(slide, x, y, 4.25, 2.6,
                     f"Cell {i + 1}. Short paragraph of body text for cell {i + 1}.", 16)
    # second slide: same grid, longer texts
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide2, 0.5, 0.3, 9.0, 0.9, "2x2 Grid, fuller", 32, bold=True)
    for i, (x, y) in enumerate(cells):
        _add_textbox(slide2, x, y, 4.25, 2.6, f"Cell {i + 1}. {LOREM[:140]}", 16)
    prs.save(str(path))


def build_flowchart(path: Path) -> None:
    """Horizontal flowchart: three rounded boxes joined by two connectors."""
    from pptx.util import Inches
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.4, 9.0, 0.9, "Simple Flowchart", 32, bold=True)
    labels = ["Start", "Process", "End"]
    boxes = []
    for i, label in enumerate(labels):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + i * 3.2), Inches(3.0), Inches(1.8), Inches(1.0),
        )
        shape.text_frame.text = label
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
        boxes.append(shape)
    for a, b in zip(boxes, boxes[1:]):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            a.left + a.width, a.top + a.height // 2,
            b.left, b.top + b.height // 2,
        )
        conn.line.width = Pt(1.5)
    prs.save(str(path))


def build_mixed_shapes(path: Path) -> None:
    """Rectangle, oval and a line: basic shape vocabulary for CeTZ."""
    from pptx.util import Inches
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.4, 9.0, 0.9, "Mixed Shapes", 32, bold=True)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(2.5), Inches(1.5))
    rect.text_frame.text = "Box"
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(6.0), Inches(4.0), Inches(2.0), Inches(2.0))
    oval.text_frame.text = "Circle"
    for shape in (rect, oval):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        rect.left + rect.width, rect.top + rect.height // 2,
        oval.left, oval.top + oval.height // 2,
    )
    prs.save(str(path))


def build_rich_text(path: Path) -> None:
    """Level 1: run styling, alignments, vertical anchoring and bullets."""
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = _add_textbox(slide, 0.5, 0.3, 9.0, 0.9, "Rich Text Features", 32, bold=True)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # mixed run styling within one paragraph
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(2.2))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    for text, bold, italic, underline, color in [
        ("Plain, ", False, False, False, None),
        ("bold, ", True, False, False, None),
        ("italic, ", False, True, False, None),
        ("underlined, ", False, False, True, None),
        ("and colored words", False, False, False, RGBColor(0xC0, 0x00, 0x00)),
    ]:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        if color is not None:
            run.font.color.rgb = color
    p2 = box.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "A second paragraph in the same box."
    run.font.size = Pt(16)

    # bulleted list with two levels
    bullets = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(2.2))
    bullets.text_frame.word_wrap = True
    items = [("First bullet item", 0), ("Second bullet item", 0),
             ("Nested sub-item", 1), ("Third bullet item", 0)]
    for i, (text, level) in enumerate(items):
        p = bullets.text_frame.paragraphs[0] if i == 0 else bullets.text_frame.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(16)
        p.level = level
        pPr = p._p.get_or_add_pPr()
        bu = pPr.makeelement(qn("a:buChar"), {"char": "•"})
        pPr.append(bu)

    # right-aligned and center-aligned paragraphs
    aligned = _add_textbox(slide, 0.5, 4.2, 4.25, 2.0,
                           "Centered line\nRight-aligned line\nLeft-aligned line", 16)
    aligned.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    aligned.text_frame.paragraphs[1].alignment = PP_ALIGN.RIGHT

    # vertically centered content
    middle = _add_textbox(slide, 5.25, 4.2, 4.25, 2.0,
                          "Vertically centered text", 16)
    middle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    prs.save(str(path))


def build_styled_shapes(path: Path) -> None:
    """Level 3: filled/stroked shapes, theme colors, elbow connector.

    - rounded rect keeps its default theme fill (accent1 + white text)
    - diamond gets an explicit red fill
    - triangle: explicit green fill, no outline
    - rect: no fill, thick blue outline
    - elbow connector and a flipped straight connector
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.3, 9.0, 0.8, "Styled Shapes", 28, bold=True)

    themed = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(2.2), Inches(1.2))
    themed.text_frame.text = "Theme fill"

    diamond = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(4.2), Inches(1.3), Inches(1.8), Inches(1.6))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = RGBColor(0xB0, 0x20, 0x20)
    diamond.text_frame.text = "?"

    triangle = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(7.2), Inches(1.4), Inches(1.6), Inches(1.4))
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(0x20, 0x80, 0x30)
    triangle.line.fill.background()

    outlined = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(4.5), Inches(1.8), Inches(1.2))
    outlined.fill.background()
    outlined.line.color.rgb = RGBColor(0x20, 0x40, 0xB0)
    outlined.line.width = Pt(3)
    outlined.text_frame.text = "Outline"

    for shape in (themed, diamond, outlined):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

    elbow = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW,
        themed.left + themed.width, themed.top + themed.height // 2,
        diamond.left, diamond.top + diamond.height // 2,
    )
    elbow.line.color.rgb = RGBColor(0x40, 0x40, 0x40)
    # flipped straight connector: drawn bottom-up (end above begin)
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        outlined.left + outlined.width // 2, outlined.top,
        diamond.left + diamond.width // 2, diamond.top + diamond.height,
    )
    prs.save(str(path))


def build_decision_flowchart(path: Path) -> None:
    """Level 4: 2D flowchart -- diamond decision, branch and back-edge.

    Grid-aligned (column centers 1.5/4.5/7.5 in, row centers 2.5/5.0 in) so
    a grid-based Fletcher translation can reproduce the geometry.
    """
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.5, 0.3, 9.0, 0.8, "Decision Flow", 28, bold=True)

    def add_node(shape_type, col_center, row_center, w, h, label):
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(col_center - w / 2), Inches(row_center - h / 2),
            Inches(w), Inches(h),
        )
        shape.text_frame.text = label
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)
        return shape

    start = add_node(MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 2.5, 1.8, 0.9, "Start")
    check = add_node(MSO_SHAPE.DIAMOND, 4.5, 2.5, 1.8, 1.4, "OK?")
    done = add_node(MSO_SHAPE.ROUNDED_RECTANGLE, 7.5, 2.5, 1.8, 0.9, "Done")
    fix = add_node(MSO_SHAPE.RECTANGLE, 4.5, 5.0, 1.8, 0.9, "Fix it")

    def connect(kind, a, a_side, b, b_side):
        def point(shape, side):
            x = {"l": shape.left, "r": shape.left + shape.width,
                 "c": shape.left + shape.width // 2}[side[0]]
            y = {"t": shape.top, "b": shape.top + shape.height,
                 "m": shape.top + shape.height // 2}[side[1]]
            return x, y
        ax, ay = point(a, a_side)
        bx, by = point(b, b_side)
        return slide.shapes.add_connector(kind, ax, ay, bx, by)

    connect(MSO_CONNECTOR.STRAIGHT, start, "rm", check, "lm")
    connect(MSO_CONNECTOR.STRAIGHT, check, "rm", done, "lm")   # yes branch
    connect(MSO_CONNECTOR.STRAIGHT, check, "cb", fix, "ct")    # no branch, vertical
    connect(MSO_CONNECTOR.ELBOW, fix, "lm", start, "cb")       # back-edge
    prs.save(str(path))


def build_rotated_shapes(path: Path) -> None:
    """Level 7: shape rotation (PPTX rot attribute, clockwise degrees)."""
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title and content
    slide.shapes.title.text = "Rotated Shapes"

    def add_rotated(col_center, row_center, w, h, deg, label, shape_type=MSO_SHAPE.RECTANGLE):
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(col_center - w / 2), Inches(row_center - h / 2),
            Inches(w), Inches(h),
        )
        shape.rotation = deg
        shape.text_frame.text = label
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
        return shape

    add_rotated(2.2, 2.8, 2.0, 1.0, 30, "30")
    add_rotated(7.2, 2.8, 2.0, 1.0, -45, "-45")
    add_rotated(4.7, 5.2, 2.5, 1.2, 90, "90", MSO_SHAPE.ROUNDED_RECTANGLE)
    prs.save(str(path))


def build_freeform_shapes(path: Path) -> None:
    """Level 7: freeform (custGeom) polygons -- a chevron arrow and a star.

    Both are pure straight-edged paths (moveTo/lnTo/close only), the case
    the emitter draws as an exact CeTZ polygon rather than falling back to
    the shape's bounding box.
    """
    import math

    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Freeform Shapes"

    scale = Emu(Inches(1)) / 1000
    fb = slide.shapes.build_freeform(start_x=0, start_y=300, scale=scale)
    fb.add_line_segments([
        (600, 300), (600, 0), (1000, 500), (600, 1000), (600, 700), (0, 700),
    ], close=True)
    arrow = fb.convert_to_shape(origin_x=Inches(1.5), origin_y=Inches(2.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xC0, 0x50, 0x20)
    arrow.line.color.rgb = RGBColor(0x40, 0x40, 0x40)

    cx, cy, r_out, r_in = 500, 500, 500, 200
    pts = []
    for i in range(10):
        r = r_out if i % 2 == 0 else r_in
        angle = -math.pi / 2 + i * math.pi / 5
        pts.append((round(cx + r * math.cos(angle)), round(cy + r * math.sin(angle))))
    fb2 = slide.shapes.build_freeform(start_x=pts[0][0], start_y=pts[0][1], scale=scale)
    fb2.add_line_segments(pts[1:], close=True)
    star = fb2.convert_to_shape(origin_x=Inches(6.5), origin_y=Inches(2.5))
    star.fill.solid()
    star.fill.fore_color.rgb = RGBColor(0x20, 0x80, 0xC0)
    star.line.color.rgb = RGBColor(0x40, 0x40, 0x40)
    prs.save(str(path))


def build_placeholder_deck(path: Path) -> None:
    """Level 2: real slide layouts; all styling inherited from the master.

    No run sets an explicit font size, alignment or anchor -- the emitter
    must resolve them through layout/master/txStyles.
    """
    prs = pptx.Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
    slide.shapes.title.text = "Inherited Styles Deck"
    slide.placeholders[1].text = "Subtitle set without any explicit formatting"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # title and content
    slide2.shapes.title.text = "Bullets from the Master"
    body = slide2.placeholders[1].text_frame
    body.text = "Top level point with a reasonable amount of words"
    p = body.add_paragraph()
    p.text = "Second top level point"
    p = body.add_paragraph()
    p.text = "An indented sub-point below it"
    p.level = 1
    p = body.add_paragraph()
    p.text = "Another sub-point on level two"
    p.level = 2
    prs.save(str(path))


# ------------------------------------------------------- Fletcher generator --

def _cluster_centers(values: list[float], tol: float = 20.0) -> list[float]:
    """Cluster 1D coordinates; returns sorted cluster centers."""
    centers: list[list[float]] = []
    for v in sorted(values):
        if centers and v - centers[-1][-1] <= tol:
            centers[-1].append(v)
        else:
            centers.append([v])
    return [sum(c) / len(c) for c in centers]


def _grid_index(value: float, centers: list[float]) -> int:
    return min(range(len(centers)), key=lambda i: abs(centers[i] - value))


def emit_fletcher_diagram(pptx_path: Path, out_path: Path) -> Path:
    """
    Emit the diagram slide of a PPTX as a 2D Fletcher diagram.

    Unlike the CeTZ emitter (absolute coordinates), Fletcher lays nodes out
    on an elastic grid. Shapes are assigned grid cells by clustering their
    centers into columns/rows; node boxes are forced to the PPTX shape
    sizes and the spacing is derived from the PPTX gaps, so positions
    should land close to the truth -- how close is exactly what the
    verification measures. Edges are inferred from connector endpoints.
    """
    from pptx.enum.shapes import MSO_SHAPE

    from typstpresenter.verify.geometry import ElementKind
    from typstpresenter.verify.method_b import PROBE_PRELUDE

    truth = extract_pptx_geometry(pptx_path)
    prs = pptx.Presentation(str(pptx_path))
    page_w = prs.slide_width / EMU_PER_PT
    page_h = prs.slide_height / EMU_PER_PT

    slide = truth.slides[0]
    shapes = [e for e in slide.elements if e.kind == ElementKind.SHAPE]
    texts = [e for e in slide.elements if e.kind == ElementKind.TEXT]

    # auto-shape types (for diamond etc.) from the pptx side, keyed by id
    auto_types = {}
    for shape in prs.slides[0].shapes:
        try:
            auto_types[f"s0-e{shape.shape_id}"] = shape.auto_shape_type
        except (ValueError, AttributeError):
            pass

    col_centers = _cluster_centers([e.bbox.center[0] for e in shapes])
    row_centers = _cluster_centers([e.bbox.center[1] for e in shapes])
    cell_of = {
        e.id: (_grid_index(e.bbox.center[0], col_centers),
               _grid_index(e.bbox.center[1], row_centers))
        for e in shapes
    }

    def _spacing(centers: list[float], extents: dict[int, float]) -> float:
        gaps = [
            centers[i + 1] - centers[i] - (extents.get(i, 0) + extents.get(i + 1, 0)) / 2
            for i in range(len(centers) - 1)
        ]
        return sum(gaps) / len(gaps) if gaps else 27.0

    # Effective fletcher node size: fletcher's diamond shape circumscribes
    # the node rect, so a diamond's node size must be the inscribed rect
    # (half the PPTX bbox) for the drawn outline to match the PPTX shape.
    def _node_size(e) -> tuple[float, float]:
        if auto_types.get(e.id) == MSO_SHAPE.DIAMOND:
            return e.bbox.w / 2, e.bbox.h / 2
        return e.bbox.w, e.bbox.h

    col_widths: dict[int, float] = {}
    row_heights: dict[int, float] = {}
    for e in shapes:
        col, row = cell_of[e.id]
        w, h = _node_size(e)
        col_widths[col] = max(col_widths.get(col, 0), w)
        row_heights[row] = max(row_heights.get(row, 0), h)
    spacing_x = _spacing(col_centers, col_widths)
    spacing_y = _spacing(row_centers, row_heights)

    # edges: match connector endpoints to the nearest shape center
    def _nearest_shape(x: float, y: float) -> str:
        return min(
            shapes,
            key=lambda e: (e.bbox.center[0] - x) ** 2 + (e.bbox.center[1] - y) ** 2,
        ).id

    shapes_by_id = {e.id: e for e in shapes}
    edges = []
    for shape in prs.slides[0].shapes:
        eid = f"s0-e{shape.shape_id}"
        if not any(e.id == eid and e.kind == ElementKind.CONNECTOR
                   for e in slide.elements):
            continue
        bx, by = shape.begin_x / EMU_PER_PT, shape.begin_y / EMU_PER_PT
        ex, ey = shape.end_x / EMU_PER_PT, shape.end_y / EMU_PER_PT
        node_a, node_b = _nearest_shape(bx, by), _nearest_shape(ex, ey)
        # L-shaped connectors need a corner vertex in the grid; whether the
        # route runs horizontally or vertically first follows from which
        # edge of the source shape the connector leaves
        (col_a, row_a), (col_b, row_b) = cell_of[node_a], cell_of[node_b]
        corner = None
        if col_a != col_b and row_a != row_b:
            a_box = shapes_by_id[node_a].bbox
            leaves_horizontally = abs(bx - a_box.center[0]) >= a_box.w / 2 - 1
            corner = (col_b, row_a) if leaves_horizontally else (col_a, row_b)
        edges.append((eid, cell_of[node_a], corner, cell_of[node_b]))

    lines = [
        "// Auto-generated Fletcher pairing (2D grid inferred from PPTX).",
        '#import "@preview/touying:0.6.1": *',
        "#import themes.simple: *",
        f'#import "@preview/fletcher:{FLETCHER_VERSION}" as fletcher: diagram, node, edge',
        "",
        "#show: simple-theme.with(",
        f"  config-page(width: {page_w:g}pt, height: {page_h:g}pt, margin: 0pt),",
        "  config-common(handout: true),",
        "  footer: none,",
        ")",
        '#set text(font: ("Calibri", "Arial", "Liberation Sans"))',
        "",
        PROBE_PRELUDE,
        "#slide[",
    ]
    for el in texts:
        b = el.bbox
        lines.append(
            f'  #tp-probe("{el.id}", {b.x:.2f}pt, {b.y:.2f}pt, {b.w:.2f}pt, {b.h:.2f}pt)'
            f"[#text(size: 28pt, weight: \"bold\")[{escape_typst(el.text)}]]"
        )
    origin_x = min(e.bbox.x for e in shapes)
    origin_y = min(e.bbox.y for e in shapes)
    lines.append(
        f"  #place(top + left, dx: {origin_x:.2f}pt + TP-CAL-DX, "
        f"dy: {origin_y:.2f}pt + TP-CAL-DY, diagram("
    )
    lines.append("    node-inset: 0pt,")
    lines.append(f"    spacing: ({spacing_x:.2f}pt, {spacing_y:.2f}pt),")
    lines.append("    node-stroke: 0.75pt,")
    for el in shapes:
        col, row = cell_of[el.id]
        is_diamond = auto_types.get(el.id) == MSO_SHAPE.DIAMOND
        # fit: 1 draws the diamond exactly circumscribing the (half-size)
        # node rect, i.e. with the PPTX shape's bounding box
        shape_arg = ("shape: fletcher.shapes.diamond.with(fit: 1)"
                     if is_diamond else "corner-radius: 4pt")
        w, h = _node_size(el)
        label = f'#tp-node-probe("{el.id}")#text(size: 12pt)[{escape_typst(el.text)}]'
        lines.append(
            f"    node(({col}, {row}), [{label}], {shape_arg}, "
            f"width: {w:.2f}pt, height: {h:.2f}pt),"
        )
    for conn_id, ca, corner, cb in edges:
        probe = f'[#tp-node-probe("{conn_id}")]'
        via = f"({corner[0]}, {corner[1]}), " if corner else ""
        lines.append(
            f'    edge(({ca[0]}, {ca[1]}), {via}({cb[0]}, {cb[1]}), "-|>", '
            f"label: {probe}, label-sep: 0pt),"
        )
    lines.append("  ))")
    lines.append("]")
    source = "\n".join(lines)

    # Calibration pass: fletcher's outer bounding box (which `place`
    # anchors) does not coincide with the node hull, so the diagram lands
    # with a constant offset. Compile once with zero correction, read the
    # probe of an anchor node (a rect node's label origin equals its node
    # rect origin) and re-emit with the measured delta.
    from typstpresenter.verify.typst_tools import query

    anchor = next((e for e in shapes
                   if auto_types.get(e.id) != MSO_SHAPE.DIAMOND), shapes[0])
    out_path.write_text(
        source.replace("TP-CAL-DX", "0pt").replace("TP-CAL-DY", "0pt"),
        encoding="utf-8", newline="\n",
    )
    probes = {p["id"]: p for p in query(out_path, "<tp-probe>").value}
    measured = probes.get(anchor.id)
    if measured is not None:
        dx = anchor.bbox.x - measured["x"]
        dy = anchor.bbox.y - measured["y"]
        out_path.write_text(
            source.replace("TP-CAL-DX", f"{dx:.2f}pt").replace("TP-CAL-DY", f"{dy:.2f}pt"),
            encoding="utf-8", newline="\n",
        )
    return out_path


# ------------------------------------------------------------ corpus driver --

# Static registry: name -> (pptx builder, case kind). Cases added by the
# autoresearch loop go here; tests parametrize over these names.
BUILDERS = {
    "layout_title_content": (build_title_content, "layout"),
    "layout_two_columns": (build_two_columns, "layout"),
    "layout_grid": (build_grid, "layout"),
    "diagram_flowchart": (build_flowchart, "diagram-cetz"),
    "diagram_mixed_shapes": (build_mixed_shapes, "diagram-cetz"),
    # autoresearch level 1: rich text
    "layout_rich_text": (build_rich_text, "layout"),
    # autoresearch level 2: placeholder layouts with inherited styles
    "layout_placeholders": (build_placeholder_deck, "layout"),
    # autoresearch level 3: styled shapes, theme colors, elbow connectors
    "diagram_styled_shapes": (build_styled_shapes, "diagram-cetz"),
    # autoresearch level 4: 2D decision flowchart with back-edge
    "diagram_decision": (build_decision_flowchart, "diagram-cetz"),
    # autoresearch level 7: shape rotation, freeform (custGeom) polygons
    "diagram_rotated": (build_rotated_shapes, "diagram-cetz"),
    "diagram_freeform": (build_freeform_shapes, "diagram-cetz"),
}

# diagram cases that additionally get a Fletcher pairing
FLETCHER_PAIRINGS = ("diagram_flowchart", "diagram_decision")

FAULT_VARIANTS = ("moved", "overflow", "resized", "missing", "extra_text")

# autoresearch level 5: real presentations from tests/data (translated as
# a whole; verification must accept them without any issue)
EXTERNAL_CASES = ("simple", "two_content", "multi_content", "media", "talk_example_a")

# autoresearch level 6: dense lecture decks; Method B gate only
EXTERNAL_CASES_B_ONLY = ("IBN_presentations/vlxN04-ibn",)


def clean_case_names(with_external: bool = True) -> list[str]:
    names = [*BUILDERS, *(f"{name}_fletcher" for name in FLETCHER_PAIRINGS)]
    if with_external:
        names += [f"data_{stem}" for stem in EXTERNAL_CASES]
        names += [f"data_{Path(rel).stem}" for rel in EXTERNAL_CASES_B_ONLY]
    return names


def fault_case_names() -> list[str]:
    return [f"layout_two_columns_faulty_{v}" for v in FAULT_VARIANTS]


def _find_id(truth: DocGeometry, text_prefix: str) -> str:
    for _, el in truth.all_elements():
        if el.text.startswith(text_prefix):
            assert el.id is not None
            return el.id
    raise KeyError(f"no element starting with '{text_prefix}'")


def generate_corpus(out_dir: Path | str,
                    external_data_dir: Path | str | None = None) -> list[CorpusCase]:
    """Build all PPTX files, paired Typst files and fault variants.

    ``external_data_dir`` points to a directory with real presentations
    (see EXTERNAL_CASES); they are translated as-is and must verify clean.
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    cases: list[CorpusCase] = []

    for name, (builder, kind) in BUILDERS.items():
        pptx_path = out_dir / f"{name}.pptx"
        builder(pptx_path)
        truth = extract_pptx_geometry(pptx_path)
        typ_path = out_dir / f"{name}.typ"
        emit_touying(pptx_path, typ_path)
        cases.append(CorpusCase(name=name, pptx_path=pptx_path, typ_path=typ_path,
                                truth=truth, kind=kind))

    # Fletcher pairings of selected diagram cases
    for name in FLETCHER_PAIRINGS:
        diagram_pptx = out_dir / f"{name}.pptx"
        fletcher_typ = out_dir / f"{name}_fletcher.typ"
        emit_fletcher_diagram(diagram_pptx, fletcher_typ)
        cases.append(CorpusCase(
            name=f"{name}_fletcher", pptx_path=diagram_pptx,
            typ_path=fletcher_typ, truth=extract_pptx_geometry(diagram_pptx),
            kind="diagram-fletcher",
        ))

    # Fault-injected variants of the two-column layout
    two_col_pptx = out_dir / "layout_two_columns.pptx"
    truth = extract_pptx_geometry(two_col_pptx)
    title_id = _find_id(truth, "Two Columns")
    left_id = _find_id(truth, "Left column.")
    right_id = _find_id(truth, "Right column.")

    fault_sets: dict[str, tuple[Fault, ...]] = {
        "faulty_moved": (Fault(title_id, dx=25.0, dy=15.0),),
        "faulty_overflow": (Fault(left_id, scale_h=0.25),),
        "faulty_resized": (Fault(right_id, scale_w=0.55),),
        "faulty_missing": (Fault(left_id, drop=True),),
        "faulty_extra_text": (Fault(
            right_id,
            extra_text="This extra sentence was injected to trigger a text "
                       "overflow beyond the designated placeholder box, with "
                       "some padding words to make absolutely sure. "
                       + LOREM + " " + LOREM,
        ),),
    }
    # real presentations, translated as-is
    if external_data_dir is not None:
        external_data_dir = Path(external_data_dir)
        externals = [(stem, True) for stem in EXTERNAL_CASES]
        externals += [(rel, False) for rel in EXTERNAL_CASES_B_ONLY]
        for rel, with_a in externals:
            src = external_data_dir / f"{rel}.pptx"
            if not src.exists():
                continue
            stem = Path(rel).stem
            pptx_path = out_dir / f"data_{stem}.pptx"
            pptx_path.write_bytes(src.read_bytes())
            typ_path = out_dir / f"data_{stem}.typ"
            emit_touying(pptx_path, typ_path)
            cases.append(CorpusCase(
                name=f"data_{stem}", pptx_path=pptx_path, typ_path=typ_path,
                truth=extract_pptx_geometry(pptx_path), kind="external",
                verify_with_a=with_a,
            ))

    for variant, faults in fault_sets.items():
        typ_path = out_dir / f"layout_two_columns_{variant}.typ"
        emit_touying(two_col_pptx, typ_path, faults=faults)
        cases.append(CorpusCase(
            name=f"layout_two_columns_{variant}", pptx_path=two_col_pptx,
            typ_path=typ_path, truth=truth,
            expected_issues_b={f.element_id: f.expected_issues() for f in faults},
            expected_issues_a={f.element_id: f.expected_issues_ink() for f in faults},
            allows_extra_text=any(f.extra_text for f in faults),
            kind="layout",
        ))
    return cases
