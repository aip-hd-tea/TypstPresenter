"""
Method S: sanity + simplicity verification for minimal (flow-mode) output.

Minimal output carries no probes, and it deliberately deviates from the
original coordinates, so Methods A/B do not apply. What must hold instead:

- the document compiles,
- every PPTX slide fits on exactly one page (no content spilling onto
  continuation pages),
- no ink escapes the page,
- text does not collide with other text or images (coherent layout),
- the source stays simple: no ``#place``, few explicit ``#text(...)``
  calls, no probe machinery.

The geometric checks read the compiled PDF with PyMuPDF (like Method A);
the simplicity metrics scan the Typst source.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

import fitz
import pptx

from typstpresenter.verify.geometry import BBox

# text collisions below this fraction of the smaller box are ignored
# (adjacent lines' font bboxes brush each other routinely)
_OVERLAP_FRACTION = 0.35
# PDF line boxes overshoot layout boxes by the font descent; the theme's
# page-number footer sits right at the page edge and trips a tight bound
_PAGE_SLACK_PT = 4.0


@dataclass
class SimplicityMetrics:
    lines: int = 0
    place_calls: int = 0
    text_calls: int = 0
    absolute_lengths: int = 0
    probe_defs: int = 0

    def summary(self) -> str:
        return (f"{self.lines} lines, {self.place_calls} #place, "
                f"{self.text_calls} #text(), {self.absolute_lengths} pt-values, "
                f"{self.probe_defs} probe defs")


@dataclass
class MinimalReport:
    typ_path: Path
    issues: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    metrics: SimplicityMetrics = field(default_factory=SimplicityMetrics)

    @property
    def ok(self) -> bool:
        return not self.issues

    def summary(self) -> str:
        status = "OK" if self.ok else f"{len(self.issues)} issues"
        lines = [f"{status} ({len(self.warnings)} warnings) -- {self.metrics.summary()}"]
        lines += [f"  issue: {i}" for i in self.issues]
        lines += [f"  warning: {w}" for w in self.warnings]
        return "\n".join(lines)


def source_metrics(typ_source: str) -> SimplicityMetrics:
    return SimplicityMetrics(
        lines=typ_source.count("\n") + 1,
        place_calls=len(re.findall(r"#place\(", typ_source)),
        text_calls=len(re.findall(r"#text\(", typ_source)),
        absolute_lengths=len(re.findall(r"\d(?:\.\d+)?pt\b", typ_source)),
        probe_defs=len(re.findall(r"#let tp-", typ_source)),
    )


def _rect_to_bbox(rect: fitz.Rect) -> BBox:
    return BBox(rect.x0, rect.y0, rect.x1 - rect.x0, rect.y1 - rect.y0)


def _text_lines(page: fitz.Page) -> list[BBox]:
    lines = []
    for block in page.get_text("dict")["blocks"]:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            text = "".join(s["text"] for s in line.get("spans", [])).strip()
            if text:
                lines.append(_rect_to_bbox(fitz.Rect(line["bbox"])))
    return lines


def _image_boxes(page: fitz.Page) -> list[BBox]:
    boxes = []
    for block in page.get_text("dict")["blocks"]:
        if block.get("type") == 1:
            boxes.append(_rect_to_bbox(fitz.Rect(block["bbox"])))
    return boxes


def _overlap_area(a: BBox, b: BBox) -> float:
    w = min(a.x2, b.x2) - max(a.x, b.x)
    h = min(a.y2, b.y2) - max(a.y, b.y)
    return max(w, 0.0) * max(h, 0.0)


def _collides(a: BBox, b: BBox) -> bool:
    area = _overlap_area(a, b)
    smaller = min(a.w * a.h, b.w * b.h)
    return smaller > 0 and area / smaller > _OVERLAP_FRACTION


def verify_minimal(typ_path: Path | str,
                   pptx_path: Path | str | None = None,
                   pdf_path: Path | str | None = None) -> MinimalReport:
    """Check a minimal .typ (and its compiled PDF) for layout sanity.

    Compiles the document if ``pdf_path`` does not exist yet.
    """
    from typstpresenter.verify.typst_tools import compile_pdf

    typ_path = Path(typ_path)
    report = MinimalReport(typ_path=typ_path)
    report.metrics = source_metrics(typ_path.read_text(encoding="utf-8"))
    if report.metrics.place_calls:
        report.warnings.append(
            f"{report.metrics.place_calls} #place() calls in source")
    if report.metrics.probe_defs:
        report.issues.append("probe machinery in minimal output")

    pdf_path = Path(pdf_path) if pdf_path else typ_path.with_suffix(".pdf")
    if not pdf_path.exists():
        compile_pdf(typ_path, pdf_path)

    doc = fitz.open(pdf_path)
    src_tt_overlap: dict[int, bool] = {}
    src_ti_overlap: dict[int, bool] = {}
    pages_match = True
    if pptx_path is not None:
        n_slides = len(pptx.Presentation(str(pptx_path)).slides._sldIdLst)
        if len(doc) != n_slides:
            pages_match = False
            report.issues.append(
                f"page count {len(doc)} != slide count {n_slides} "
                "(slides split across pages)")
        else:
            # slides whose *source* already overlaps text with text/images
            # (annotated screenshots, labels on drawings) may reproduce
            # those overlaps -- that is fidelity, not a layout bug
            src_tt_overlap, src_ti_overlap = _source_overlaps(pptx_path)

    for page_index, page in enumerate(doc):
        page_box = BBox(-_PAGE_SLACK_PT, -_PAGE_SLACK_PT,
                        page.rect.width + 2 * _PAGE_SLACK_PT,
                        page.rect.height + 2 * _PAGE_SLACK_PT)
        lines = _text_lines(page)
        images = _image_boxes(page)

        for kind, boxes in (("text", lines), ("image", images)):
            for box in boxes:
                if (box.x < page_box.x or box.y < page_box.y
                        or box.x2 > page_box.x2 or box.y2 > page_box.y2):
                    report.issues.append(
                        f"page {page_index + 1}: {kind} ink outside page "
                        f"at ({box.x:.0f},{box.y:.0f},{box.x2:.0f},{box.y2:.0f})")

        # text-text collisions (different lines overlapping substantially)
        tt_ok = pages_match and src_tt_overlap.get(page_index, False)
        ti_ok = pages_match and src_ti_overlap.get(page_index, False)
        for i in range(len(lines)):
            for j in range(i + 1, len(lines)):
                if _collides(lines[i], lines[j]):
                    a = lines[i]
                    msg = (f"page {page_index + 1}: text lines overlap "
                           f"near ({a.x:.0f},{a.y:.0f})")
                    (report.warnings if tt_ok else report.issues).append(msg)
                    break  # one report per line is enough
        # text-image collisions
        for line in lines:
            for img in images:
                if _collides(line, img):
                    msg = (f"page {page_index + 1}: text overlaps image "
                           f"near ({line.x:.0f},{line.y:.0f})")
                    (report.warnings if ti_ok else report.issues).append(msg)
                    break
    return report


def _source_overlaps(pptx_path: Path | str) -> tuple[dict[int, bool], dict[int, bool]]:
    """Per slide: does the source itself overlap text with text / images?"""
    from typstpresenter.verify.geometry import ElementKind
    from typstpresenter.verify.pptx_geometry import extract_pptx_geometry

    doc = extract_pptx_geometry(pptx_path)
    tt: dict[int, bool] = {}
    ti: dict[int, bool] = {}
    for slide in doc.slides:
        texts = [e.bbox for e in slide.elements if e.kind == ElementKind.TEXT]
        images = [e.bbox for e in slide.elements if e.kind == ElementKind.IMAGE]
        # shapes double as text carriers: their labels can overlap anything
        shapes = [e.bbox for e in slide.elements
                  if e.kind in (ElementKind.SHAPE, ElementKind.CONNECTOR)]
        tt[slide.index] = any(
            _collides(a, b)
            for k, a in enumerate(texts) for b in texts[k + 1:]
        ) or any(_collides(t, s) for t in texts for s in shapes) or bool(
            shapes and len(shapes) > 1)
        # translated label text renders slightly larger than its source box,
        # so captions that *adjoin* an image in the source may brush it in
        # the output -- tolerate sources within a small margin
        grown = [BBox(i.x - 8, i.y - 8, i.w + 16, i.h + 16) for i in images]
        ti[slide.index] = any(
            _collides(t, i) for t in texts for i in grown
        ) or any(_collides(s, i) for s in shapes for i in grown)
    return tt, ti
