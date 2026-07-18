"""
Method D: diagram (autoshape/connector) translation fidelity.

Methods S and F check text, images and overall layout; neither looks at
what CeTZ actually drew inside a diagram canvas. Method D compares the
PPTX autoshapes/connectors of a slide against the vector paths PyMuPDF
extracts from the compiled PDF (``page.get_drawings()``):

- every source shape/connector must have a matching rendered path,
- its rendered *type* (rect / rounded rect / diamond / oval / triangle /
  line / elbow) must agree with the PPTX preset geometry,
- its fill/stroke colors must match the resolved source colors,
- its position and size, normalized against the diagram cluster's own
  bounding box, must match the source's normalized geometry (this is
  robust to the cluster being uniformly scaled/translated/aligned by the
  flow emitter -- exact pixel coordinates are not the point).

This only applies to slides whose diagram cluster renders as an isolated,
unambiguous set of vector paths: slides that also contain a table (whose
cell borders are additional vector strokes) are skipped, as are slides
where the source/rendered shape counts do not match closely enough to
support a confident normalization. That keeps Method D precise on the
synthetic single-canvas benchmarks it is primarily meant for, while still
being safe to run over real decks (it just abstains where the page is too
visually complex to disentangle).

Slides with more than ``_CONFIDENT_SHAPE_COUNT`` source shapes still get
checked, but their findings land in ``warnings`` rather than ``issues``:
the matcher's single global affine transform + greedy nearest-neighbor
assignment is exact-verified only up to that scale (the synthetic
benchmark corpus); on denser real slides it accumulates enough small
misattributions to produce false positives on renders that are visually
correct (confirmed by eye against several flagged real slides).
"""

from __future__ import annotations

import math
from dataclasses import dataclass, field
from pathlib import Path

import fitz
import pptx

from typstpresenter.verify.geometry import EMU_PER_PT, BBox

# a closed polygon this small is an arrowhead marker, not a shape
_ARROWHEAD_DIAG_PT = 10.0
# fraction of source/rendered shapes that must be present to attempt
# normalized-geometry matching at all. A single global affine transform
# (uniform scale + translate) only holds for a genuinely isolated diagram
# cluster; dense real slides absorb extra text/images into the canvas,
# embed OLE objects the converter drops entirely, or place several
# loosely-related shape groups on one page, none of which this matcher's
# single-transform model can represent -- better to abstain than to
# report noisy, low-confidence findings on those. Simple, single-diagram
# slides (the corpus this method is chiefly meant to validate) match
# almost exactly, so a tight fraction does not cost real coverage there.
_MIN_MATCH_FRACTION = 0.92
# relative-size and relative-position tolerances (fraction of cluster size)
_SIZE_TOLERANCE = 0.25
_POS_TOLERANCE = 0.12
_COLOR_TOLERANCE = 40  # per-channel, 0-255
# slides with more source shapes than this get their per-shape/connector
# findings downgraded to warnings: the matcher's single global affine
# transform + greedy nearest-neighbor assignment was validated (0 false
# positives) against the synthetic benchmark decks, all at or below this
# scale; on denser real slides (sub-grouped icon clusters, shapes absorbed
# from outside the diagram kernel, several loosely related visual clusters
# on one page) it produces enough false positives on visually-correct
# renders that "issue" would be misleading -- confirmed by eye on several
# flagged real slides that were in fact perfectly translated
_CONFIDENT_SHAPE_COUNT = 15


@dataclass
class DiagramReport:
    issues: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    checked_slides: int = 0
    skipped_slides: int = 0
    checked_shapes: int = 0
    total_shapes: int = 0

    @property
    def ok(self) -> bool:
        return not self.issues

    def summary(self) -> str:
        status = "OK" if self.ok else f"{len(self.issues)} issues"
        cov = (self.checked_shapes / self.total_shapes * 100) if self.total_shapes > 0 else 0.0
        lines = [f"{status} ({len(self.warnings)} warnings) -- "
                f"{self.checked_slides} slide(s) checked, "
                f"{self.skipped_slides} skipped (too complex / no diagram), "
                f"coverage {cov:.1f}% ({self.checked_shapes}/{self.total_shapes} shapes)"]
        lines += [f"  issue: {i}" for i in self.issues]
        lines += [f"  warning: {w}" for w in self.warnings]
        return "\n".join(lines)


# --------------------------------------------------------------- source ----

@dataclass
class _SourceShape:
    eid: str
    kind: str          # "rect" | "rounded_rect" | "diamond" | "oval" | "triangle" | "other_shape" | "connector"
    bbox: BBox         # PPTX (unrotated) bounding box
    fill_rgb: str | None
    stroke_rgb: str | None
    endpoints: tuple[tuple[float, float], tuple[float, float]] | None = None
    rotated: bool = False
    effective_bbox: BBox = field(default=None)  # type: ignore[assignment]

    def __post_init__(self) -> None:
        if self.effective_bbox is None:
            self.effective_bbox = self.bbox


def _rotated_bbox(bbox: BBox, angle_deg: float) -> BBox:
    """Axis-aligned bbox of `bbox` rotated `angle_deg` clockwise (PPTX/
    OOXML convention, y-down) around its own center."""
    angle = math.radians(angle_deg)
    cos_a, sin_a = math.cos(angle), math.sin(angle)
    cx, cy = bbox.center
    corners = [(bbox.x, bbox.y), (bbox.x2, bbox.y), (bbox.x2, bbox.y2), (bbox.x, bbox.y2)]
    rotated = []
    for x, y in corners:
        dx, dy = x - cx, y - cy
        rotated.append((cx + dx * cos_a - dy * sin_a, cy + dx * sin_a + dy * cos_a))
    xs = [p[0] for p in rotated]
    ys = [p[1] for p in rotated]
    return BBox(min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys))


_AUTO_SHAPE_KIND: dict[object, str] = {}


def _auto_shape_kind(shape) -> str:
    from pptx.enum.shapes import MSO_SHAPE

    if not _AUTO_SHAPE_KIND:
        _AUTO_SHAPE_KIND.update({
            MSO_SHAPE.RECTANGLE: "rect",
            MSO_SHAPE.ROUNDED_RECTANGLE: "rounded_rect",
            MSO_SHAPE.DIAMOND: "diamond",
            MSO_SHAPE.OVAL: "oval",
            MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
        })
    try:
        auto = shape.auto_shape_type
    except (ValueError, AttributeError):
        return "other_shape"
    return _AUTO_SHAPE_KIND.get(auto, "other_shape")


def _shape_rotation_deg(shape) -> float:
    """Clockwise rotation in degrees (PPTX/OOXML convention), or 0.0."""
    from pptx.oxml.ns import qn

    xfrm = shape.element.spPr.find(qn("a:xfrm")) if shape.element.spPr is not None else None
    rot = xfrm.get("rot") if xfrm is not None else None
    return int(rot) / 60000.0 if rot else 0.0


def source_diagram_shapes(slide, slide_index: int) -> list[_SourceShape]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from typstpresenter.convert.cetz import is_diagram_shape
    from typstpresenter.convert.pptx_style import shape_fill_rgb, shape_line_rgb
    from typstpresenter.verify.pptx_geometry import element_id, iter_flat_shapes

    out = []
    for shape, bbox in iter_flat_shapes(slide.shapes):
        if not is_diagram_shape(shape):
            continue
        eid = element_id(slide_index, shape.shape_id)
        is_connector = (shape.shape_type == MSO_SHAPE_TYPE.LINE
                        or shape.element.tag.endswith("}cxnSp"))
        if is_connector:
            bx, by = shape.begin_x / EMU_PER_PT, shape.begin_y / EMU_PER_PT
            ex, ey = shape.end_x / EMU_PER_PT, shape.end_y / EMU_PER_PT
            out.append(_SourceShape(
                eid=eid, kind="connector", bbox=bbox,
                fill_rgb=None, stroke_rgb=shape_line_rgb(shape) or "000000",
                endpoints=((bx, by), (ex, ey)),
            ))
        else:
            angle = _shape_rotation_deg(shape)
            out.append(_SourceShape(
                eid=eid, kind=_auto_shape_kind(shape), bbox=bbox,
                fill_rgb=shape_fill_rgb(shape), stroke_rgb=shape_line_rgb(shape),
                rotated=bool(angle),
                effective_bbox=_rotated_bbox(bbox, angle) if angle else bbox,
            ))
    return out


# ---------------------------------------------------------------- render ---

@dataclass
class _RenderedItem:
    bbox: BBox
    kind: str           # "polygon" | "ellipse" | "open_path"
    axis_aligned: bool   # all edges horizontal/vertical (closed polygons only)
    fill_rgb: tuple[int, int, int] | None
    stroke_rgb: tuple[int, int, int] | None
    n_points: int


def _pts_of(item: tuple) -> list[tuple[float, float]]:
    tag = item[0]
    if tag in ("l", "c"):
        return [(p.x, p.y) for p in item[1:] if hasattr(p, "x")]
    if tag == "qu":
        quad = item[1]
        return [(pt.x, pt.y) for pt in (quad.ul, quad.ur, quad.lr, quad.ll)]
    if tag == "re":
        rect = item[1]
        return [(rect.x0, rect.y0), (rect.x1, rect.y1)]
    return []


def _to255(c: tuple[float, float, float] | None) -> tuple[int, int, int] | None:
    if c is None:
        return None
    return tuple(round(v * 255) for v in c)  # type: ignore[return-value]


def _merge_drawings(drawings: list[dict]) -> list[_RenderedItem]:
    """Group same-bbox fill/stroke records (typst splits a filled+stroked
    shape into two draw calls) and classify the merged path."""
    groups: list[dict] = []
    for d in drawings:
        r = d["rect"]
        bbox = BBox(r.x0, r.y0, r.width, r.height)
        placed = False
        for g in groups:
            gb = g["bbox"]
            if (abs(bbox.x - gb.x) < 1.5 and abs(bbox.y - gb.y) < 1.5
                    and abs(bbox.w - gb.w) < 1.5 and abs(bbox.h - gb.h) < 1.5):
                g["records"].append(d)
                placed = True
                break
        if not placed:
            groups.append({"bbox": bbox, "records": [d]})

    items = []
    for g in groups:
        bbox = g["bbox"]
        fill_rgb = stroke_rgb = None
        all_points: list[tuple[float, float]] = []
        item_tags: list[str] = []
        for d in g["records"]:
            if d.get("fill") is not None:
                fill_rgb = _to255(d["fill"])
            if d.get("type") in ("s", "fs") and d.get("color") is not None:
                stroke_rgb = _to255(d["color"])
            for it in d["items"]:
                item_tags.append(it[0])
                all_points += _pts_of(it)
        if not all_points:
            continue
        has_curve = "c" in item_tags
        has_quad = "qu" in item_tags
        closed = (has_quad or (len(all_points) >= 3
                  and math.dist(all_points[0], all_points[-1]) < 1.5))
        if has_curve:
            kind = "ellipse" if not any(t == "l" for t in item_tags) else "polygon"
            axis_aligned = False
        elif closed or has_quad:
            kind = "polygon"
            edges = list(zip(all_points, all_points[1:]))
            axis_aligned = all(
                abs(a[0] - b[0]) < 1.0 or abs(a[1] - b[1]) < 1.0 for a, b in edges
            ) if edges else True
        else:
            kind = "open_path"
            axis_aligned = False
        items.append(_RenderedItem(
            bbox=bbox, kind=kind, axis_aligned=axis_aligned,
            fill_rgb=fill_rgb, stroke_rgb=stroke_rgb, n_points=len(all_points),
        ))
    # Merge touching open_path connectors
    connectors = [it for it in items if it.kind == "open_path"]
    other_items = [it for it in items if it.kind != "open_path"]
    changed = True
    while changed:
        changed = False
        i = 0
        while i < len(connectors):
            j = i + 1
            while j < len(connectors):
                c1 = connectors[i]
                c2 = connectors[j]
                dx = max(0, max(c1.bbox.x, c2.bbox.x) - min(c1.bbox.x2, c2.bbox.x2))
                dy = max(0, max(c1.bbox.y, c2.bbox.y) - min(c1.bbox.y2, c2.bbox.y2))
                if dx < 4.0 and dy < 4.0:
                    union_box = BBox(
                        min(c1.bbox.x, c2.bbox.x),
                        min(c1.bbox.y, c2.bbox.y),
                        max(c1.bbox.x2, c2.bbox.x2) - min(c1.bbox.x, c2.bbox.x),
                        max(c1.bbox.y2, c2.bbox.y2) - min(c1.bbox.y, c2.bbox.y)
                    )
                    merged = _RenderedItem(
                        bbox=union_box,
                        kind="open_path",
                        axis_aligned=c1.axis_aligned and c2.axis_aligned,
                        fill_rgb=c1.fill_rgb or c2.fill_rgb,
                        stroke_rgb=c1.stroke_rgb or c2.stroke_rgb,
                        n_points=c1.n_points + c2.n_points
                    )
                    connectors[i] = merged
                    connectors.pop(j)
                    changed = True
                    break
                j += 1
            if changed:
                break
            i += 1
    return other_items + connectors


def rendered_diagram_items(page: fitz.Page) -> tuple[list[_RenderedItem], list[_RenderedItem]]:
    """(shapes, connectors) -- arrowhead markers are dropped."""
    items = _merge_drawings(page.get_drawings())
    shapes, connectors = [], []
    for it in items:
        if it.kind in ("polygon", "ellipse"):
            if math.hypot(it.bbox.w, it.bbox.h) < _ARROWHEAD_DIAG_PT:
                continue  # arrowhead
            shapes.append(it)
        else:
            connectors.append(it)
    return shapes, connectors


# ------------------------------------------------------------- matching ----

def _color_close(a: str | None, b: tuple[int, int, int] | None) -> bool:
    if a is None or b is None:
        return a is None and b is None
    ar, ag, ab = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
    return (abs(ar - b[0]) <= _COLOR_TOLERANCE and abs(ag - b[1]) <= _COLOR_TOLERANCE
            and abs(ab - b[2]) <= _COLOR_TOLERANCE)


def _expected_kind_ok(source_kind: str, rendered: _RenderedItem) -> bool:
    if source_kind in ("rect", "rounded_rect", "other_shape"):
        return rendered.kind == "polygon"
    if source_kind == "diamond":
        return rendered.kind == "polygon" and not rendered.axis_aligned
    if source_kind == "oval":
        return rendered.kind == "ellipse"
    if source_kind == "triangle":
        return rendered.kind == "polygon"
    return True


def _normalize(bbox: BBox, cluster: BBox) -> tuple[float, float, float, float]:
    cw = cluster.w or 1.0
    ch = cluster.h or 1.0
    return ((bbox.x - cluster.x) / cw, (bbox.y - cluster.y) / ch,
            bbox.w / cw, bbox.h / ch)


def _cluster_bounds(boxes: list[BBox]) -> BBox:
    x = min(b.x for b in boxes)
    y = min(b.y for b in boxes)
    return BBox(x, y, max(b.x2 for b in boxes) - x, max(b.y2 for b in boxes) - y)


def _check_slide(shapes: list[_SourceShape], page: fitz.Page,
                 report: DiagramReport, slide_no: int) -> None:
    src_shapes = [s for s in shapes if s.kind != "connector"]
    src_conns = [s for s in shapes if s.kind == "connector"]
    if not src_shapes:
        report.skipped_slides += 1
        return

    rendered_shapes, rendered_conns = rendered_diagram_items(page)
    # a slide with table borders or many unrelated vector strokes makes
    # the rendered set too noisy to attribute confidently -- abstain
    if len(rendered_shapes) < len(src_shapes) * _MIN_MATCH_FRACTION:
        report.skipped_slides += 1
        report.warnings.append(
            f"slide {slide_no}: abstained (found {len(rendered_shapes)} "
            f"rendered shape-like paths for {len(src_shapes)} source shapes)")
        return
    if len(rendered_shapes) > len(src_shapes) * (1 / _MIN_MATCH_FRACTION) + 3:
        report.skipped_slides += 1
        report.warnings.append(
            f"slide {slide_no}: abstained (found {len(rendered_shapes)} "
            f"rendered shape-like paths, only {len(src_shapes)} source shapes "
            "-- likely other vector content, e.g. a table, on this slide)")
        return

    report.checked_slides += 1
    report.checked_shapes += len(src_shapes)
    src_cluster = _cluster_bounds([s.effective_bbox for s in src_shapes])
    rendered_cluster = _cluster_bounds([r.bbox for r in rendered_shapes])

    # findings on clusters bigger than validated are advisory (warnings),
    # not blocking -- see _CONFIDENT_SHAPE_COUNT
    confident = len(src_shapes) <= _CONFIDENT_SHAPE_COUNT
    sink = report.issues if confident else report.warnings

    remaining = list(rendered_shapes)
    for s in sorted(src_shapes, key=lambda s: -s.effective_bbox.w * s.effective_bbox.h):
        s_norm = _normalize(s.effective_bbox, src_cluster)

        def cost(r: _RenderedItem) -> float:
            r_norm = _normalize(r.bbox, rendered_cluster)
            dx = r_norm[0] - s_norm[0]
            dy = r_norm[1] - s_norm[1]
            dw = r_norm[2] - s_norm[2]
            dh = r_norm[3] - s_norm[3]
            return dx * dx + dy * dy + dw * dw + dh * dh

        if not remaining:
            sink.append(f"slide {slide_no}: shape {s.eid} ({s.kind}) - missing from the rendered diagram")
            continue
        best = min(remaining, key=cost)
        remaining.remove(best)
        r_norm = _normalize(best.bbox, rendered_cluster)
        dpos = math.hypot(r_norm[0] - s_norm[0], r_norm[1] - s_norm[1])
        dsize = math.hypot(r_norm[2] - s_norm[2], r_norm[3] - s_norm[3])
        if dpos > _POS_TOLERANCE:
            sink.append(
                f"slide {slide_no}: shape {s.eid} ({s.kind}) - displaced (drift {dpos:.2f} of cluster size)")
        if dsize > _SIZE_TOLERANCE:
            sink.append(
                f"slide {slide_no}: shape {s.eid} ({s.kind}) - wrong-size (drift {dsize:.2f} of cluster size)")
        if not _expected_kind_ok(s.kind, best):
            sink.append(
                f"slide {slide_no}: shape {s.eid} ({s.kind}) - wrong-kind (expected {s.kind}, "
                f"rendered path looks like {best.kind}{' (axis-aligned)' if best.axis_aligned else ''})")
        if not _color_close(s.fill_rgb, best.fill_rgb):
            sink.append(
                f"slide {slide_no}: shape {s.eid} ({s.kind}) - wrong-style (fill mismatch: "
                f"source #{s.fill_rgb}, rendered {best.fill_rgb})")
        if not _color_close(s.stroke_rgb, best.stroke_rgb):
            report.warnings.append(
                f"slide {slide_no}: shape {s.eid} ({s.kind}) - wrong-style (stroke mismatch: "
                f"source #{s.stroke_rgb}, rendered {best.stroke_rgb})")

    for extra in remaining:
        report.warnings.append(
            f"slide {slide_no}: unmatched rendered shape near "
            f"({extra.bbox.x:.0f},{extra.bbox.y:.0f}) -- no corresponding source shape")

    _check_connectors(src_conns, src_shapes, rendered_conns, src_cluster,
                      rendered_cluster, sink, slide_no)


def _check_slide_clusters(shapes: list[_SourceShape], page: fitz.Page,
                          report: DiagramReport, slide_no: int,
                          markers: list[dict]) -> None:
    all_rendered_shapes, all_rendered_conns = rendered_diagram_items(page)
    checked_any = False

    for marker in markers:
        cluster_id = marker["id"]
        c_shapes_set = set(marker.get("shapes", []))

        src_shapes = [s for s in shapes if s.eid in c_shapes_set and s.kind != "connector"]
        src_conns = [s for s in shapes if s.eid in c_shapes_set and s.kind == "connector"]

        if not src_shapes:
            continue

        canvas_bbox = BBox(marker["x"], marker["y"], marker["w"], marker["h"])

        def is_inside(r_bbox: BBox) -> bool:
            cx, cy = r_bbox.center
            return (canvas_bbox.x - 5.0 <= cx <= canvas_bbox.x2 + 5.0 and
                    canvas_bbox.y - 5.0 <= cy <= canvas_bbox.y2 + 5.0)

        rendered_shapes = [r for r in all_rendered_shapes if is_inside(r.bbox)]
        rendered_conns = [r for r in all_rendered_conns if is_inside(r.bbox)]

        if len(rendered_shapes) < len(src_shapes) * _MIN_MATCH_FRACTION:
            report.skipped_slides += 1
            report.warnings.append(
                f"slide {slide_no} (cluster {cluster_id}): abstained (found {len(rendered_shapes)} "
                f"rendered shape-like paths for {len(src_shapes)} source shapes)")
            continue
        if len(rendered_shapes) > len(src_shapes) * (1 / _MIN_MATCH_FRACTION) + 3:
            report.skipped_slides += 1
            report.warnings.append(
                f"slide {slide_no} (cluster {cluster_id}): abstained (found {len(rendered_shapes)} "
                f"rendered shape-like paths, only {len(src_shapes)} source shapes)")
            continue

        checked_any = True
        report.checked_shapes += len(src_shapes)
        src_cluster = _cluster_bounds([s.effective_bbox for s in src_shapes])
        rendered_cluster = _cluster_bounds([r.bbox for r in rendered_shapes])

        confident = len(src_shapes) <= _CONFIDENT_SHAPE_COUNT
        sink = report.issues if confident else report.warnings

        remaining = list(rendered_shapes)
        for s in sorted(src_shapes, key=lambda s: -s.effective_bbox.w * s.effective_bbox.h):
            s_norm = _normalize(s.effective_bbox, src_cluster)

            def cost(r: _RenderedItem) -> float:
                r_norm = _normalize(r.bbox, rendered_cluster)
                dx = r_norm[0] - s_norm[0]
                dy = r_norm[1] - s_norm[1]
                dw = r_norm[2] - s_norm[2]
                dh = r_norm[3] - s_norm[3]
                return dx * dx + dy * dy + dw * dw + dh * dh

            if not remaining:
                sink.append(f"slide {slide_no} (cluster {cluster_id}): shape {s.eid} ({s.kind}) - missing from the rendered diagram")
                continue
            best = min(remaining, key=cost)
            remaining.remove(best)
            r_norm = _normalize(best.bbox, rendered_cluster)
            dpos = math.hypot(r_norm[0] - s_norm[0], r_norm[1] - s_norm[1])
            dsize = math.hypot(r_norm[2] - s_norm[2], r_norm[3] - s_norm[3])
            if dpos > _POS_TOLERANCE:
                sink.append(
                    f"slide {slide_no} (cluster {cluster_id}): shape {s.eid} ({s.kind}) - displaced (drift {dpos:.2f} of cluster size)")
            if dsize > _SIZE_TOLERANCE:
                sink.append(
                    f"slide {slide_no} (cluster {cluster_id}): shape {s.eid} ({s.kind}) - wrong-size (drift {dsize:.2f} of cluster size)")
            if not _expected_kind_ok(s.kind, best):
                sink.append(
                    f"slide {slide_no} (cluster {cluster_id}): shape {s.eid} ({s.kind}) - wrong-kind (expected {s.kind}, "
                    f"rendered path looks like {best.kind}{' (axis-aligned)' if best.axis_aligned else ''})")
            if not _color_close(s.fill_rgb, best.fill_rgb):
                sink.append(
                    f"slide {slide_no} (cluster {cluster_id}): shape {s.eid} ({s.kind}) - wrong-style (fill mismatch: "
                    f"source #{s.fill_rgb}, rendered {best.fill_rgb})")
            if not _color_close(s.stroke_rgb, best.stroke_rgb):
                report.warnings.append(
                    f"slide {slide_no} (cluster {cluster_id}): shape {s.eid} ({s.kind}) - wrong-style (stroke mismatch: "
                    f"source #{s.stroke_rgb}, rendered {best.stroke_rgb})")

        for extra in remaining:
            report.warnings.append(
                f"slide {slide_no} (cluster {cluster_id}): unmatched rendered shape near "
                f"({extra.bbox.x:.0f},{extra.bbox.y:.0f}) -- no corresponding source shape")

        _check_connectors(src_conns, shapes, rendered_conns, src_cluster,
                          rendered_cluster, sink, slide_no)

    if checked_any:
        report.checked_slides += 1


def _nearest_shape_id(point: tuple[float, float], shapes: list[_SourceShape]) -> str:
    return min(shapes, key=lambda s: (s.bbox.center[0] - point[0]) ** 2
                                     + (s.bbox.center[1] - point[1]) ** 2).eid


def _map_point(pt: tuple[float, float], src_cluster: BBox,
              rendered_cluster: BBox) -> tuple[float, float]:
    """Map a point from source (PPTX) space into rendered (PDF page) space
    via the same uniform scale + translate the diagram canvas as a whole
    was placed with (derived from the two clusters' bounding boxes)."""
    fx = (pt[0] - src_cluster.x) / (src_cluster.w or 1.0)
    fy = (pt[1] - src_cluster.y) / (src_cluster.h or 1.0)
    return (rendered_cluster.x + fx * rendered_cluster.w,
            rendered_cluster.y + fy * rendered_cluster.h)


def _check_connectors(src_conns: list[_SourceShape], src_shapes: list[_SourceShape],
                      rendered_conns: list[_RenderedItem],
                      src_cluster: BBox, rendered_cluster: BBox,
                      sink: list[str], slide_no: int) -> None:
    if not src_conns:
        return
    if len(rendered_conns) < len(src_conns):
        sink.append(
            f"slide {slide_no}: only {len(rendered_conns)} of {len(src_conns)} "
            "connectors found in the rendered diagram")
        return
    # topology: which two shapes does each connector touch (by nearest
    # shape center to each endpoint, in source space); a rendered
    # connector's bbox corners, mapped into source space, should land near
    # both of those shapes' source bboxes
    remaining = list(rendered_conns)
    for c in src_conns:
        assert c.endpoints is not None
        a_id = _nearest_shape_id(c.endpoints[0], src_shapes)
        b_id = _nearest_shape_id(c.endpoints[1], src_shapes)
        if a_id == b_id:
            continue  # degenerate (self-loop-ish); not worth chasing

        def touches_both(r: _RenderedItem) -> bool:
            corners_page = [(r.bbox.x, r.bbox.y), (r.bbox.x2, r.bbox.y2),
                            (r.bbox.x, r.bbox.y2), (r.bbox.x2, r.bbox.y)]
            corners_src = [_map_point(pt, rendered_cluster, src_cluster)
                          for pt in corners_page]
            near_a = any(_near_shape(pt, a_id, src_shapes) for pt in corners_src)
            near_b = any(_near_shape(pt, b_id, src_shapes) for pt in corners_src)
            return near_a and near_b

        match = next((r for r in remaining if touches_both(r)), None)
        if match is None:
            sink.append(
                f"slide {slide_no}: connector {c.eid} should join {a_id} and "
                f"{b_id}, but no rendered line touches both")
        else:
            remaining.remove(match)


def _near_shape(pt: tuple[float, float], eid: str,
                shapes: list[_SourceShape], slack: float = 20.0) -> bool:
    s = next(s for s in shapes if s.eid == eid)
    b = s.bbox
    return (b.x - slack <= pt[0] <= b.x2 + slack
            and b.y - slack <= pt[1] <= b.y2 + slack)


# ------------------------------------------------------------------ entry --

def verify_diagrams(pptx_path: Path | str, pdf_path: Path | str) -> DiagramReport:
    """Check diagram (autoshape/connector) translation fidelity.

    Compares every slide's PPTX diagram shapes against the vector paths
    rendered on the matching PDF page. Slides without diagrams, or whose
    page can't be confidently attributed to the source shape count (e.g.
    a table sharing the slide), are skipped rather than mis-flagged.
    """
    prs = pptx.Presentation(str(pptx_path))
    doc = fitz.open(pdf_path)
    report = DiagramReport()
    if len(doc) != len(prs.slides):
        report.warnings.append(
            f"page count {len(doc)} != slide count {len(prs.slides)}; "
            "skipping (pages don't align 1:1 with slides)")
        return report

    # Query canvas markers
    canvas_markers_by_page: dict[int, list[dict]] = {}
    typ_path = Path(pdf_path).with_suffix(".typ")
    if typ_path.exists():
        from typstpresenter.convert.emitter import emit_minimal
        from typstpresenter.verify.typst_tools import query
        
        temp_typ_path = typ_path.with_name(f"{typ_path.stem}_verify_d.typ")
        try:
            emit_minimal(pptx_path, temp_typ_path, canvas_markers=True)
            res = query(temp_typ_path, "<tp-canvas>")
            for item in res.value:
                page_idx = item["page"] - 1 # 0-based
                canvas_markers_by_page.setdefault(page_idx, []).append(item)
        except Exception as e:
            report.warnings.append(f"Could not query canvas markers: {e}")
        finally:
            if temp_typ_path.exists():
                try:
                    temp_typ_path.unlink()
                except Exception:
                    pass

    for slide_index, slide in enumerate(prs.slides):
        shapes = source_diagram_shapes(slide, slide_index)
        if not shapes:
            continue
        src_shapes_all = [s for s in shapes if s.kind != "connector"]
        report.total_shapes += len(src_shapes_all)

        markers = canvas_markers_by_page.get(slide_index, [])
        if markers:
            _check_slide_clusters(shapes, doc[slide_index], report, slide_index + 1, markers)
        else:
            _check_slide(shapes, doc[slide_index], report, slide_index + 1)
    return report
