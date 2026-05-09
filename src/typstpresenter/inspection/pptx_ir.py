"""Convert PPTX slides to the inspection IR, using the existing model layer."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

import pptx as python_pptx

from typstpresenter.model.Element import PlacedElement
from typstpresenter.model.Image import Image
from typstpresenter.model.List import List
from typstpresenter.model.Presentation import Presentation
from typstpresenter.model.PresentationTitle import PresentationTitle
from typstpresenter.model.Slide import Slide
from typstpresenter.model.Title import Title
from typstpresenter.model.text.Text import Text

from .ir import BoundingBox, ElementIR, PresentationIR, SlideIR

# Standard PowerPoint slide dimensions in EMU (914400 EMU = 1 inch)
_DEFAULT_SLIDE_WIDTH_EMU = 9144000   # 10 inches
_DEFAULT_SLIDE_HEIGHT_EMU = 6858000  # 7.5 inches


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _rel(value: int | float | None, dimension: int) -> float:
    """Convert an EMU value to a relative slide fraction."""
    if value is None:
        return 0.0
    return float(value) / float(dimension)


def _bounds(placed: PlacedElement, w: int, h: int) -> BoundingBox | None:
    if placed.left is None or placed.top is None:
        return None
    return BoundingBox(
        left=_rel(placed.left, w),
        top=_rel(placed.top, h),
        width=_rel(placed.width or 0, w),
        height=_rel(placed.height or 0, h),
    )


def _list_to_flat_items(lst: List, indent: int = 0) -> list[str]:
    """Recursively flatten a nested List to indented strings."""
    result: list[str] = []
    prefix = "  " * indent
    for item in lst.items:
        if isinstance(item, List):
            result.extend(_list_to_flat_items(item, indent + 1))
        else:
            result.append(f"{prefix}{item}")
    return result


def _placed_to_element_ir(placed: PlacedElement, slide_w: int, slide_h: int) -> ElementIR | None:
    element = placed.element
    box = _bounds(placed, slide_w, slide_h)

    if isinstance(element, Title):
        return ElementIR(kind="title", bounds=box, text=str(element.text))
    if isinstance(element, PresentationTitle):
        return ElementIR(kind="presentation_title", bounds=box, text=str(element.text))
    if isinstance(element, Text):
        return ElementIR(kind="text", bounds=box, text=str(element))
    if isinstance(element, List):
        return ElementIR(kind="list", bounds=box, items=_list_to_flat_items(element))
    if isinstance(element, Image):
        return ElementIR(kind="image", bounds=box, image_name=element.name)
    return None


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def slide_to_ir(
    slide: Slide,
    slide_index: int,
    slide_width: int = _DEFAULT_SLIDE_WIDTH_EMU,
    slide_height: int = _DEFAULT_SLIDE_HEIGHT_EMU,
) -> SlideIR:
    """Convert an existing :class:`Slide` model object to a :class:`SlideIR`.

    Parameters
    ----------
    slide:
        The parsed slide (from ``Slide.from_pptx_slide``).
    slide_index:
        0-based slide index within the presentation (title slide excluded).
    slide_width / slide_height:
        Slide dimensions in EMU; defaults to the standard 10 × 7.5 inch canvas.
    """
    elements: list[ElementIR] = []
    title_text: str | None = None

    for placed in slide.elements:
        ir = _placed_to_element_ir(placed, slide_width, slide_height)
        if ir is None:
            continue
        if ir.kind == "title" and title_text is None:
            title_text = ir.text
        elements.append(ir)

    return SlideIR(index=slide_index, title=title_text, elements=elements)


def pptx_to_presentation_ir(path: Path) -> PresentationIR:
    """Load a ``.pptx`` file and return a :class:`PresentationIR`.

    The first slide is treated as the title slide (matching the existing
    :class:`~typstpresenter.model.Presentation` convention) and is only
    used to populate ``PresentationIR.title``; it is not included in
    ``PresentationIR.slides``.
    """
    raw_prs = python_pptx.Presentation(str(path))
    slide_w = int(raw_prs.slide_width)
    slide_h = int(raw_prs.slide_height)

    prs = Presentation.from_file(path)
    title_str = str(prs.title) if prs.title else None

    slide_irs = [
        slide_to_ir(slide, i, slide_w, slide_h)
        for i, slide in enumerate(prs.slides)
    ]

    return PresentationIR(title=title_str, slides=slide_irs)
