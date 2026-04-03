from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.shapes import Subshape
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.text.text import _Paragraph, TextFrame, _Run

from typstpresenter.model.Element import Element
from typstpresenter.model.List import List
from typstpresenter.model.PresentationTitle import PresentationTitle
from typstpresenter.model.Title import Title
from typstpresenter.model.text.Link import Link
from typstpresenter.model.text.Subscript import Subscript
from typstpresenter.model.text.Superscript import Superscript
from typstpresenter.model.text.Text import Text, Atom
from typstpresenter.powerpoint.Ignore import Ignore


def _interpret_placeholder(shape: SlidePlaceholder) -> Element | Ignore | None:
    """
    Interpret a shape of type SlidePlaceholder.

    These are somewhat nicer to interpret, because some of them have a type assigned which is actually semantic.

    See: https://python-pptx.readthedocs.io/en/latest/user/understanding-shapes.html
    """

    match shape.placeholder_format.type:
        case PP_PLACEHOLDER_TYPE.TITLE:
            return Title(text=_interpret_text_frame(shape.text_frame))
        case PP_PLACEHOLDER_TYPE.CENTER_TITLE:
            return PresentationTitle(text=_interpret_text_frame(shape.text_frame))
        case PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
            return Ignore()
        case PP_PLACEHOLDER_TYPE.OBJECT:
            # Just pretend that object means a bunch of text, and nothing else.
            return _interpret_text_frame(shape.text_frame)
        case _:
            return None


def _interpret_text_frame(text_frame: TextFrame) -> Text | List:
    if len(text_frame.paragraphs) == 1:
        return _interpret_paragraph(text_frame.paragraphs[0])
    else:
        # Just pretend any multi-paragraph text is a list
        return List(items=tuple(_interpret_paragraph(p) for p in text_frame.paragraphs))


def _interpret_paragraph(paragraph: _Paragraph) -> Text:
    return Text(tuple(_interpret_run(run) for run in paragraph.runs))


# See https://stackoverflow.com/questions/61329224/how-do-i-add-superscript-subscript-text-to-powerpoint-using-python-pptx
_SUBSCRIPT = "-25000"
_SUPERSCRIPT = "30000"


def _interpret_run(run: _Run) -> Atom:
    if run.hyperlink.address is not None:
        return Link(text=Text(run.text), target=run.hyperlink.address)

    # The text baseline indicates how high in a line the run is positioned.
    baseline_position = run.font._element.get("baseline")

    if baseline_position == _SUBSCRIPT:
        return Subscript(text=Text(run.text))
    elif baseline_position == _SUPERSCRIPT:
        return Superscript(text=Text(run.text))

    return run.text


def interpret(shape: BaseShape | Subshape) -> Element | None:
    """
    Interpret a shape (or a subshape, which is an element introduced as child of a shape) coming from a PowerPoint.

    Interpreting a shape means converting it to an Element, i.e., an instance from our abstraction layer which
    we control. If a shape cannot be interpreted, this returns None.
    """
    match shape:
        case SlidePlaceholder():
            return _interpret_placeholder(shape)
        case _:
            # TODO Implement as needed
            return None
