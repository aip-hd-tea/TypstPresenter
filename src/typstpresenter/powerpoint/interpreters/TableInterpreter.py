from pptx.shapes import Subshape
from pptx.shapes.base import BaseShape
from pptx.table import _Row

from typstpresenter.model.Element import Element
from typstpresenter.model.Table import Table, Row
from typstpresenter.powerpoint.Ignore import Ignore
from typstpresenter.powerpoint.interpreters.SlidePlaceholderInterpreter import interpret_text_frame


class TableInterpreter:

    def can_interpret(self, shape: BaseShape | Subshape) -> bool:
        return shape.has_table

    def __call__(self, shape: BaseShape | Subshape, context: dict | None = None) -> Element | Ignore | None:
        table = shape.table
        return Table(
            rows=tuple(self.__interpret_row(row) for row in table.rows)
        )

    def __interpret_row(self, row: _Row) -> Row:
        return tuple(interpret_text_frame(cell.text_frame) for cell in row.cells)
