from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes import Subshape
from pptx.shapes.base import BaseShape

from typstpresenter.model.Element import Element
from typstpresenter.model.Image import Image
from typstpresenter.powerpoint.Ignore import Ignore
from typstpresenter.powerpoint.image_naming import name_for_image


class PictureInterpreter:
    def can_interpret(self, shape: BaseShape | Subshape) -> bool:
        return hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE

    def __call__(self, shape: BaseShape | Subshape, context: dict | None = None) -> Element | Ignore | None:
        ext = shape.image.ext
        blob = shape.image.blob
        name = name_for_image(blob, ext)
        width = getattr(shape, 'width', None)
        height = getattr(shape, 'height', None)
        width_pt = getattr(width, 'pt', None) if width else None
        height_pt = getattr(height, 'pt', None) if height else None
        return Image(name=name, ext=ext, blob=blob, width_pt=width_pt, height_pt=height_pt)
