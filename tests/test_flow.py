"""
Tests for the idiomatic flow emitter (minimal mode) and Method S.

The contract of minimal mode: human-editable Typst (native lists and
headings, no #place, no probes) that compiles to exactly one page per
slide with no text/image collisions.
"""

import shutil
from pathlib import Path

import pytest

from typstpresenter.convert.emitter import emit_minimal
from typstpresenter.convert.flow import (
    _paired_value_table,
    _paragraph_inline,
    _render_diagram_cluster,
    escape_flow,
    guard_markup_start,
)
from typstpresenter.verify.geometry import BBox
from typstpresenter.verify.method_s import source_metrics, verify_minimal

pytestmark = pytest.mark.skipif(
    shutil.which("typst") is None, reason="typst CLI not on PATH"
)

DATA = Path(__file__).parent / "data"


def test_escape_flow_keeps_text_readable():
    assert escape_flow("Metadaten-Extraktion") == "Metadaten-Extraktion"
    assert escape_flow("C# and $x$") == "C\\# and \\$x\\$"


def test_escape_flow_neutralizes_comments():
    assert "//" not in escape_flow("int x; // counter")
    assert "//" not in escape_flow("http://example.org".replace(":", ""))


def test_guard_markup_start():
    assert guard_markup_start("- item") == "\\- item"
    assert guard_markup_start("/ term") == "\\/ term"
    assert guard_markup_start("= heading") == "\\= heading"
    assert guard_markup_start("1. numbered") == "1\\. numbered"
    assert guard_markup_start("plain") == "plain"
    assert guard_markup_start("  + indented") == "  \\+ indented"


def test_inline_hash_expression_guarded_against_call_chain():
    red = (18.0, False, False, False, "FF0000", None, None)
    plain = (18.0, False, False, False, None, None, None)
    markup = _paragraph_inline([("Hosts ", red), ("(Endsysteme)", plain)], 18.0)
    # without the ';' the '(' would be parsed as a call on the #text result
    assert "];(" in markup


def test_inline_hyperlink_becomes_link():
    linked = (18.0, False, False, False, None, "https://example.org/x?a=1", None)
    markup = _paragraph_inline([("see here", linked)], 18.0)
    assert markup == '#link("https://example.org/x?a=1")[see here]'


def test_inline_highlight_becomes_highlight():
    highlighted = (18.0, False, False, False, None, None, "00FF00")
    markup = _paragraph_inline([("see here", highlighted)], 18.0)
    assert markup == '#highlight(fill: rgb("#00FF00"))[see here]'


def test_paired_value_table_merges_label_and_value_columns():
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Pt(37), Pt(315), Pt(374), Pt(210))
    a.text_frame.text = "In Millionen SLOC"
    for label in ["OpenSolaris:", "Linux kernel:", "Mac OS X:", "Debian:"]:
        a.text_frame.add_paragraph().text = label
    b = slide.shapes.add_textbox(Pt(411), Pt(345), Pt(179), Pt(155))
    b.text_frame.text = "9.7"
    for value in ["12.6", "86", "283 / 324"]:
        b.text_frame.add_paragraph().text = value
    markup = _paired_value_table(a, BBox(37, 315, 374, 210),
                                 b, BBox(411, 345, 179, 155),
                                 18.0, 18.0, 1.0)
    assert markup is not None
    assert "stroke: none" in markup
    # header stays above the table; each label shares a row with its value
    head, table = markup.split("#table", 1)
    assert "In Millionen SLOC" in head
    assert "[OpenSolaris:], [9\\.7]," in table
    assert "[Debian:], [283 / 324]," in table


def test_brace_renders_as_stroke_not_filled_bar():
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt

    from typstpresenter.convert.cetz import emit_cetz_shape

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.LEFT_BRACE, Pt(382), Pt(298), Pt(24), Pt(188))
    markup, _ = emit_cetz_shape(shape, "", BBox(382, 298, 24, 188),
                                probes=False)
    assert "line(" in markup
    assert "rect(" not in markup  # no theme-filled bar over the content


def test_rotated_absorbed_text_gets_canvas_angle(tmp_path):
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    label = slide.shapes.add_textbox(Pt(277), Pt(377), Pt(170), Pt(32))
    label.text_frame.text = "Rahmen (Frame)"
    label.rotation = 270
    _, markup = _render_diagram_cluster(
        [], [("text", label, BBox(277, 377, 170, 32), "s1-e1")],
        720.0, 18.0, 18.0, tmp_path)
    assert "angle: 90deg" in markup
    assert "Rahmen" in markup
def test_simple_deck_is_idiomatic_and_sane(tmp_path):
    pptx_path = tmp_path / "simple.pptx"
    pptx_path.write_bytes((DATA / "simple.pptx").read_bytes())
    typ_path = tmp_path / "simple.typ"
    emit_minimal(pptx_path, typ_path)
    source = typ_path.read_text(encoding="utf-8")

    assert "== Slide 1 Title" in source
    assert "- This is a bullet point" in source
    metrics = source_metrics(source)
    assert metrics.place_calls == 0
    assert metrics.probe_defs == 0

    report = verify_minimal(typ_path, pptx_path)
    assert report.ok, report.summary()


@pytest.mark.parametrize("deck", ["two_content.pptx", "multi_content.pptx",
                                  "media.pptx", "talk_example_a.pptx"])
def test_data_decks_flow_clean(tmp_path, deck):
    src = DATA / deck
    if not src.exists():
        pytest.skip(f"{deck} not available")
    pptx_path = tmp_path / deck
    pptx_path.write_bytes(src.read_bytes())
    typ_path = pptx_path.with_suffix(".typ")
    emit_minimal(pptx_path, typ_path)
    report = verify_minimal(typ_path, pptx_path)
    assert report.ok, report.summary()
