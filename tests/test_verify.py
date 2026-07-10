"""
Tests for the verification tools (Methods A and B).

The corpus (PPTX + paired Typst files, clean and fault-injected) is
generated once per session into a tmp directory; both verification
methods must produce zero issues on clean cases and detect every
injected fault they can see in principle.

Requires the ``typst`` CLI on PATH (skipped otherwise).
"""

import shutil

import pytest

from typstpresenter.verify.compare import compare_by_id, compare_spatial
from typstpresenter.verify.corpus import (
    clean_case_names,
    fault_case_names,
    generate_corpus,
)
from typstpresenter.verify.geometry import BBox
from typstpresenter.verify.method_a import run_method_a
from typstpresenter.verify.method_b import run_method_b

pytestmark = pytest.mark.skipif(
    shutil.which("typst") is None, reason="typst CLI not on PATH"
)


@pytest.fixture(scope="session")
def corpus(tmp_path_factory):
    return generate_corpus(tmp_path_factory.mktemp("verify_corpus"))


def _case(corpus, name):
    return next(c for c in corpus if c.name == name)


def _report_b(case):
    slide0 = case.truth.slides[0]
    result = run_method_b(case.typ_path, slide0.width, slide0.height)
    return compare_by_id(case.truth, result.geometry, overflows=result.overflows)


def _report_a(case):
    result = run_method_a(case.typ_path)
    return compare_spatial(case.truth, result.geometry)


# Derived from the corpus registry so every case added during autoresearch
# iterations is covered automatically.
CLEAN_CASES = clean_case_names()
FAULT_CASES = fault_case_names()


@pytest.mark.parametrize("name", CLEAN_CASES)
def test_clean_case_has_no_issues_method_b(corpus, name):
    report = _report_b(_case(corpus, name))
    assert report.ok, report.summary()


@pytest.mark.parametrize("name", CLEAN_CASES)
def test_clean_case_has_no_issues_method_a(corpus, name):
    report = _report_a(_case(corpus, name))
    assert report.ok, report.summary()


@pytest.mark.parametrize("name", FAULT_CASES)
def test_fault_detected_by_method_b(corpus, name):
    case = _case(corpus, name)
    report = _report_b(case)
    reported = {}
    for issue in report.issues:
        reported.setdefault(issue.element_id, set()).add(issue.kind)
    for element_id, expected in case.expected_issues_b.items():
        assert expected <= reported.get(element_id, set()), (
            f"{element_id}: expected {expected}, reported "
            f"{reported.get(element_id, set())}\n{report.summary()}"
        )
    # no issues on elements without faults
    unexpected = set(reported) - set(case.expected_issues_b)
    assert not unexpected, report.summary()


@pytest.mark.parametrize("name", FAULT_CASES)
def test_fault_detected_by_method_a(corpus, name):
    case = _case(corpus, name)
    report = _report_a(case)
    reported = {}
    for issue in report.issues:
        reported.setdefault(issue.element_id, set()).add(issue.kind)
    for element_id, expected in case.expected_issues_a.items():
        assert expected <= reported.get(element_id, set()), (
            f"{element_id}: expected {expected}, reported "
            f"{reported.get(element_id, set())}\n{report.summary()}"
        )
    unexpected = set(reported) - set(case.expected_issues_a)
    if case.allows_extra_text:
        unexpected -= {None}
    assert not unexpected, report.summary()


def test_method_b_reports_exact_positions(corpus):
    """Probes must return the designed coordinates to sub-pt precision."""
    case = _case(corpus, "layout_two_columns")
    slide0 = case.truth.slides[0]
    result = run_method_b(case.typ_path, slide0.width, slide0.height)
    found = {e.id: e for s in result.geometry.slides for e in s.elements}
    for _, truth_el in case.truth.all_elements():
        probe = found[truth_el.id]
        assert abs(probe.bbox.x - truth_el.bbox.x) < 0.02
        assert abs(probe.bbox.y - truth_el.bbox.y) < 0.02
        assert abs(probe.bbox.w - truth_el.bbox.w) < 0.02
        assert abs(probe.bbox.h - truth_el.bbox.h) < 0.02


def test_inherited_styles_resolved_from_master(corpus, tmp_path):
    """Placeholder text without explicit formatting must resolve to the
    Office-default master styles (title 44pt centered, body 32/28/24pt)."""
    import pptx
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    from typstpresenter.verify.pptx_inherit import (
        resolve_alignment,
        resolve_anchor,
        resolve_bullet,
        resolve_font_size_pt,
    )

    case = _case(corpus, "layout_placeholders")
    prs = pptx.Presentation(str(case.pptx_path))

    title = prs.slides[0].shapes.title
    p = title.text_frame.paragraphs[0]
    assert resolve_font_size_pt(p.runs[0], p, title) == 44.0
    assert resolve_alignment(p, title) == PP_ALIGN.CENTER
    assert resolve_anchor(title) == MSO_ANCHOR.MIDDLE

    body = prs.slides[1].placeholders[1]
    sizes = [
        resolve_font_size_pt(p.runs[0] if p.runs else None, p, body)
        for p in body.text_frame.paragraphs
    ]
    assert sizes == [32.0, 32.0, 28.0, 24.0]
    assert resolve_bullet(body.text_frame.paragraphs[0], body) is not None


def test_bbox_iou_and_union():
    a = BBox(0, 0, 10, 10)
    b = BBox(5, 5, 10, 10)
    assert a.intersection_area(b) == 25.0
    assert a.iou(b) == pytest.approx(25.0 / 175.0)
    u = a.union(b)
    assert (u.x, u.y, u.w, u.h) == (0, 0, 15, 15)
    assert BBox(0, 0, 1, 1).iou(BBox(5, 5, 1, 1)) == 0.0
