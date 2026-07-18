"""Contract tests for the diagram2svg package (see docs/diagram-to-svg-plan.md).

The complexity-ladder decks are regenerated on the fly (deterministic);
every level built so far must pass the structural gate with 0 findings.
"""

from __future__ import annotations

import math

import pytest

from typstpresenter.diagram2svg.convert import pptx_to_svgs
from typstpresenter.diagram2svg.presets import (
    GuideContext,
    evaluate_preset,
    load_presets,
)
from typstpresenter.diagram2svg.structural import check_svg_against_pptx

import sys

sys.path.insert(0, str(__import__("pathlib").Path(__file__).parent))
from generate_diagram_data import build_l1, build_l2, build_l3  # noqa: E402

# ------------------------------------------------------- formula interpreter


def test_presets_load_complete():
    presets = load_presets()
    assert len(presets) == 187
    for name in ("rect", "ellipse", "roundRect", "diamond", "triangle", "line",
                 "rightArrow", "chevron", "wedgeRectCallout", "can", "cloud"):
        assert name in presets


def test_guide_formula_ops():
    ctx = GuideContext(w=200.0, h=100.0)
    assert ctx.evaluate("val 42") == 42
    assert ctx.evaluate("*/ w 1 2") == 100  # w*1/2
    assert ctx.evaluate("+- r 0 10") == 190
    assert ctx.evaluate("+/ w h 2") == 150
    assert ctx.evaluate("pin 10 5 20") == 10
    assert ctx.evaluate("pin 10 15 20") == 15
    assert ctx.evaluate("pin 10 50 20") == 20
    assert ctx.evaluate("max w h") == 200
    assert ctx.evaluate("min w h") == 100
    assert ctx.evaluate("mod 3 4 0") == 5
    assert ctx.evaluate("sqrt 81") == 9
    assert ctx.evaluate("?: 1 7 9") == 7
    assert ctx.evaluate("?: -1 7 9") == 9
    # trig: angles in 1/60000 deg; cos(60°) = 0.5
    assert ctx.evaluate("cos 100 3600000") == pytest.approx(50)
    assert ctx.evaluate("sin 100 1800000") == pytest.approx(50)
    # built-in tokens
    assert ctx.resolve("wd2") == 100
    assert ctx.resolve("hd4") == 25
    assert ctx.resolve("ss") == 100
    assert ctx.resolve("cd4") == 5400000
    assert ctx.resolve("3cd4") == 16200000


def test_rect_preset_corners():
    (path,) = evaluate_preset("rect", 200.0, 100.0)
    pts = [(s[1], s[2]) for s in path.segments if s[0] in ("M", "L")]
    assert pts == [(0, 0), (200, 0), (200, 100), (0, 100)]
    assert path.segments[-1][0] == "Z"


def test_diamond_preset_midpoints():
    (path,) = evaluate_preset("diamond", 200.0, 100.0)
    pts = [(s[1], s[2]) for s in path.segments if s[0] in ("M", "L")]
    assert pts == [(0, 50), (100, 0), (200, 50), (100, 100)]


def test_ellipse_preset_is_arcs_and_spans_bbox():
    (path,) = evaluate_preset("ellipse", 200.0, 100.0)
    arc_ends = [(s[4], s[5]) for s in path.segments if s[0] == "A"]
    assert len(arc_ends) >= 4
    xs = [p[0] for p in arc_ends]
    ys = [p[1] for p in arc_ends]
    assert math.isclose(min(xs), 0, abs_tol=1e-6) or math.isclose(min(ys), 0, abs_tol=1e-6)
    # 90°-arc endpoints hit the bbox extremes
    assert math.isclose(max(xs), 200, abs_tol=1e-6)
    assert math.isclose(max(ys), 100, abs_tol=1e-6)


def test_roundrect_adjustment_shrinks_corner():
    (default,) = evaluate_preset("roundRect", 200.0, 200.0)
    (tight,) = evaluate_preset("roundRect", 200.0, 200.0, {"adj": "val 5000"})
    # first moveTo y == corner radius x1 = ss*adj/100000
    assert default.segments[0][2] == pytest.approx(200 * 16667 / 100000)
    assert tight.segments[0][2] == pytest.approx(200 * 5000 / 100000)


# ----------------------------------------------------------- L1 end-to-end


@pytest.fixture(scope="module")
def l1_results(tmp_path_factory):
    pptx_path = build_l1()
    out = tmp_path_factory.mktemp("l1_svg")
    return pptx_path, pptx_to_svgs(pptx_path, out)


def test_l1_structural_gate(l1_results):
    pptx_path, results = l1_results
    assert len(results) == 7
    all_findings = []
    for i, (svg_path, res) in enumerate(results):
        assert res.shape_count > 0
        assert not res.fallbacks, f"slide {i + 1}: unexpected fallbacks {res.fallbacks}"
        report = check_svg_against_pptx(
            svg_path.read_text(encoding="utf-8"), pptx_path, i
        )
        assert report.checked == res.shape_count
        all_findings += [(i + 1, f) for f in report.findings]
    assert not all_findings, "\n".join(
        f"slide {s}: {f.shape_id} {f.kind}: {f.detail}" for s, f in all_findings
    )


def test_l2_structural_and_text_gate(tmp_path):
    from typstpresenter.diagram2svg.structural import check_text_and_markers

    pptx_path = build_l2()
    results = pptx_to_svgs(pptx_path, tmp_path)
    assert len(results) == 4
    all_findings = []
    for i, (svg_path, res) in enumerate(results):
        svg = svg_path.read_text(encoding="utf-8")
        assert not res.fallbacks
        rep1 = check_svg_against_pptx(svg, pptx_path, i)
        rep2 = check_text_and_markers(svg, pptx_path, i)
        all_findings += [(i + 1, f) for f in rep1.findings + rep2.findings]
    assert not all_findings, "\n".join(
        f"slide {s}: {f.shape_id} {f.kind}: {f.detail}" for s, f in all_findings
    )


def test_l2_theme_color_resolution(tmp_path):
    """Slide 1 shapes carry no explicit fill; accent1 of the template's
    theme (Office 2007: 4F81BD) must be resolved through the p:style
    fillRef, and the fontRef default (lt1 white) must reach the label."""
    pptx_path = build_l2()
    results = pptx_to_svgs(pptx_path, tmp_path)
    svg = results[0][0].read_text(encoding="utf-8")
    assert '#4F81BD' in svg
    assert 'fill="#FFFFFF"' in svg


def test_l2_wrap_and_marker_emission(tmp_path):
    pptx_path = build_l2()
    results = pptx_to_svgs(pptx_path, tmp_path)
    svg2 = results[1][0].read_text(encoding="utf-8")
    # long text must be broken into several tspans (wrapped lines)
    assert svg2.count("<text>") >= 4
    svg3 = results[2][0].read_text(encoding="utf-8")
    assert "marker-end" in svg3 and "marker-start" in svg3
    assert "<defs>" in svg3


def test_l3_structural_gate(tmp_path):
    from typstpresenter.diagram2svg.structural import check_text_and_markers

    pptx_path = build_l3()
    results = pptx_to_svgs(pptx_path, tmp_path)
    assert len(results) == 6
    all_findings = []
    for i, (svg_path, res) in enumerate(results):
        svg = svg_path.read_text(encoding="utf-8")
        assert not res.fallbacks, f"slide {i + 1}: {res.fallbacks}"
        rep1 = check_svg_against_pptx(svg, pptx_path, i)
        rep2 = check_text_and_markers(svg, pptx_path, i)
        all_findings += [(i + 1, f) for f in rep1.findings + rep2.findings]
    assert not all_findings, "\n".join(
        f"slide {s}: {f.shape_id} {f.kind}: {f.detail}" for s, f in all_findings
    )


def test_l3_feature_emission(tmp_path):
    pptx_path = build_l3()
    results = pptx_to_svgs(pptx_path, tmp_path)
    svgs = [p.read_text(encoding="utf-8") for p, _ in results]
    # rotation and flip transforms
    assert "rotate(30)" in svgs[0] and "rotate(180)" in svgs[0]
    assert "scale(-1 1)" in svgs[0]
    # curved connector uses cubic beziers, elbow only lines
    assert "C" in svgs[3]
    # gradients and dashes
    assert "<linearGradient" in svgs[4]
    assert "stroke-dasharray" in svgs[4]
    # freeform custGeom is a real polygon, not a bbox fallback
    assert results[5][1].fallbacks == []


def test_custgeom_freeform_polygon(tmp_path):
    """The freeform's SVG path must trace the drawn pentagon, not a rect."""
    import xml.etree.ElementTree as ET

    pptx_path = build_l3()
    results = pptx_to_svgs(pptx_path, tmp_path)
    root = ET.fromstring(results[5][0].read_text(encoding="utf-8"))
    ns = {"svg": "http://www.w3.org/2000/svg"}
    groups = root.findall("svg:g", ns)
    d = groups[-1].find("svg:path", ns).get("d")
    # pentagon: 1 moveTo + 4 lnTo + close
    assert d.count("L") == 4 and d.count("Z") == 1


def test_symbol_font_mapping():
    from typstpresenter.diagram2svg.symbols import map_symbol_text

    assert map_symbol_text("l", "Symbol") == "λ"
    assert map_symbol_text("m", "Symbol") == "μ"
    assert map_symbol_text("2m", "Symbol") == "2μ"
    # PUA-encoded Symbol char without declared font
    assert map_symbol_text("", None) == "μ"
    # Wingdings bullet square
    assert map_symbol_text("", "Wingdings") == "▪"
    # plain text untouched
    assert map_symbol_text("Hello", "Calibri") == "Hello"


def test_l4_l5_extracts_structural_gate(tmp_path):
    """Real-deck extracts: every geometric shape translates without
    bbox fallback and passes the structural comparison."""
    from generate_diagram_data import DATA_DIR

    from typstpresenter.diagram2svg.structural import check_text_and_markers

    extracts = sorted(DATA_DIR.glob("L4-*.pptx")) + sorted(DATA_DIR.glob("L5-*.pptx"))
    if not extracts:
        pytest.skip("extract decks not generated (source decks unavailable)")
    all_findings = []
    for pptx_path in extracts:
        results = pptx_to_svgs(pptx_path, tmp_path / pptx_path.stem)
        for i, (svg_path, res) in enumerate(results):
            svg = svg_path.read_text(encoding="utf-8")
            assert not res.fallbacks, f"{pptx_path.stem}: {res.fallbacks}"
            rep1 = check_svg_against_pptx(svg, pptx_path, i)
            rep2 = check_text_and_markers(svg, pptx_path, i)
            all_findings += [
                (pptx_path.stem, f) for f in rep1.findings + rep2.findings
            ]
    assert not all_findings, "\n".join(
        f"{d}: {f.shape_id} {f.kind}: {f.detail}" for d, f in all_findings
    )


def test_svg_backend_flow_integration(tmp_path, monkeypatch):
    """flow.py emits #image(...svg) clusters when the backend flag is set."""
    import typstpresenter.convert.flow as flow
    from typstpresenter.convert.emitter import emit_touying

    monkeypatch.setattr(flow, "DIAGRAM_BACKEND", "svg")
    pptx_path = build_l2()
    typ_path = tmp_path / "L2.typ"
    emit_touying(pptx_path, typ_path, minimal=True)
    text = typ_path.read_text(encoding="utf-8")
    assert ".svg" in text and "image(" in text
    media = tmp_path / "L2_media"
    assert list(media.glob("*.svg")), "no SVG cluster files written"
    # slide 3's flowchart is lifted to Fletcher before the backend runs
    # (idiomatic output outranks both raster paths); slide 1's cluster
    # must carry the theme-styled boxes with their labels
    svg = next(media.glob("s1-*.svg")).read_text(encoding="utf-8")
    assert "<tspan" in svg and "#4F81BD" in svg


def test_g3_normautofit_shrinks_text(tmp_path):
    """A normAutofit shape with too much text shrinks it into the box."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(100), Pt(100), Pt(200), Pt(60))
    box.text_frame.text = ("Sehr viel Text " * 12).strip()
    for para in box.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
    bodyPr = box.text_frame._txBody.find(qn("a:bodyPr"))
    bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))
    path = tmp_path / "autofit.pptx"
    prs.save(path)

    results = pptx_to_svgs(path, tmp_path)
    svg = results[0][0].read_text(encoding="utf-8")
    import re

    sizes = [float(m) for m in re.findall(r'font-size="([\d.]+)"', svg)]
    assert sizes and max(sizes) < 24 * 0.95, f"text not shrunk: {sizes}"


def test_g5_spatial_components():
    from typstpresenter.convert.flow import _spatial_components
    from typstpresenter.verify.geometry import BBox

    def shp(x, y, w=100, h=60):
        return (object(), BBox(x, y, w, h))

    # two far-apart pairs split; a connector bridging them keeps one group
    left = [shp(50, 100), shp(50, 180)]
    right = [shp(600, 100), shp(600, 180)]
    comps = _spatial_components(left + right, [])
    assert len(comps) == 2
    bridge = [shp(140, 120, 460, 20)]  # overlaps both sides
    comps = _spatial_components(left + right + bridge, [])
    assert len(comps) == 1
    # a lone tiny speck merges into the nearest substantial group
    speck = [shp(300, 500, 10, 10)]
    comps = _spatial_components(left + right + speck, [])
    assert len(comps) == 2
    assert sum(len(c[0]) for c in comps) == 5


def test_g6_ole_preview_extraction():
    from generate_diagram_data import DATA_DIR

    import pptx as _pptx

    from typstpresenter.verify.pptx_geometry import iter_flat_shapes, picture_image

    deck = DATA_DIR / "L6-vlN01-s22-ole.pptx"
    if not deck.exists():
        pytest.skip("L6 extract not generated")
    prs = _pptx.Presentation(deck)
    previews = [
        picture_image(sh)
        for sh, _ in iter_flat_shapes(list(prs.slides)[0].shapes)
        if sh.shape_id in (28, 29)
    ]
    assert len(previews) == 2 and all(p is not None for p in previews)


def test_l1_svg_wellformed_and_sized(l1_results):
    import xml.etree.ElementTree as ET

    _, results = l1_results
    for svg_path, _res in results:
        root = ET.fromstring(svg_path.read_text(encoding="utf-8"))
        assert root.get("viewBox") == "0 0 960 540"
        assert root.get("width") == "960pt"
