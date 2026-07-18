"""Generate the diagram2svg complexity-ladder test decks (L1..L3).

Run: uv run python tests/generate_diagram_data.py
Writes tests/diagram2svg_data/L{n}.pptx (deterministic, regenerable).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Emu, Pt

DATA_DIR = Path(__file__).parent / "diagram2svg_data"

# 16:9 slide, 960x540 pt
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def _new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _style(shape, fill: str | None, line: str | None, width_pt: float = 2.0):
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor.from_string(line)
        shape.line.width = Pt(width_pt)
    if shape.has_text_frame:
        shape.text_frame.clear()  # no text in L1
    shape.shadow.inherit = False


def build_l1() -> Path:
    """One primitive per slide + one composite slide. No text, no rotation."""
    prs = _new_deck()

    # slide 1: rectangle
    s = _blank(prs)
    _style(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(120), Pt(90), Pt(300), Pt(180)),
           "4472C4", "1F3864")

    # slide 2: rounded rectangle, default + adjusted corner radius
    s = _blank(prs)
    _style(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(90), Pt(120), Pt(260), Pt(150)),
           "ED7D31", "833C00")
    rr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(480), Pt(120), Pt(260), Pt(150))
    rr.adjustments[0] = 0.45
    _style(rr, "FFC000", "7F6000")

    # slide 3: ellipse + circle
    s = _blank(prs)
    _style(s.shapes.add_shape(MSO_SHAPE.OVAL, Pt(100), Pt(100), Pt(320), Pt(180)),
           "70AD47", "375623")
    _style(s.shapes.add_shape(MSO_SHAPE.OVAL, Pt(560), Pt(100), Pt(180), Pt(180)),
           "A9D18E", None)

    # slide 4: diamond
    s = _blank(prs)
    _style(s.shapes.add_shape(MSO_SHAPE.DIAMOND, Pt(180), Pt(110), Pt(240), Pt(200)),
           "FF0000", "7F0000")

    # slide 5: isosceles triangle (default apex + moved apex)
    s = _blank(prs)
    _style(s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Pt(120), Pt(120), Pt(240), Pt(190)),
           "9E5FBE", "4B2064")
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Pt(480), Pt(120), Pt(240), Pt(190))
    tri.adjustments[0] = 0.25
    _style(tri, "C5A0D5", "4B2064")

    # slide 6: straight lines (plain, thick, and upward = flipV)
    s = _blank(prs)
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(120), Pt(140), Pt(420), Pt(300))
    ln.line.color.rgb = RGBColor.from_string("1F3864")
    ln.line.width = Pt(2.0)
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(520), Pt(140), Pt(820), Pt(140))
    ln.line.color.rgb = RGBColor.from_string("C00000")
    ln.line.width = Pt(4.5)
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(120), Pt(460), Pt(420), Pt(340))
    ln.line.color.rgb = RGBColor.from_string("375623")
    ln.line.width = Pt(2.0)

    # slide 7: composite — all primitives together, overlap-free
    s = _blank(prs)
    _style(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(60), Pt(60), Pt(200), Pt(120)),
           "4472C4", "1F3864")
    _style(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(320), Pt(60), Pt(200), Pt(120)),
           "ED7D31", "833C00")
    _style(s.shapes.add_shape(MSO_SHAPE.OVAL, Pt(580), Pt(60), Pt(200), Pt(120)),
           "70AD47", "375623")
    _style(s.shapes.add_shape(MSO_SHAPE.DIAMOND, Pt(60), Pt(280), Pt(200), Pt(160)),
           "FF0000", "7F0000")
    _style(s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Pt(320), Pt(280), Pt(200), Pt(160)),
           "9E5FBE", "4B2064")
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(580), Pt(300), Pt(830), Pt(420))
    ln.line.color.rgb = RGBColor.from_string("1F3864")
    ln.line.width = Pt(3.0)

    DATA_DIR.mkdir(parents=True, exist_ok=True)
    out = DATA_DIR / "L1.pptx"
    prs.save(out)
    return out


def _arrow(conn, end_tag: str, kind: str, w: str = "med", length: str = "med"):
    from pptx.oxml.ns import qn

    ln = conn.line._get_or_add_ln()
    el = ln.makeelement(qn(f"a:{end_tag}"), {"type": kind, "w": w, "len": length})
    ln.append(el)


def _text(shape, value: str, size_pt: float = 18.0, color: str | None = None,
          bold: bool = False):
    tf = shape.text_frame
    tf.text = value
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            if color:
                run.font.color.rgb = RGBColor.from_string(color)


def build_l2() -> Path:
    """Small assemblies: text labels, theme colors, connectors with
    arrowheads, z-order overlap."""
    prs = _new_deck()

    # slide 1: theme-styled boxes (no explicit colors -> accent1 via p:style)
    s = _blank(prs)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(120), Pt(120), Pt(220), Pt(110))
    box.shadow.inherit = False
    _text(box, "Eingabe")
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(420), Pt(120), Pt(220), Pt(110))
    box.shadow.inherit = False
    _text(box, "Verarbeitung")
    box = s.shapes.add_shape(MSO_SHAPE.OVAL, Pt(720), Pt(120), Pt(220), Pt(110))
    box.shadow.inherit = False
    _text(box, "Ausgabe")

    # slide 2: styled text — wrapping, alignments, mixed runs
    s = _blank(prs)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(80), Pt(80), Pt(300), Pt(200))
    _style(box, "DEEBF7", "2E75B6", 1.5)
    _text(box, "Dieser etwas laengere Text muss innerhalb der Box "
               "auf mehrere Zeilen umgebrochen werden", size_pt=16, color="1F3864")
    box2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(460), Pt(80), Pt(320), Pt(200))
    _style(box2, "FBE5D6", "C55A11", 1.5)
    tf = box2.text_frame
    tf.text = "Fett und wichtig"
    p = tf.paragraphs[0]
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor.from_string("C00000")
    from pptx.enum.text import PP_ALIGN as _AL
    p2 = tf.add_paragraph()
    p2.text = "kursiv daneben"
    p2.alignment = _AL.LEFT
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.italic = True
    p2.runs[0].font.color.rgb = RGBColor.from_string("1F3864")
    p3 = tf.add_paragraph()
    p3.text = "rechtsbuendig"
    p3.alignment = _AL.RIGHT
    p3.runs[0].font.size = Pt(16)
    p3.runs[0].font.color.rgb = RGBColor.from_string("375623")

    # slide 3: mini flowchart with arrowhead connectors
    s = _blank(prs)
    a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(100), Pt(200), Pt(180), Pt(90))
    _style(a, "4472C4", "1F3864")
    _text(a, "Start", color="FFFFFF")
    b = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Pt(400), Pt(180), Pt(180), Pt(130))
    _style(b, "FFC000", "7F6000")
    _text(b, "ok?", color="000000")
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(700), Pt(200), Pt(180), Pt(90))
    _style(c, "70AD47", "375623")
    _text(c, "Ende", color="FFFFFF")
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(280), Pt(245), Pt(400), Pt(245))
    ln.line.color.rgb = RGBColor.from_string("404040")
    ln.line.width = Pt(2.25)
    _arrow(ln, "tailEnd", "triangle")
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(580), Pt(245), Pt(700), Pt(245))
    ln.line.color.rgb = RGBColor.from_string("404040")
    ln.line.width = Pt(2.25)
    _arrow(ln, "tailEnd", "stealth")
    _arrow(ln, "headEnd", "oval", "sm", "sm")

    # slide 4: z-order — overlapping shapes paint in document order
    s = _blank(prs)
    r1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(200), Pt(120), Pt(300), Pt(200))
    _style(r1, "4472C4", "1F3864")
    r2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Pt(350), Pt(200), Pt(300), Pt(200))
    _style(r2, "ED7D31", "833C00")
    r3 = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Pt(500), Pt(140), Pt(280), Pt(200))
    _style(r3, "70AD47", "375623")

    out = DATA_DIR / "L2.pptx"
    prs.save(out)
    return out


def build_l3() -> Path:
    """Hard single-shape features: rotation/flips, census presets, elbow +
    curved connectors, gradients, dashes, groups, freeform custGeom."""
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    prs = _new_deck()

    # slide 1: rotation and flips
    s = _blank(prs)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(120), Pt(140), Pt(220), Pt(110))
    _style(r, "4472C4", "1F3864")
    r.rotation = 30
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Pt(430), Pt(140), Pt(180), Pt(140))
    _style(tri, "ED7D31", "833C00")
    tri.rotation = 180
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Pt(700), Pt(140), Pt(200), Pt(90))
    _style(arrow, "70AD47", "375623")
    arrow.rotation = 90
    r2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Pt(120), Pt(360), Pt(180), Pt(120))
    _style(r2, "FFC000", "7F6000")
    from pptx.oxml.ns import qn as _qn
    r2.element.find(_qn("p:spPr")).find(_qn("a:xfrm")).set("flipH", "1")

    # slide 2: gap-census presets via the generic interpreter
    s = _blank(prs)
    presets = [
        (MSO_SHAPE.RIGHT_ARROW, "rightArrow"),
        (MSO_SHAPE.CHEVRON, "chevron"),
        (MSO_SHAPE.PARALLELOGRAM, "parallelogram"),
        (MSO_SHAPE.CAN, "can"),
        (MSO_SHAPE.CLOUD, "cloud"),
        (MSO_SHAPE.HEXAGON, "hexagon"),
        (MSO_SHAPE.LIGHTNING_BOLT, "lightningBolt"),
        (MSO_SHAPE.MATH_PLUS, "mathPlus"),
    ]
    colors = ["4472C4", "ED7D31", "70AD47", "FFC000", "FF0000", "9E5FBE",
              "2E75B6", "C55A11"]
    for i, (mso, _name) in enumerate(presets):
        x = 60 + (i % 4) * 230
        y = 80 + (i // 4) * 220
        sh = s.shapes.add_shape(mso, Pt(x), Pt(y), Pt(180), Pt(140))
        _style(sh, colors[i], "404040", 1.5)

    # slide 3: callouts with adjusted leader + arc (open shape)
    s = _blank(prs)
    co = s.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, Pt(120), Pt(100), Pt(240), Pt(120))
    _style(co, "FBE5D6", "C55A11", 1.5)
    co.adjustments[0] = -0.3
    co.adjustments[1] = 1.2
    ar = s.shapes.add_shape(MSO_SHAPE.ARC, Pt(500), Pt(100), Pt(200), Pt(200))
    ar.fill.background()
    ar.line.color.rgb = RGBColor.from_string("1F3864")
    ar.line.width = Pt(3)
    ar.shadow.inherit = False

    # slide 4: elbow and curved connectors with arrowheads
    s = _blank(prs)
    a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(100), Pt(100), Pt(160), Pt(80))
    _style(a, "4472C4", "1F3864")
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(600), Pt(320), Pt(160), Pt(80))
    _style(b, "70AD47", "375623")
    ln = s.shapes.add_connector(MSO_CONNECTOR.ELBOW, Pt(260), Pt(140), Pt(600), Pt(360))
    ln.line.color.rgb = RGBColor.from_string("404040")
    ln.line.width = Pt(2.25)
    _arrow(ln, "tailEnd", "triangle")
    ln = s.shapes.add_connector(MSO_CONNECTOR.CURVE, Pt(180), Pt(180), Pt(680), Pt(320))
    ln.line.color.rgb = RGBColor.from_string("C00000")
    ln.line.width = Pt(2.25)
    _arrow(ln, "tailEnd", "stealth")

    # slide 5: gradient fills and dashed strokes
    s = _blank(prs)
    g1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(100), Pt(120), Pt(260), Pt(160))
    g1.fill.gradient()
    stops = g1.fill.gradient_stops
    stops[0].color.rgb = RGBColor.from_string("4472C4")
    stops[1].color.rgb = RGBColor.from_string("FFFFFF")
    g1.line.color.rgb = RGBColor.from_string("1F3864")
    g1.line.width = Pt(1.5)
    g1.shadow.inherit = False
    g2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Pt(460), Pt(120), Pt(260), Pt(160))
    g2.fill.gradient()
    stops = g2.fill.gradient_stops
    stops[0].color.rgb = RGBColor.from_string("FFC000")
    stops[1].color.rgb = RGBColor.from_string("C00000")
    try:
        g2.fill.gradient_angle = 0
    except (ValueError, AttributeError):
        pass
    g2.line.fill.background()
    g2.shadow.inherit = False
    d1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(100), Pt(360), Pt(260), Pt(120))
    _style(d1, None, "1F3864", 2.25)
    d1.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    d2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Pt(460), Pt(360), Pt(260), Pt(120))
    _style(d2, None, "C00000", 2.25)
    d2.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

    # slide 6: group with children + freeform polygon
    s = _blank(prs)
    grp = s.shapes.add_group_shape()
    ga = grp.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(120), Pt(120), Pt(140), Pt(90))
    _style(ga, "4472C4", "1F3864")
    gb = grp.shapes.add_shape(MSO_SHAPE.OVAL, Pt(300), Pt(120), Pt(140), Pt(90))
    _style(gb, "ED7D31", "833C00")
    fb = s.shapes.build_freeform(Pt(600), Pt(120), scale=1.0)
    fb.add_line_segments(
        [(Pt(760), Pt(180)), (Pt(720), Pt(300)), (Pt(640), Pt(300)), (Pt(560), Pt(200))],
        close=True,
    )
    ff = fb.convert_to_shape()
    _style(ff, "9E5FBE", "4B2064")

    out = DATA_DIR / "L3.pptx"
    prs.save(out)
    return out


# --------------------------------------------------- real-deck extracts --

_DATA = Path(__file__).parent / "data"

# (level, source deck relative to tests/data, 1-based slide number, slug)
_EXTRACTS = [
    ("L4", "IBN_presentations2/vl04-ibn (processes, threads, race cond A).pptx", 9, "vl04-s9-states"),
    ("L4", "IBN_presentations2/vl15-ibn (deadlocks A).pptx", 12, "vl15-s12"),
    ("L4", "IBN_presentations2/vl09-ibn (mem mgt B).pptx", 24, "vl09-s24"),
    ("L4", "IBN_presentations2/vl03-ibn (win APIs, processes).pptx", 31, "vl03-s31"),
    ("L4", "IBN_presentations2/vl16-ibn (deadlocks B + scheduling A).pptx", 2, "vl16-s2"),
    ("L4", "IBN_presentations2/vl10-ibn (mem mgt C).pptx", 14, "vl10-s14"),
    ("L5", "IBN_presentations2/vl17-ibn (scheduling B + IO + security).pptx", 19, "vl17-s19-scheduling"),
    ("L5", "IBN_presentations/vlN01-ibn.pptx", 31, "vlN01-s31"),
    ("L5", "IBN_presentations/vlxN03-ibn.pptx", 17, "vlxN03-s17"),
    ("L5", "IBN_presentations/vlxN04-ibn.pptx", 30, "vlxN04-s30"),
    ("L6", "IBN_presentations/vlN01-ibn.pptx", 22, "vlN01-s22-ole"),
]


def _extract_slide(src: Path, slide_no: int, dest: Path):
    """Copy the deck and delete every slide except `slide_no` (1-based)."""
    import shutil

    shutil.copy(src, dest)
    prs = Presentation(dest)
    sld_id_lst = prs.slides._sldIdLst
    for idx, sld_id in reversed(list(enumerate(sld_id_lst))):
        if idx != slide_no - 1:
            sld_id_lst.remove(sld_id)
    prs.save(dest)


def build_l4_l5() -> list[Path]:
    """One-slide extracts of real lecture-deck diagram slides."""
    out_paths = []
    for level, rel, slide_no, slug in _EXTRACTS:
        src = _DATA / rel
        if not src.exists():
            print(f"skip (missing source): {src}")
            continue
        dest = DATA_DIR / f"{level}-{slug}.pptx"
        _extract_slide(src, slide_no, dest)
        out_paths.append(dest)
    return out_paths


if __name__ == "__main__":
    print(build_l1())
    print(build_l2())
    print(build_l3())
    for p in build_l4_l5():
        print(p)
